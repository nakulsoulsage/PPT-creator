#!/usr/bin/env python3
"""
Tier-2/3 Healthcare Disruption - Professional 3-Slide Presentation
Clean, non-overlapping design with proper spacing and visual hierarchy
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import FancyBboxPatch, Circle, Rectangle, FancyArrowPatch
import numpy as np
import io

class HealthcarePresentationCreator:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 widescreen
        self.prs.slide_height = Inches(7.5)

        # Professional color palette
        self.colors = {
            'primary_blue': RGBColor(0, 102, 204),      # Professional blue
            'accent_teal': RGBColor(0, 153, 153),       # Teal accent
            'dark_navy': RGBColor(25, 42, 86),          # Dark navy
            'success_green': RGBColor(76, 175, 80),     # Success green
            'warning_orange': RGBColor(255, 152, 0),    # Warning orange
            'danger_red': RGBColor(244, 67, 54),        # Danger red
            'dark_gray': RGBColor(66, 66, 66),          # Dark gray text
            'medium_gray': RGBColor(117, 117, 117),     # Medium gray
            'light_gray': RGBColor(189, 189, 189),      # Light gray
            'very_light_gray': RGBColor(245, 245, 245), # Background gray
            'white': RGBColor(255, 255, 255),           # White
        }

    def _add_title(self, slide, main_title, subtitle=None):
        """Add title with optional subtitle"""
        # Main title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.3), Inches(0.6)
        )
        tf = title_box.text_frame
        tf.clear()
        tf.margin_left = tf.margin_right = Inches(0.1)
        tf.margin_top = tf.margin_bottom = Inches(0.05)

        p = tf.paragraphs[0]
        p.text = main_title
        p.font.name = 'Segoe UI'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors['dark_navy']
        p.alignment = PP_ALIGN.CENTER

        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.95),
                Inches(12.3), Inches(0.8)
            )
            tf = subtitle_box.text_frame
            tf.clear()
            tf.margin_all = Inches(0.1)
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.name = 'Segoe UI'
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors['medium_gray']
            p.alignment = PP_ALIGN.CENTER
            p.font.italic = True

    def _add_footer(self, slide, slide_number):
        """Add footer with presenters' names"""
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(7.15),
            Inches(12.3), Inches(0.3)
        )
        tf = footer_box.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        p.text = "Presented By — Nakul Nandanwar, Vaishnavi Bhangale, Rahul Kumbhare"
        p.font.name = 'Segoe UI'
        p.font.size = Pt(10)
        p.font.color.rgb = self.colors['medium_gray']
        p.alignment = PP_ALIGN.CENTER

        # Slide number
        num_box = slide.shapes.add_textbox(
            Inches(12.5), Inches(7.15),
            Inches(0.3), Inches(0.3)
        )
        tf = num_box.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        p.text = str(slide_number)
        p.font.name = 'Segoe UI'
        p.font.size = Pt(10)
        p.font.color.rgb = self.colors['medium_gray']
        p.alignment = PP_ALIGN.CENTER

    def _add_bottom_banner(self, slide, text):
        """Add bottom insight banner"""
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(6.9),
            Inches(13.333), Inches(0.25)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = self.colors['primary_blue']
        banner.line.fill.background()

        tf = banner.text_frame
        tf.clear()
        tf.margin_all = Inches(0.02)

        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Segoe UI'
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    def _create_india_map_visual(self):
        """Create India map highlighting Tier-2/3 cities"""
        fig, ax = plt.subplots(figsize=(6, 5), dpi=150)
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 10)
        ax.axis('off')

        # Simplified India outline
        india_outline = plt.Polygon([
            (3, 2), (4, 1.5), (5, 1.5), (6, 2),
            (7, 3), (7.5, 4), (7, 5), (6.5, 6),
            (6, 7), (5.5, 7.5), (5, 8), (4, 7.5),
            (3, 7), (2.5, 6), (2, 5), (2, 4),
            (2.5, 3), (3, 2)
        ], facecolor='#f0f0f0', edgecolor='#333333', linewidth=2)
        ax.add_patch(india_outline)

        # Convert RGBColor to tuple for matplotlib
        teal_color = (0/255, 153/255, 153/255)
        blue_color = (0/255, 102/255, 204/255)

        # Highlight Tier-2/3 regions
        tier23_regions = [
            Circle((3.5, 5), 0.4, color=teal_color, alpha=0.6),
            Circle((5, 4), 0.5, color=teal_color, alpha=0.6),
            Circle((4, 6), 0.4, color=teal_color, alpha=0.6),
            Circle((5.5, 5.5), 0.45, color=teal_color, alpha=0.6),
            Circle((3, 3.5), 0.35, color=teal_color, alpha=0.6),
        ]

        for region in tier23_regions:
            ax.add_patch(region)

        # Metro cities (smaller, different color)
        metros = [
            Circle((3, 6.5), 0.25, color=blue_color, alpha=0.8),  # Delhi
            Circle((3.5, 4), 0.25, color=blue_color, alpha=0.8),  # Mumbai
            Circle((5.5, 3), 0.25, color=blue_color, alpha=0.8),  # Chennai
            Circle((6, 5), 0.25, color=blue_color, alpha=0.8),    # Kolkata
        ]

        for metro in metros:
            ax.add_patch(metro)

        # Legend
        ax.text(1, 9, "● Metro Cities", fontsize=10, color=blue_color)
        ax.text(1, 8.5, "● Tier-2/3 Cities", fontsize=10, color=teal_color)

        plt.title("India: Tier-2/3 Market Distribution", fontsize=14, fontweight='bold', pad=20)

        buf = io.BytesIO()
        plt.savefig(buf, format='png', bbox_inches='tight', facecolor='white', edgecolor='none')
        plt.close()
        buf.seek(0)
        return buf

    def _create_infrastructure_comparison_chart(self):
        """Create comparison chart for urban vs Tier-2/3"""
        fig, ax = plt.subplots(figsize=(7, 4), dpi=150)

        categories = ['Doctors\nper 1,000', 'Bank branches\nper 100k', 'Internet\npenetration (%)', 'Smartphone\npenetration (%)']
        urban_values = [1.8, 18, 78, 85]
        tier23_values = [0.5, 6, 62, 60]

        x = np.arange(len(categories))
        width = 0.35

        bars1 = ax.bar(x - width/2, urban_values, width, label='Urban',
                      color=(0/255, 102/255, 204/255))
        bars2 = ax.bar(x + width/2, tier23_values, width, label='Tier-2/3',
                      color=(0/255, 153/255, 153/255))

        # Add value labels on bars
        for bar in bars1:
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height + 1,
                   f'{height}', ha='center', va='bottom', fontsize=10)

        for bar in bars2:
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height + 1,
                   f'{height}', ha='center', va='bottom', fontsize=10)

        ax.set_ylabel('Value', fontsize=12)
        ax.set_title('Urban vs Tier-2/3 Infrastructure Gap', fontsize=14, fontweight='bold', pad=20)
        ax.set_xticks(x)
        ax.set_xticklabels(categories, fontsize=10)
        ax.legend(fontsize=10)
        ax.set_ylim(0, max(max(urban_values), max(tier23_values)) * 1.15)
        ax.grid(axis='y', alpha=0.3, linestyle='--')

        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', bbox_inches='tight', facecolor='white', edgecolor='none')
        plt.close()
        buf.seek(0)
        return buf

    def _create_healthcare_financing_pie(self):
        """Create healthcare financing pie chart"""
        fig, ax = plt.subplots(figsize=(5, 4), dpi=150)

        sizes = [62, 30, 8]
        labels = ['Out of Pocket\n(62%)', 'Government\n(30%)', 'Insurance\n(8%)']
        colors = [(244/255, 67/255, 54/255),
                 (0/255, 102/255, 204/255),
                 (76/255, 175/255, 80/255)]

        wedges, texts, autotexts = ax.pie(sizes, labels=labels, colors=colors,
                                          autopct='', startangle=90,
                                          explode=(0.05, 0, 0))

        # Enhance text
        for text in texts:
            text.set_fontsize(11)
            text.set_fontweight('bold')

        ax.set_title('Healthcare Financing Split in India', fontsize=13, fontweight='bold', pad=20)

        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', bbox_inches='tight', facecolor='white', edgecolor='none')
        plt.close()
        buf.seek(0)
        return buf

    def _create_competitive_matrix(self):
        """Create 2x2 competitive positioning matrix"""
        fig, ax = plt.subplots(figsize=(6, 5), dpi=150)

        # Set up the axes
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 10)
        ax.set_xlabel('Geography Focus →\n(Urban to Rural)', fontsize=11)
        ax.set_ylabel('Service Breadth →\n(Narrow to Broad)', fontsize=11)

        # Add grid lines at midpoint
        ax.axhline(y=5, color='gray', linestyle='--', alpha=0.5)
        ax.axvline(x=5, color='gray', linestyle='--', alpha=0.5)

        # Quadrant labels (subtle background)
        quadrants = [
            Rectangle((0, 5), 5, 5, facecolor='#f8f8f8', alpha=0.3),
            Rectangle((5, 5), 5, 5, facecolor='#f0f8ff', alpha=0.3),
            Rectangle((0, 0), 5, 5, facecolor='#fff8f0', alpha=0.3),
            Rectangle((5, 0), 5, 5, facecolor='#f0fff0', alpha=0.3),
        ]
        for q in quadrants:
            ax.add_patch(q)

        # Plot competitors
        competitors = {
            'Practo/Tata Health': (2, 7, (0/255, 102/255, 204/255)),
            '1mg/PharmEasy': (2.5, 3, (255/255, 152/255, 0/255)),
            'eSanjeevani': (7, 6, (76/255, 175/255, 80/255)),
            'MediChain\n(Opportunity)': (7.5, 8, (244/255, 67/255, 54/255)),
        }

        for name, (x, y, color) in competitors.items():
            if 'MediChain' in name:
                # Highlight opportunity gap
                circle = Circle((x, y), 0.5, color=color, alpha=0.3)
                ax.add_patch(circle)
                ax.scatter(x, y, s=200, c=[color], marker='*', edgecolors='black', linewidth=2)
            else:
                ax.scatter(x, y, s=150, c=[color], alpha=0.8, edgecolors='black', linewidth=1)

            ax.annotate(name, (x, y), xytext=(5, 5), textcoords='offset points',
                       fontsize=9, fontweight='bold')

        ax.set_title('Competitive Landscape: Healthcare Platforms', fontsize=13, fontweight='bold', pad=20)
        ax.set_xticks([0, 2.5, 5, 7.5, 10])
        ax.set_xticklabels(['Urban', '', 'Mixed', '', 'Rural'], fontsize=9)
        ax.set_yticks([0, 2.5, 5, 7.5, 10])
        ax.set_yticklabels(['Narrow', '', 'Medium', '', 'Broad'], fontsize=9)

        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', bbox_inches='tight', facecolor='white', edgecolor='none')
        plt.close()
        buf.seek(0)
        return buf

    def _create_conversion_funnel(self):
        """Create conversion funnel diagram"""
        fig, ax = plt.subplots(figsize=(5, 6), dpi=150)

        # Funnel data
        stages = [
            ('Addressable Population', 650000, (0/255, 102/255, 204/255)),
            ('Reach (Awareness)', 75000, (0/255, 153/255, 153/255)),
            ('Active Users', 50000, (76/255, 175/255, 80/255)),
            ('Teleconsults', 20000, (255/255, 152/255, 0/255)),
            ('Paid Subscribers', 5000, (244/255, 67/255, 54/255)),
        ]

        y_pos = 5
        max_width = 8

        for i, (stage, value, color) in enumerate(stages):
            width = max_width * (value / stages[0][1])
            x_pos = (10 - width) / 2

            # Draw trapezoid/rectangle
            if i < len(stages) - 1:
                next_width = max_width * (stages[i+1][1] / stages[0][1])
                next_x = (10 - next_width) / 2

                trapezoid = plt.Polygon([
                    (x_pos, y_pos),
                    (x_pos + width, y_pos),
                    (next_x + next_width, y_pos - 1),
                    (next_x, y_pos - 1)
                ], facecolor=color, alpha=0.7, edgecolor='black', linewidth=1)
                ax.add_patch(trapezoid)
            else:
                rect = Rectangle((x_pos, y_pos - 1), width, 1,
                               facecolor=color, alpha=0.7, edgecolor='black', linewidth=1)
                ax.add_patch(rect)

            # Add text
            ax.text(5, y_pos - 0.5, f"{stage}\n{value:,}",
                   ha='center', va='center', fontsize=10, fontweight='bold', color='white')

            # Conversion rate
            if i > 0:
                conv_rate = (value / stages[i-1][1]) * 100
                ax.text(9.5, y_pos + 0.25, f"{conv_rate:.0f}%",
                       ha='right', va='center', fontsize=9, color='gray')

            y_pos -= 1

        ax.set_xlim(0, 10)
        ax.set_ylim(-1, 6)
        ax.axis('off')
        ax.set_title('Pilot Conversion Funnel', fontsize=13, fontweight='bold', pad=20)

        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', bbox_inches='tight', facecolor='white', edgecolor='none')
        plt.close()
        buf.seek(0)
        return buf

    def _create_roadmap_visual(self):
        """Create 5-year roadmap arrows"""
        fig, ax = plt.subplots(figsize=(10, 3), dpi=150)

        ax.set_xlim(0, 10)
        ax.set_ylim(0, 3)
        ax.axis('off')

        # Timeline arrow
        arrow = FancyArrowPatch((0.5, 1.5), (9.5, 1.5),
                               arrowstyle='->', mutation_scale=20,
                               linewidth=2, color=(0/255, 102/255, 204/255))
        ax.add_patch(arrow)

        # Milestones
        milestones = [
            (2, 'Year 1', '100 kiosks\n50K users', (0/255, 153/255, 153/255)),
            (5, 'Year 3', '25 cities\n1M users', (76/255, 175/255, 80/255)),
            (8, 'Year 5', 'Pan-India\n100M lives', (244/255, 67/255, 54/255)),
        ]

        for x, year, details, color in milestones:
            # Milestone circle
            circle = Circle((x, 1.5), 0.2, color=color, zorder=3)
            ax.add_patch(circle)

            # Year label
            ax.text(x, 2.3, year, ha='center', fontsize=11, fontweight='bold')

            # Details
            ax.text(x, 0.7, details, ha='center', fontsize=9, color='gray')

        ax.set_title('5-Year Scale Roadmap', fontsize=13, fontweight='bold', y=0.95)

        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', bbox_inches='tight', facecolor='white', edgecolor='none')
        plt.close()
        buf.seek(0)
        return buf

    def create_slide1_opportunity(self):
        """Slide 1: Opportunity Landscape"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title and subtitle
        self._add_title(
            slide,
            "Why Tier-2 & Tier-3 India are Ripe for Disruption",
            "India's Tier-2 and Tier-3 cities are no longer peripheral — they are fast becoming the engines of economic growth. "
            "With rising digital penetration and UPI-led inclusion, these markets are digitally ready yet structurally underserved, "
            "creating fertile ground for technology-led disruption."
        )

        # Left section - Macro Trends
        macro_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(5.5), Inches(0.4))
        tf = macro_title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "📊 MACRO TRENDS"
        p.font.name = 'Segoe UI Semibold'
        p.font.size = Pt(16)
        p.font.color.rgb = self.colors['primary_blue']
        p.font.bold = True

        # Macro trend items
        macro_items = [
            ("💰", "Contribute ~45% of India's GDP by 2025"),
            ("👥", "Population base ~650M"),
            ("📱", "Smartphone penetration ~60% (vs 78% urban)"),
            ("💳", "UPI → >12B transactions/month (2025)"),
            ("📡", "Mobile data = cheapest globally (~$0.17/GB)")
        ]

        y_pos = 2.4
        for icon, text in macro_items:
            # Icon
            icon_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos), Inches(0.4), Inches(0.35))
            tf = icon_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = icon
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.CENTER

            # Text
            text_box = slide.shapes.add_textbox(Inches(1.1), Inches(y_pos), Inches(4.8), Inches(0.35))
            tf = text_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = text
            p.font.name = 'Segoe UI'
            p.font.size = Pt(12)
            p.font.color.rgb = self.colors['dark_gray']

            y_pos += 0.45

        # Right section - Underserved Sectors
        sectors_title = slide.shapes.add_textbox(Inches(6.5), Inches(1.9), Inches(6), Inches(0.4))
        tf = sectors_title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "🎯 UNDERSERVED SECTORS"
        p.font.name = 'Segoe UI Semibold'
        p.font.size = Pt(16)
        p.font.color.rgb = self.colors['accent_teal']
        p.font.bold = True

        # Sector items
        sector_items = [
            ("🏥", "Healthcare", "600M underserved; 50+ km avg travel"),
            ("📚", "Education", "Teacher-student ratio 1:60 vs 1:30 norm"),
            ("🏦", "Finance", "190M unbanked; <5% insured"),
            ("🌾", "Agriculture", "Post-harvest losses ~₹90,000 Cr annually")
        ]

        y_pos = 2.4
        for icon, title, desc in sector_items:
            # Container box with light background
            container = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(6.6), Inches(y_pos), Inches(6.2), Inches(0.5)
            )
            container.fill.solid()
            container.fill.fore_color.rgb = self.colors['very_light_gray']
            container.line.fill.background()

            # Icon
            icon_box = slide.shapes.add_textbox(Inches(6.7), Inches(y_pos + 0.05), Inches(0.4), Inches(0.4))
            tf = icon_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = icon
            p.font.size = Pt(16)

            # Title and description
            text_box = slide.shapes.add_textbox(Inches(7.2), Inches(y_pos + 0.05), Inches(5.5), Inches(0.4))
            tf = text_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"{title}: "
            run.font.name = 'Segoe UI Semibold'
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = self.colors['dark_navy']

            run = p.add_run()
            run.text = desc
            run.font.name = 'Segoe UI'
            run.font.size = Pt(10)
            run.font.color.rgb = self.colors['dark_gray']

            y_pos += 0.6

        # Add visuals
        # India map
        map_img = self._create_india_map_visual()
        slide.shapes.add_picture(map_img, Inches(0.5), Inches(4.8), Inches(3), Inches(2))

        # Infrastructure comparison chart
        chart_img = self._create_infrastructure_comparison_chart()
        slide.shapes.add_picture(chart_img, Inches(4), Inches(4.8), Inches(4.5), Inches(2))

        # Bottom banner
        self._add_bottom_banner(slide, "Digital readiness + structural gaps = fertile ground for disruption")
        self._add_footer(slide, 1)

    def create_slide2_healthcare(self):
        """Slide 2: Healthcare Focus"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        self._add_title(
            slide,
            "Healthcare in Tier-2/3 India: Urgent Problem, Large Market"
        )

        # Problem Statement - 3 Pillars
        y_start = 1.2
        pillar_width = 4
        pillar_spacing = 0.2

        pillars = [
            ("🏥", "Accessibility",
             ["75% doctors urban → 600M underserved", "50+ km avg travel"]),
            ("💰", "Affordability",
             ["OOP = 62% of spend", "60M fall into poverty yearly"]),
            ("🤝", "Awareness & Trust",
             ["Preventive stigma", "Reliance on unqualified practitioners"])
        ]

        x_pos = 0.5
        for icon, title, points in pillars:
            # Pillar container
            container = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(y_start), Inches(pillar_width), Inches(1.3)
            )
            container.fill.solid()
            container.fill.fore_color.rgb = self.colors['very_light_gray']
            container.line.color.rgb = self.colors['primary_blue']
            container.line.width = Pt(1)

            # Icon and title
            title_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.1), Inches(y_start + 0.1),
                Inches(pillar_width - 0.2), Inches(0.4)
            )
            tf = title_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = f"{icon} {title}"
            p.font.name = 'Segoe UI Semibold'
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.colors['dark_navy']
            p.alignment = PP_ALIGN.CENTER

            # Points
            points_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.2), Inches(y_start + 0.5),
                Inches(pillar_width - 0.4), Inches(0.7)
            )
            tf = points_box.text_frame
            tf.clear()
            for point in points:
                p = tf.add_paragraph()
                p.text = f"• {point}"
                p.font.name = 'Segoe UI'
                p.font.size = Pt(11)
                p.font.color.rgb = self.colors['dark_gray']
                p.space_before = Pt(2)

            x_pos += pillar_width + pillar_spacing

        # Market Potential section
        market_title = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(6), Inches(0.4))
        tf = market_title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "📈 MARKET POTENTIAL"
        p.font.name = 'Segoe UI Semibold'
        p.font.size = Pt(16)
        p.font.color.rgb = self.colors['success_green']
        p.font.bold = True

        # Market stats in grid
        market_stats = [
            ("Healthcare market", "USD 372B by 2025", "CAGR ~22%"),
            ("Telemedicine", "USD 5.4B by 2025", "Growing rapidly"),
            ("Diagnostics", "CAGR ~20%", "High demand"),
            ("eSanjeevani", "160M+ teleconsults", "Proof of adoption")
        ]

        y_pos = 3.3
        for i in range(0, len(market_stats), 2):
            for j in range(2):
                if i + j < len(market_stats):
                    stat = market_stats[i + j]
                    x = 0.5 + (j * 3.2)

                    stat_box = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(x), Inches(y_pos), Inches(3), Inches(0.7)
                    )
                    stat_box.fill.solid()
                    stat_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
                    stat_box.line.fill.background()

                    text_box = slide.shapes.add_textbox(
                        Inches(x + 0.1), Inches(y_pos + 0.05),
                        Inches(2.8), Inches(0.6)
                    )
                    tf = text_box.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = stat[0] + "\n"
                    run.font.name = 'Segoe UI Semibold'
                    run.font.size = Pt(11)
                    run.font.bold = True

                    run = p.add_run()
                    run.text = stat[1] + " • " + stat[2]
                    run.font.name = 'Segoe UI'
                    run.font.size = Pt(10)
                    run.font.color.rgb = self.colors['medium_gray']

            y_pos += 0.85

        # Add visuals
        # Competitive matrix
        matrix_img = self._create_competitive_matrix()
        slide.shapes.add_picture(matrix_img, Inches(6.8), Inches(2.8), Inches(6), Inches(3.8))

        # Healthcare financing pie
        pie_img = self._create_healthcare_financing_pie()
        slide.shapes.add_picture(pie_img, Inches(0.3), Inches(5.2), Inches(3.5), Inches(1.5))

        # Bottom banner
        self._add_bottom_banner(slide, "Healthcare = burning platform → unmet need + adoption proof + policy push")
        self._add_footer(slide, 2)

    def create_slide3_medichain(self):
        """Slide 3: MediChain Solution"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        self._add_title(
            slide,
            "MediChain — Tech-enabled Primary Care & Diagnostics"
        )

        # Solution Components - 4 horizontal boxes
        components = [
            ("🤖", "AI Symptom Triage", "Vernacular chatbot\nTriage <₹20"),
            ("🔬", "IoT Diagnostic Kiosks", "BP, sugar, ECG, SPO₂\nCost ~₹1L per kiosk"),
            ("🔗", "Blockchain Health Records", "NDHM aligned\nSecure, portable"),
            ("💊", "Phygital Linkages", "Local pharmacies\nLast-mile medicine")
        ]

        x_pos = 0.5
        component_width = 3.1
        for icon, title, desc in components:
            # Component box
            comp_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(1.2), Inches(component_width), Inches(1.2)
            )
            comp_box.fill.solid()
            comp_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
            comp_box.line.color.rgb = self.colors['primary_blue']
            comp_box.line.width = Pt(2)

            # Icon
            icon_box = slide.shapes.add_textbox(
                Inches(x_pos + component_width/2 - 0.3), Inches(1.3),
                Inches(0.6), Inches(0.4)
            )
            tf = icon_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = icon
            p.font.size = Pt(24)
            p.alignment = PP_ALIGN.CENTER

            # Title and description
            text_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.1), Inches(1.7),
                Inches(component_width - 0.2), Inches(0.4)
            )
            tf = text_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = 'Segoe UI Semibold'
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.colors['dark_navy']
            p.alignment = PP_ALIGN.CENTER

            p = tf.add_paragraph()
            p.text = desc
            p.font.name = 'Segoe UI'
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors['medium_gray']
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(2)

            x_pos += component_width + 0.15

        # Differentiators section
        diff_title = slide.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(5), Inches(0.4))
        tf = diff_title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "⭐ KEY DIFFERENTIATORS"
        p.font.name = 'Segoe UI Semibold'
        p.font.size = Pt(14)
        p.font.color.rgb = self.colors['accent_teal']
        p.font.bold = True

        # Differentiator points
        diff_points = [
            "🗣️ Vernacular-first UX for regional adoption",
            "💰 Affordable: <₹100 consults, ₹499/year family plan",
            "🤝 Trust: Kiosk placement in pharmacies + NGO/state tie-ups"
        ]

        y_pos = 3.0
        for point in diff_points:
            point_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos), Inches(5), Inches(0.35))
            tf = point_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = point
            p.font.name = 'Segoe UI'
            p.font.size = Pt(11)
            p.font.color.rgb = self.colors['dark_gray']
            y_pos += 0.4

        # Impact & Scale section
        impact_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(5), Inches(0.4))
        tf = impact_title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "🎯 IMPACT & SCALE"
        p.font.name = 'Segoe UI Semibold'
        p.font.size = Pt(14)
        p.font.color.rgb = self.colors['success_green']
        p.font.bold = True

        # Impact points in two columns
        impact_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(4.95), Inches(5.5), Inches(1.8)
        )
        impact_box.fill.solid()
        impact_box.fill.fore_color.rgb = RGBColor(245, 255, 245)
        impact_box.line.fill.background()

        impact_text = slide.shapes.add_textbox(Inches(0.7), Inches(5.05), Inches(5.2), Inches(1.6))
        tf = impact_text.text_frame
        tf.clear()

        p = tf.add_paragraph()
        p.text = "Economic Impact:"
        p.font.name = 'Segoe UI Semibold'
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.colors['dark_navy']

        economic_points = [
            "• Low OOP burden for families",
            "• Scalable subscription revenue",
            "• Job creation in rural areas"
        ]

        for point in economic_points:
            p = tf.add_paragraph()
            p.text = point
            p.font.name = 'Segoe UI'
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors['dark_gray']
            p.space_before = Pt(1)

        p = tf.add_paragraph()
        p.text = "\nSocial Impact:"
        p.font.name = 'Segoe UI Semibold'
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.colors['dark_navy']
        p.space_before = Pt(4)

        social_points = [
            "• Access for 100M+ underserved",
            "• Preventive care culture",
            "• SDG-3 alignment"
        ]

        for point in social_points:
            p = tf.add_paragraph()
            p.text = point
            p.font.name = 'Segoe UI'
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors['dark_gray']
            p.space_before = Pt(1)

        # Add visuals
        # Conversion funnel
        funnel_img = self._create_conversion_funnel()
        slide.shapes.add_picture(funnel_img, Inches(6.2), Inches(2.6), Inches(3.5), Inches(4))

        # Roadmap
        roadmap_img = self._create_roadmap_visual()
        slide.shapes.add_picture(roadmap_img, Inches(9.8), Inches(3.8), Inches(3.3), Inches(2.8))

        # Bottom banner
        self._add_bottom_banner(slide, "MediChain = Vernacular, trust-first, affordable healthcare pathway for Bharat")
        self._add_footer(slide, 3)

    def generate_presentation(self):
        """Generate the complete presentation"""
        self.create_slide1_opportunity()
        self.create_slide2_healthcare()
        self.create_slide3_medichain()

        # Save to PPT Generated folder
        output_path = "/mnt/e/AI and Projects/Case Comp PPT/PPT Generated/Tier23_Healthcare_Disruption.pptx"
        self.prs.save(output_path)
        print(f"Presentation saved to: {output_path}")
        return output_path

if __name__ == "__main__":
    creator = HealthcarePresentationCreator()
    output_file = creator.generate_presentation()
    print(f"\n✅ Tier-2/3 Healthcare Disruption presentation created successfully!")
    print(f"📍 Location: {output_file}")