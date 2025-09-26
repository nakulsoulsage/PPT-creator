from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import numpy as np
from io import BytesIO
import seaborn as sns

class MediChainRichVisualPresentation:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)

        # Rich color palette
        self.colors = {
            'primary': RGBColor(0, 150, 170),     # Medical Teal
            'secondary': RGBColor(0, 83, 159),    # Deep Blue
            'accent1': RGBColor(0, 180, 130),     # Health Green
            'accent2': RGBColor(255, 140, 0),     # Orange
            'accent3': RGBColor(237, 28, 36),     # Alert Red
            'dark': RGBColor(45, 45, 45),         # Dark Gray
            'light': RGBColor(245, 245, 245),     # Light Gray
            'white': RGBColor(255, 255, 255)      # White
        }

    def add_footer(self, slide):
        """Add footer to all slides"""
        footer = slide.shapes.add_textbox(
            Inches(0.5), Inches(8.3), Inches(15), Inches(0.5)
        )
        footer.text_frame.text = "Presented By — Nakul Nandanwar, Vaishnavi Bhangale, Rahul Kumbhare"
        footer.text_frame.paragraphs[0].font.size = Pt(10)
        footer.text_frame.paragraphs[0].font.color.rgb = self.colors['dark']
        footer.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        footer.text_frame.paragraphs[0].font.italic = True

    def create_slide1_opportunity(self):
        """Slide 1: Opportunity Landscape"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        # Title with gradient background
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Inches(1.2)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = self.colors['primary']
        title_bg.line.fill.background()

        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(0.8))
        title.text_frame.text = "Why Tier-2 & Tier-3 India are Ripe for Disruption"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['white']
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Narrative
        narrative = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(15), Inches(0.6))
        narrative.text_frame.text = "Tier-2 and Tier-3 cities are the new economic frontier.\nWith rising digital penetration but weak physical infrastructure, these markets are digitally ready yet structurally underserved."
        narrative.text_frame.paragraphs[0].font.size = Pt(14)
        narrative.text_frame.paragraphs[0].font.color.rgb = self.colors['dark']
        narrative.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Macro Trends Section
        trends_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(2.1), Inches(7.7), Inches(3)
        )
        trends_bg.fill.solid()
        trends_bg.fill.fore_color.rgb = self.colors['light']

        trends_title = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(7.3), Inches(0.4))
        trends_title.text_frame.text = "📊 MACRO TRENDS"
        trends_title.text_frame.paragraphs[0].font.size = Pt(18)
        trends_title.text_frame.paragraphs[0].font.bold = True
        trends_title.text_frame.paragraphs[0].font.color.rgb = self.colors['secondary']

        # Trend boxes
        trend_data = [
            ("💰", "45% of GDP", "by 2025"),
            ("👥", "650M", "Population"),
            ("📱", "60%", "Smartphone"),
            ("💳", "12B+", "UPI/month"),
            ("📡", "$0.17/GB", "Data Cost")
        ]

        x_start = 0.5
        for i, (icon, value, label) in enumerate(trend_data):
            x = x_start + (i % 3) * 2.5
            y = 2.7 + (i // 3) * 1.1

            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.3), Inches(0.9)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.colors['white']

            text = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.1), Inches(2.1), Inches(0.7))
            text.text_frame.text = f"{icon} {value}\n{label}"
            text.text_frame.paragraphs[0].font.size = Pt(12)
            text.text_frame.paragraphs[0].font.bold = True
            text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Underserved Sectors
        sectors_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(2.1), Inches(7.5), Inches(3)
        )
        sectors_bg.fill.solid()
        sectors_bg.fill.fore_color.rgb = self.colors['accent1']
        sectors_bg.fill.transparency = 0.9

        sectors_title = slide.shapes.add_textbox(Inches(8.4), Inches(2.2), Inches(7.1), Inches(0.4))
        sectors_title.text_frame.text = "🚨 UNDERSERVED SECTORS"
        sectors_title.text_frame.paragraphs[0].font.size = Pt(18)
        sectors_title.text_frame.paragraphs[0].font.bold = True
        sectors_title.text_frame.paragraphs[0].font.color.rgb = self.colors['secondary']

        sector_data = [
            ("🏥", "Healthcare", "600M underserved"),
            ("📚", "Education", "Ratio 1:60 vs 1:30"),
            ("💰", "Finance", "190M unbanked"),
            ("🌾", "Agriculture", "₹90,000 Cr losses")
        ]

        for i, (icon, sector, stat) in enumerate(sector_data):
            y = 2.7 + i * 0.55

            text = slide.shapes.add_textbox(Inches(8.5), Inches(y), Inches(7), Inches(0.5))
            text.text_frame.text = f"{icon} {sector}: {stat}"
            text.text_frame.paragraphs[0].font.size = Pt(13)
            text.text_frame.paragraphs[0].font.color.rgb = self.colors['dark']

        # Visual chart - Urban vs Tier-2/3 Access
        fig, ax = plt.subplots(figsize=(7, 2.5))
        categories = ['Doctors', 'Banks', 'Internet', 'Smartphones']
        urban = [85, 78, 78, 78]
        tier23 = [25, 35, 60, 60]

        x = np.arange(len(categories))
        width = 0.35

        bars1 = ax.bar(x - width/2, urban, width, label='Urban', color='#00539F')
        bars2 = ax.bar(x + width/2, tier23, width, label='Tier-2/3', color='#0096AA')

        ax.set_xlabel('Access Categories', fontsize=10)
        ax.set_ylabel('Penetration %', fontsize=10)
        ax.set_title('Urban vs Tier-2/3 Access Gap', fontsize=12, fontweight='bold')
        ax.set_xticks(x)
        ax.set_xticklabels(categories)
        ax.legend()
        ax.grid(True, alpha=0.3)

        for bars in [bars1, bars2]:
            for bar in bars:
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2., height,
                       f'{height}%', ha='center', va='bottom', fontsize=9)

        img_stream = BytesIO()
        plt.tight_layout()
        plt.savefig(img_stream, format='PNG', dpi=150, bbox_inches='tight')
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, Inches(4), Inches(5.3), height=Inches(2.5))
        plt.close()

        # Bottom banner
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(8), Inches(16), Inches(0.3)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = self.colors['secondary']

        banner_text = slide.shapes.add_textbox(Inches(0.5), Inches(7.95), Inches(15), Inches(0.4))
        banner_text.text_frame.text = "Digital readiness + structural gaps = fertile ground for disruption."
        banner_text.text_frame.paragraphs[0].font.size = Pt(14)
        banner_text.text_frame.paragraphs[0].font.bold = True
        banner_text.text_frame.paragraphs[0].font.color.rgb = self.colors['white']
        banner_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        self.add_footer(slide)

    def create_slide2_healthcare(self):
        """Slide 2: Healthcare Focus"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        # Title
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Inches(1)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = self.colors['accent3']
        title_bg.line.fill.background()

        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(0.6))
        title.text_frame.text = "Why Healthcare is the Burning Platform"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['white']
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Narrative
        narrative = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(15), Inches(0.4))
        narrative.text_frame.text = "Among all underserved sectors, healthcare stands out due to its scale, urgency, and proven readiness for digital disruption."
        narrative.text_frame.paragraphs[0].font.size = Pt(14)
        narrative.text_frame.paragraphs[0].font.color.rgb = self.colors['dark']
        narrative.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Three Pillars - Key Barriers
        pillars = [
            ("🚫 ACCESSIBILITY", "75% doctors in urban\n600M underserved\nAvg 50+ km travel", self.colors['secondary']),
            ("💸 AFFORDABILITY", "OOP = 62% spend\n60M fall into poverty\nyearly from healthcare", self.colors['accent2']),
            ("🧠 AWARENESS & TRUST", "Stigma around preventive\n& mental care\nReliance on quacks", self.colors['accent1'])
        ]

        for i, (title_text, content, color) in enumerate(pillars):
            x = 0.5 + i * 5.2

            # Pillar shape
            pillar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.7), Inches(5), Inches(2.2)
            )
            pillar.fill.solid()
            pillar.fill.fore_color.rgb = color
            pillar.fill.transparency = 0.9

            # Pillar title
            pillar_title = slide.shapes.add_textbox(Inches(x + 0.1), Inches(1.8), Inches(4.8), Inches(0.5))
            pillar_title.text_frame.text = title_text
            pillar_title.text_frame.paragraphs[0].font.size = Pt(16)
            pillar_title.text_frame.paragraphs[0].font.bold = True
            pillar_title.text_frame.paragraphs[0].font.color.rgb = color
            pillar_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Pillar content
            pillar_content = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.4), Inches(4.6), Inches(1.3))
            pillar_content.text_frame.text = content
            pillar_content.text_frame.paragraphs[0].font.size = Pt(12)
            pillar_content.text_frame.paragraphs[0].font.color.rgb = self.colors['dark']
            pillar_content.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Data Highlights
        highlights = [
            ("Healthcare Market", "$372B by 2025", "22% CAGR"),
            ("Telemedicine", "$5.4B by 2025", "45% CAGR"),
            ("Diagnostics", "20% CAGR", "Growing rapidly"),
            ("eSanjeevani", "160M+ teleconsults", "Govt success")
        ]

        for i, (label, value, subtext) in enumerate(highlights):
            x = 1 + (i % 2) * 7.5
            y = 4.2 + (i // 2) * 1.2

            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(7), Inches(1)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.colors['light']

            text = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.1), Inches(6.6), Inches(0.8))
            text.text_frame.text = f"📊 {label}: {value}\n{subtext}"
            text.text_frame.paragraphs[0].font.size = Pt(13)
            text.text_frame.paragraphs[0].font.bold = True
            text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Pie Chart - Healthcare Spend Split
        fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 3))

        # Pie chart
        sizes = [62, 30, 8]
        labels = ['Out-of-Pocket (62%)', 'Government (30%)', 'Insurance (8%)']
        colors = ['#ED1C24', '#0096AA', '#00B482']
        explode = (0.1, 0, 0)

        ax1.pie(sizes, explode=explode, labels=labels, colors=colors, autopct='%1.0f%%',
                shadow=True, startangle=90)
        ax1.set_title('Healthcare Spending Split', fontsize=12, fontweight='bold')

        # Bar chart - Doctor Density
        categories = ['Urban', 'Tier-2', 'Tier-3', 'Rural']
        doctors = [1.2, 0.5, 0.3, 0.2]
        colors_bar = ['#00539F', '#0096AA', '#00B482', '#FF8C00']

        bars = ax2.bar(categories, doctors, color=colors_bar)
        ax2.set_ylabel('Doctors per 1000 people', fontsize=10)
        ax2.set_title('Doctor Density by Region', fontsize=12, fontweight='bold')
        ax2.grid(True, alpha=0.3, axis='y')

        for bar in bars:
            height = bar.get_height()
            ax2.text(bar.get_x() + bar.get_width()/2., height,
                    f'{height}', ha='center', va='bottom', fontsize=10)

        plt.tight_layout()
        img_stream = BytesIO()
        plt.savefig(img_stream, format='PNG', dpi=150, bbox_inches='tight')
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, Inches(3), Inches(6.5), height=Inches(1.3))
        plt.close()

        # Bottom banner
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(8), Inches(16), Inches(0.3)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = self.colors['accent3']

        banner_text = slide.shapes.add_textbox(Inches(0.5), Inches(7.95), Inches(15), Inches(0.4))
        banner_text.text_frame.text = "Healthcare = urgent problem + massive market potential."
        banner_text.text_frame.paragraphs[0].font.size = Pt(14)
        banner_text.text_frame.paragraphs[0].font.bold = True
        banner_text.text_frame.paragraphs[0].font.color.rgb = self.colors['white']
        banner_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        self.add_footer(slide)

    def create_slide3_competition(self):
        """Slide 3: Market Insights & Competition"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        # Title
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(0.6))
        title.text_frame.text = "Competitive Landscape & White Space"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['secondary']
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Narrative
        narrative = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(15), Inches(0.4))
        narrative.text_frame.text = "Existing players have not solved the affordability + accessibility gap in Tier-2/3 markets, leaving a white space for a vernacular, trust-based model."
        narrative.text_frame.paragraphs[0].font.size = Pt(13)
        narrative.text_frame.paragraphs[0].font.color.rgb = self.colors['dark']
        narrative.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 2x2 Competitive Matrix
        matrix_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(5)
        )
        matrix_bg.fill.solid()
        matrix_bg.fill.fore_color.rgb = self.colors['light']

        # Draw matrix lines
        h_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4), Inches(9), Inches(0.02)
        )
        h_line.fill.solid()
        h_line.fill.fore_color.rgb = self.colors['dark']

        v_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(5), Inches(1.5), Inches(0.02), Inches(5)
        )
        v_line.fill.solid()
        v_line.fill.fore_color.rgb = self.colors['dark']

        # Axis labels
        x_label = slide.shapes.add_textbox(Inches(4.5), Inches(6.6), Inches(1), Inches(0.3))
        x_label.text_frame.text = "Geography →"
        x_label.text_frame.paragraphs[0].font.size = Pt(11)
        x_label.text_frame.paragraphs[0].font.bold = True

        y_label = slide.shapes.add_textbox(Inches(0.1), Inches(3.8), Inches(0.4), Inches(0.4))
        y_label.text_frame.text = "Service Breadth →"
        y_label.text_frame.paragraphs[0].font.size = Pt(11)
        y_label.text_frame.paragraphs[0].font.bold = True

        # Quadrant labels
        quadrants = [
            (Inches(1), Inches(1.7), "Urban + Narrow", ["Practo", "Apollo 24/7"]),
            (Inches(5.5), Inches(1.7), "Rural + Narrow", ["Local clinics"]),
            (Inches(1), Inches(4.3), "Urban + Broad", ["1mg", "PharmEasy", "Tata Health"]),
            (Inches(5.5), Inches(4.3), "Rural + Broad", ["eSanjeevani", "WHITE SPACE"])
        ]

        for x, y, label, companies in quadrants:
            # Label
            label_box = slide.shapes.add_textbox(x, y, Inches(3.5), Inches(0.3))
            label_box.text_frame.text = label
            label_box.text_frame.paragraphs[0].font.size = Pt(10)
            label_box.text_frame.paragraphs[0].font.bold = True
            label_box.text_frame.paragraphs[0].font.color.rgb = self.colors['accent1']

            # Companies
            for i, company in enumerate(companies):
                if "WHITE SPACE" in company:
                    comp_bg = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.2, y + 0.4 + i*0.5, Inches(3), Inches(0.4)
                    )
                    comp_bg.fill.solid()
                    comp_bg.fill.fore_color.rgb = self.colors['accent2']

                comp_box = slide.shapes.add_textbox(x + 0.3, y + 0.4 + i*0.5, Inches(2.8), Inches(0.4))
                comp_box.text_frame.text = company
                comp_box.text_frame.paragraphs[0].font.size = Pt(11)
                if "WHITE SPACE" in company:
                    comp_box.text_frame.paragraphs[0].font.bold = True
                    comp_box.text_frame.paragraphs[0].font.color.rgb = self.colors['white']

        # Insights Section
        insights_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10), Inches(1.5), Inches(5.5), Inches(5)
        )
        insights_bg.fill.solid()
        insights_bg.fill.fore_color.rgb = self.colors['accent1']
        insights_bg.fill.transparency = 0.95

        insights_title = slide.shapes.add_textbox(Inches(10.2), Inches(1.6), Inches(5.1), Inches(0.4))
        insights_title.text_frame.text = "🔍 KEY INSIGHTS"
        insights_title.text_frame.paragraphs[0].font.size = Pt(16)
        insights_title.text_frame.paragraphs[0].font.bold = True
        insights_title.text_frame.paragraphs[0].font.color.rgb = self.colors['secondary']

        insights = [
            "• India spends only 3% of GDP on healthcare vs 9% global avg",
            "• 60% of rural Indians rely on informal practitioners",
            "• Telemedicine adoption surged 10x post-COVID",
            "• Digital health market CAGR of 39% till 2027",
            "• 75% patients willing to use digital health if affordable"
        ]

        for i, insight in enumerate(insights):
            insight_box = slide.shapes.add_textbox(Inches(10.3), Inches(2.2 + i*0.8), Inches(5.1), Inches(0.7))
            insight_box.text_frame.text = insight
            insight_box.text_frame.paragraphs[0].font.size = Pt(11)
            insight_box.text_frame.word_wrap = True

        # White Space Highlight
        white_space = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(7), Inches(12), Inches(0.8)
        )
        white_space.fill.solid()
        white_space.fill.fore_color.rgb = self.colors['accent2']
        white_space.fill.transparency = 0.1

        white_text = slide.shapes.add_textbox(Inches(2.2), Inches(7.1), Inches(11.6), Inches(0.6))
        white_text.text_frame.text = "🎯 WHITE SPACE = Affordable vernacular model with physical touchpoints for Tier-2/3"
        white_text.text_frame.paragraphs[0].font.size = Pt(14)
        white_text.text_frame.paragraphs[0].font.bold = True
        white_text.text_frame.paragraphs[0].font.color.rgb = self.colors['white']
        white_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Bottom banner
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(8), Inches(16), Inches(0.3)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = self.colors['secondary']

        banner_text = slide.shapes.add_textbox(Inches(0.5), Inches(7.95), Inches(15), Inches(0.4))
        banner_text.text_frame.text = "White space = scalable, affordable healthcare model for Tier-2/3 India."
        banner_text.text_frame.paragraphs[0].font.size = Pt(14)
        banner_text.text_frame.paragraphs[0].font.bold = True
        banner_text.text_frame.paragraphs[0].font.color.rgb = self.colors['white']
        banner_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        self.add_footer(slide)

    def create_slide4_solution(self):
        """Slide 4: MediChain Solution"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        # Title with gradient
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Inches(1)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = self.colors['accent1']

        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(0.6))
        title.text_frame.text = "MediChain — Tech-enabled Primary Care & Diagnostics"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['white']
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Solution Components - 4 pillars
        components = [
            ("🤖", "AI Symptom Triage", "Vernacular chatbot\nTriage <₹20", self.colors['secondary']),
            ("📊", "IoT Diagnostic Kiosks", "BP, ECG, SPO₂, sugar\n~₹1L per kiosk", self.colors['primary']),
            ("🔐", "Blockchain Records", "NDHM aligned\nSecure, portable", self.colors['accent1']),
            ("🏪", "Phygital Linkages", "Local pharmacies\nLast-mile delivery", self.colors['accent2'])
        ]

        for i, (icon, title_text, desc, color) in enumerate(components):
            x = 0.5 + (i % 2) * 8
            y = 1.3 + (i // 2) * 2.5

            # Component box
            comp_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(7.5), Inches(2.2)
            )
            comp_bg.fill.solid()
            comp_bg.fill.fore_color.rgb = color
            comp_bg.fill.transparency = 0.92

            # Icon circle
            icon_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(x + 0.3), Inches(y + 0.3), Inches(1.5), Inches(1.5)
            )
            icon_circle.fill.solid()
            icon_circle.fill.fore_color.rgb = color

            icon_text = slide.shapes.add_textbox(Inches(x + 0.5), Inches(y + 0.6), Inches(1.1), Inches(0.9))
            icon_text.text_frame.text = icon
            icon_text.text_frame.paragraphs[0].font.size = Pt(36)
            icon_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Component text
            comp_text = slide.shapes.add_textbox(Inches(x + 2), Inches(y + 0.4), Inches(5.2), Inches(1.4))
            comp_text.text_frame.text = f"{title_text}\n\n{desc}"
            p1 = comp_text.text_frame.paragraphs[0]
            p1.font.size = Pt(16)
            p1.font.bold = True
            p1.font.color.rgb = color

            if len(comp_text.text_frame.paragraphs) > 1:
                p2 = comp_text.text_frame.paragraphs[1]
                p2.font.size = Pt(12)
                p2.font.color.rgb = self.colors['dark']

        # Differentiators
        diff_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.2), Inches(7.2), Inches(1.5)
        )
        diff_bg.fill.solid()
        diff_bg.fill.fore_color.rgb = self.colors['light']

        diff_title = slide.shapes.add_textbox(Inches(0.7), Inches(6.3), Inches(7), Inches(0.3))
        diff_title.text_frame.text = "🎯 KEY DIFFERENTIATORS"
        diff_title.text_frame.paragraphs[0].font.size = Pt(14)
        diff_title.text_frame.paragraphs[0].font.bold = True
        diff_title.text_frame.paragraphs[0].font.color.rgb = self.colors['secondary']

        diffs = [
            "✓ Vernacular-first UX for mass adoption",
            "✓ Affordable: <₹100 consults, ₹499/year family plan",
            "✓ Trust via pharmacy kiosks + NGO/govt tie-ups"
        ]

        for i, diff in enumerate(diffs):
            diff_text = slide.shapes.add_textbox(Inches(0.8), Inches(6.7 + i*0.3), Inches(6.8), Inches(0.3))
            diff_text.text_frame.text = diff
            diff_text.text_frame.paragraphs[0].font.size = Pt(11)

        # Funnel Diagram
        fig, ax = plt.subplots(figsize=(6, 3))

        stages = ['Total Market\n650K', 'Addressable\n75K', 'Engaged\n50K', 'Active\n20K', 'Subscribers\n5K']
        values = [650, 75, 50, 20, 5]
        colors_funnel = ['#00539F', '#0096AA', '#00B482', '#FF8C00', '#ED1C24']

        y_pos = np.arange(len(stages))

        for i, (stage, value, color) in enumerate(zip(stages, values, colors_funnel)):
            width = value / values[0]
            left = (1 - width) / 2
            ax.barh(i, width, left=left, color=color, height=0.8)
            ax.text(0.5, i, stage, ha='center', va='center', fontsize=10, fontweight='bold', color='white')

        ax.set_xlim(0, 1)
        ax.set_ylim(-0.5, len(stages) - 0.5)
        ax.axis('off')
        ax.set_title('User Acquisition Funnel', fontsize=12, fontweight='bold', pad=20)

        img_stream = BytesIO()
        plt.tight_layout()
        plt.savefig(img_stream, format='PNG', dpi=150, bbox_inches='tight', facecolor='white')
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, Inches(8.3), Inches(6), height=Inches(1.7))
        plt.close()

        # Bottom banner
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(8), Inches(16), Inches(0.3)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = self.colors['accent1']

        banner_text = slide.shapes.add_textbox(Inches(0.5), Inches(7.95), Inches(15), Inches(0.4))
        banner_text.text_frame.text = "MediChain = vernacular, trust-first, affordable healthcare disruption."
        banner_text.text_frame.paragraphs[0].font.size = Pt(14)
        banner_text.text_frame.paragraphs[0].font.bold = True
        banner_text.text_frame.paragraphs[0].font.color.rgb = self.colors['white']
        banner_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        self.add_footer(slide)

    def create_slide5_impact(self):
        """Slide 5: Impact & Roadmap"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        # Title
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(0.6))
        title.text_frame.text = "Scalable Impact Pathway"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['secondary']
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Narrative
        narrative = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(15), Inches(0.4))
        narrative.text_frame.text = "MediChain balances profitability with social good, delivering both sustainable returns and healthcare access at scale."
        narrative.text_frame.paragraphs[0].font.size = Pt(13)
        narrative.text_frame.paragraphs[0].font.color.rgb = self.colors['dark']
        narrative.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Impact Sections - Two columns
        # Economic Impact
        econ_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(7.5), Inches(2.8)
        )
        econ_bg.fill.solid()
        econ_bg.fill.fore_color.rgb = self.colors['secondary']
        econ_bg.fill.transparency = 0.95

        econ_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(7.1), Inches(0.4))
        econ_title.text_frame.text = "💰 ECONOMIC IMPACT"
        econ_title.text_frame.paragraphs[0].font.size = Pt(16)
        econ_title.text_frame.paragraphs[0].font.bold = True
        econ_title.text_frame.paragraphs[0].font.color.rgb = self.colors['secondary']

        econ_points = [
            "• Subscription revenues: ₹499/year per family",
            "• Lower out-of-pocket burden for patients",
            "• Scalable unit economics with kiosk model",
            "• Break-even in 18 months per district",
            "• 35% EBITDA margins at scale"
        ]

        for i, point in enumerate(econ_points):
            point_text = slide.shapes.add_textbox(Inches(0.8), Inches(2.1 + i*0.4), Inches(7), Inches(0.4))
            point_text.text_frame.text = point
            point_text.text_frame.paragraphs[0].font.size = Pt(11)

        # Social Impact
        social_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.5), Inches(7), Inches(2.8)
        )
        social_bg.fill.solid()
        social_bg.fill.fore_color.rgb = self.colors['accent1']
        social_bg.fill.transparency = 0.95

        social_title = slide.shapes.add_textbox(Inches(8.7), Inches(1.6), Inches(6.6), Inches(0.4))
        social_title.text_frame.text = "🌍 SOCIAL IMPACT"
        social_title.text_frame.paragraphs[0].font.size = Pt(16)
        social_title.text_frame.paragraphs[0].font.bold = True
        social_title.text_frame.paragraphs[0].font.color.rgb = self.colors['accent1']

        social_points = [
            "• 100M+ underserved gain access by Y5",
            "• Preventive health adoption increases 5x",
            "• Mental health stigma reduction",
            "• Women's health focus (50% users)",
            "• SDG-3 alignment: Good Health for All"
        ]

        for i, point in enumerate(social_points):
            point_text = slide.shapes.add_textbox(Inches(8.8), Inches(2.1 + i*0.4), Inches(6.5), Inches(0.4))
            point_text.text_frame.text = point
            point_text.text_frame.paragraphs[0].font.size = Pt(11)

        # 5-Year Roadmap Timeline
        roadmap_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(15), Inches(3)
        )
        roadmap_bg.fill.solid()
        roadmap_bg.fill.fore_color.rgb = self.colors['light']

        roadmap_title = slide.shapes.add_textbox(Inches(0.7), Inches(4.6), Inches(14.6), Inches(0.4))
        roadmap_title.text_frame.text = "📈 5-YEAR GROWTH ROADMAP"
        roadmap_title.text_frame.paragraphs[0].font.size = Pt(16)
        roadmap_title.text_frame.paragraphs[0].font.bold = True
        roadmap_title.text_frame.paragraphs[0].font.color.rgb = self.colors['secondary']

        # Timeline milestones
        milestones = [
            ("Y1", "PILOT", "100 kiosks\n50K users\n5 districts", self.colors['secondary']),
            ("Y2", "EXPAND", "500 kiosks\n300K users\n3 states", self.colors['primary']),
            ("Y3", "SCALE", "2500 kiosks\n1M users\n25 cities", self.colors['accent1']),
            ("Y4", "GROWTH", "10K kiosks\n10M users\n100 cities", self.colors['accent2']),
            ("Y5", "NATIONAL", "50K kiosks\n100M users\nPan-India", self.colors['accent3'])
        ]

        for i, (year, phase, details, color) in enumerate(milestones):
            x = 1 + i * 2.8

            # Year circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(x), Inches(5.1), Inches(0.8), Inches(0.8)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = color

            year_text = slide.shapes.add_textbox(Inches(x + 0.15), Inches(5.25), Inches(0.5), Inches(0.5))
            year_text.text_frame.text = year
            year_text.text_frame.paragraphs[0].font.size = Pt(14)
            year_text.text_frame.paragraphs[0].font.bold = True
            year_text.text_frame.paragraphs[0].font.color.rgb = self.colors['white']
            year_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Phase and details
            phase_text = slide.shapes.add_textbox(Inches(x - 0.3), Inches(6), Inches(1.4), Inches(1.2))
            phase_text.text_frame.text = f"{phase}\n{details}"
            phase_text.text_frame.paragraphs[0].font.size = Pt(10)
            phase_text.text_frame.paragraphs[0].font.bold = True
            phase_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Connect with arrow
            if i < len(milestones) - 1:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW, Inches(x + 0.9), Inches(5.35), Inches(1.8), Inches(0.3)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.colors['dark']
                arrow.fill.transparency = 0.7

        # Bottom banner
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(8), Inches(16), Inches(0.3)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = self.colors['primary']

        banner_text = slide.shapes.add_textbox(Inches(0.5), Inches(7.95), Inches(15), Inches(0.4))
        banner_text.text_frame.text = "Scalable, sustainable, socially impactful disruption for Bharat."
        banner_text.text_frame.paragraphs[0].font.size = Pt(14)
        banner_text.text_frame.paragraphs[0].font.bold = True
        banner_text.text_frame.paragraphs[0].font.color.rgb = self.colors['white']
        banner_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        self.add_footer(slide)

    def generate_presentation(self):
        """Generate all slides and save presentation"""
        print("Creating MediChain Rich Visual Presentation...")

        self.create_slide1_opportunity()
        print("✓ Slide 1: Opportunity Landscape created")

        self.create_slide2_healthcare()
        print("✓ Slide 2: Healthcare Focus created")

        self.create_slide3_competition()
        print("✓ Slide 3: Market Insights & Competition created")

        self.create_slide4_solution()
        print("✓ Slide 4: MediChain Solution created")

        self.create_slide5_impact()
        print("✓ Slide 5: Impact & Roadmap created")

        # Save presentation
        filename = "PPT Generated/MediChain_Rich_Visual_5Slides.pptx"
        self.prs.save(filename)
        print(f"\n✅ Presentation saved as '{filename}'")
        print(f"📊 Total slides: 5")
        print("🎨 Style: Rich Visual with Maximum Infographics")

        return filename

if __name__ == "__main__":
    creator = MediChainRichVisualPresentation()
    creator.generate_presentation()