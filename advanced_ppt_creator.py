#!/usr/bin/env python3
"""
Advanced Case Competition PowerPoint Creator
Complete system for creating professional presentations with charts, infographics, and tables
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd
import numpy as np
from datetime import datetime
import io
from PIL import Image

class CaseCompetitionPPT:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Professional color scheme (McKinsey style)
        self.colors = {
            'primary': RGBColor(0, 51, 141),      # Dark Blue
            'secondary': RGBColor(0, 114, 206),   # Light Blue
            'accent': RGBColor(0, 176, 240),      # Sky Blue
            'success': RGBColor(0, 176, 80),      # Green
            'warning': RGBColor(255, 192, 0),     # Yellow
            'danger': RGBColor(255, 0, 0),        # Red
            'dark': RGBColor(64, 64, 64),         # Dark Gray
            'light': RGBColor(191, 191, 191),     # Light Gray
            'white': RGBColor(255, 255, 255),     # White
            'black': RGBColor(0, 0, 0)            # Black
        }
        
    def add_title_slide(self, title, subtitle, team_members):
        """Create professional title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Blank layout
        
        # Add background shape
        left = Inches(0)
        top = Inches(0)
        width = self.prs.slide_width
        height = self.prs.slide_height
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['primary']
        shape.line.fill.background()
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['white']
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = self.colors['accent']
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add team members
        team_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1.5))
        team_frame = team_box.text_frame
        team_frame.text = "Team Members"
        p = team_frame.add_paragraph()
        p.text = " | ".join(team_members)
        team_frame.paragraphs[0].font.size = Pt(16)
        team_frame.paragraphs[0].font.bold = True
        team_frame.paragraphs[0].font.color.rgb = self.colors['accent']
        p.font.size = Pt(14)
        p.font.color.rgb = self.colors['white']
        team_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        p.alignment = PP_ALIGN.CENTER
        
        # Add date
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        date_frame = date_box.text_frame
        date_frame.text = datetime.now().strftime("%B %Y")
        date_frame.paragraphs[0].font.size = Pt(12)
        date_frame.paragraphs[0].font.color.rgb = self.colors['light']
        date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_executive_summary(self, problem, solution, impact, timeline):
        """Create executive summary slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = "Executive Summary"
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # Create four quadrants
        quadrants = [
            {"title": "Problem Statement", "content": problem, "icon": MSO_SHAPE.RECTANGLE},
            {"title": "Proposed Solution", "content": solution, "icon": MSO_SHAPE.ROUNDED_RECTANGLE},
            {"title": "Expected Impact", "content": impact, "icon": MSO_SHAPE.CHEVRON},
            {"title": "Timeline", "content": timeline, "icon": MSO_SHAPE.PENTAGON}
        ]
        
        positions = [
            (0.5, 1.5),  # Top left
            (5, 1.5),    # Top right
            (0.5, 4),    # Bottom left
            (5, 4)       # Bottom right
        ]
        
        for i, (quad, pos) in enumerate(zip(quadrants, positions)):
            # Add shape
            shape = slide.shapes.add_shape(
                quad["icon"], Inches(pos[0]), Inches(pos[1]), Inches(4), Inches(2)
            )
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = self.colors['accent'] if i % 2 == 0 else self.colors['secondary']
            
            # Add text
            text_frame = shape.text_frame
            text_frame.clear()
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)
            text_frame.margin_top = Inches(0.2)
            text_frame.margin_bottom = Inches(0.2)
            
            p = text_frame.paragraphs[0]
            p.text = quad["title"]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            
            p = text_frame.add_paragraph()
            p.text = quad["content"]
            p.font.size = Pt(12)
            p.font.color.rgb = self.colors['white']
            p.line_spacing = 1.2
        
        return slide
    
    def add_pie_chart(self, title, data_dict):
        """Create slide with pie chart"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # Create matplotlib pie chart
        plt.figure(figsize=(8, 6))
        colors = ['#003399', '#0072CE', '#00B0F0', '#00B050', '#FFC000']
        
        # Create pie chart
        wedges, texts, autotexts = plt.pie(
            data_dict.values(), 
            labels=data_dict.keys(), 
            autopct='%1.1f%%',
            colors=colors[:len(data_dict)],
            startangle=90
        )
        
        # Enhance text
        for text in texts:
            text.set_fontsize(12)
            text.set_fontweight('bold')
        
        for autotext in autotexts:
            autotext.set_color('white')
            autotext.set_fontsize(11)
            autotext.set_fontweight('bold')
        
        plt.title('')
        
        # Save to memory
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=300, bbox_inches='tight')
        plt.close()
        
        # Add to slide
        img_stream.seek(0)
        pic = slide.shapes.add_picture(img_stream, Inches(2), Inches(1.5), width=Inches(6))
        
        return slide
    
    def add_bar_chart(self, title, categories, values, y_label="Value"):
        """Create slide with bar chart"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # Create bar chart data
        chart_data = ChartData()
        chart_data.categories = categories
        chart_data.add_series('Series 1', values)
        
        # Add chart to slide
        x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        
        # Format chart
        chart.has_legend = False
        
        # Format bars
        series = chart.series[0]
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['secondary']
        
        return slide
    
    def add_table(self, title, headers, data):
        """Create slide with professional table"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # Create table
        rows = len(data) + 1
        cols = len(headers)
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(0.5 * rows)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set column widths
        for i in range(cols):
            table.columns[i].width = Inches(8 / cols)
        
        # Add headers
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].font.color.rgb = self.colors['white']
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.colors['primary']
        
        # Add data
        for row_idx, row_data in enumerate(data):
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value)
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Alternate row colors
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
        
        return slide
    
    def add_timeline(self, title, milestones):
        """Create timeline infographic slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # Create timeline
        timeline_y = Inches(3.5)
        timeline_start = Inches(1)
        timeline_end = Inches(9)
        timeline_width = timeline_end - timeline_start
        
        # Draw timeline line
        line = slide.shapes.add_connector(
            1, timeline_start, timeline_y, timeline_end, timeline_y
        )
        line.line.color.rgb = self.colors['dark']
        line.line.width = Pt(3)
        
        # Add milestones
        num_milestones = len(milestones)
        spacing = timeline_width / (num_milestones + 1)
        
        for i, (date, milestone) in enumerate(milestones):
            x = timeline_start + spacing * (i + 1)
            
            # Add circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, x - Inches(0.25), timeline_y - Inches(0.25), 
                Inches(0.5), Inches(0.5)
            )
            fill = circle.fill
            fill.solid()
            fill.fore_color.rgb = self.colors['secondary']
            
            # Add date above
            date_box = slide.shapes.add_textbox(
                x - Inches(1), timeline_y - Inches(1.5), Inches(2), Inches(0.5)
            )
            date_frame = date_box.text_frame
            date_frame.text = date
            date_frame.paragraphs[0].font.size = Pt(12)
            date_frame.paragraphs[0].font.bold = True
            date_frame.paragraphs[0].font.color.rgb = self.colors['dark']
            date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Add milestone below
            milestone_box = slide.shapes.add_textbox(
                x - Inches(1.5), timeline_y + Inches(0.5), Inches(3), Inches(1)
            )
            milestone_frame = milestone_box.text_frame
            milestone_frame.text = milestone
            milestone_frame.paragraphs[0].font.size = Pt(11)
            milestone_frame.paragraphs[0].font.color.rgb = self.colors['dark']
            milestone_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            milestone_frame.word_wrap = True
        
        return slide
    
    def add_swot_analysis(self, strengths, weaknesses, opportunities, threats):
        """Create SWOT analysis slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = "SWOT Analysis"
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # Create SWOT quadrants
        swot_data = [
            ("Strengths", strengths, self.colors['success'], 0.5, 1.5),
            ("Weaknesses", weaknesses, self.colors['danger'], 5, 1.5),
            ("Opportunities", opportunities, self.colors['secondary'], 0.5, 4),
            ("Threats", threats, self.colors['warning'], 5, 4)
        ]
        
        for title, items, color, x, y in swot_data:
            # Add box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4), Inches(2)
            )
            fill = box.fill
            fill.solid()
            fill.fore_color.rgb = color
            
            # Add text
            text_frame = box.text_frame
            text_frame.clear()
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)
            text_frame.margin_top = Inches(0.2)
            
            p = text_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            
            for item in items:
                p = text_frame.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(11)
                p.font.color.rgb = self.colors['white']
                p.line_spacing = 1.2
        
        return slide
    
    def add_financial_analysis(self, title, data_df):
        """Create financial analysis slide with multiple charts"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # Create figure with subplots
        fig, ((ax1, ax2), (ax3, ax4)) = plt.subplots(2, 2, figsize=(10, 7))
        fig.suptitle('')
        
        # Style settings
        sns.set_style("whitegrid")
        colors = ['#003399', '#0072CE', '#00B0F0', '#00B050']
        
        # Plot 1: Revenue trend
        if 'Year' in data_df.columns and 'Revenue' in data_df.columns:
            ax1.plot(data_df['Year'], data_df['Revenue'], marker='o', color=colors[0], linewidth=2)
            ax1.set_title('Revenue Trend', fontsize=12, fontweight='bold')
            ax1.set_xlabel('Year')
            ax1.set_ylabel('Revenue ($M)')
            ax1.grid(True, alpha=0.3)
        
        # Plot 2: Cost breakdown
        if 'Costs' in data_df.columns:
            cost_categories = data_df.iloc[0]['Costs'] if isinstance(data_df.iloc[0]['Costs'], dict) else {'Total': 100}
            ax2.pie(cost_categories.values(), labels=cost_categories.keys(), autopct='%1.1f%%', colors=colors)
            ax2.set_title('Cost Breakdown', fontsize=12, fontweight='bold')
        
        # Plot 3: Profit margins
        if 'Profit_Margin' in data_df.columns:
            ax3.bar(data_df.index, data_df['Profit_Margin'], color=colors[1])
            ax3.set_title('Profit Margins', fontsize=12, fontweight='bold')
            ax3.set_ylabel('Margin (%)')
            ax3.grid(True, alpha=0.3, axis='y')
        
        # Plot 4: ROI projection
        if 'ROI' in data_df.columns:
            ax4.bar(data_df.index, data_df['ROI'], color=colors[2])
            ax4.set_title('ROI Projection', fontsize=12, fontweight='bold')
            ax4.set_ylabel('ROI (%)')
            ax4.grid(True, alpha=0.3, axis='y')
        
        plt.tight_layout()
        
        # Save to memory
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=300, bbox_inches='tight')
        plt.close()
        
        # Add to slide
        img_stream.seek(0)
        pic = slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.2), width=Inches(9))
        
        return slide
    
    def add_next_steps(self, immediate_actions, milestones_30day, long_term):
        """Create next steps slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = "Next Steps & Implementation"
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # Three columns
        columns = [
            ("Immediate Actions (Week 1)", immediate_actions, self.colors['danger']),
            ("30-Day Milestones", milestones_30day, self.colors['warning']),
            ("Long-term Goals", long_term, self.colors['success'])
        ]
        
        for i, (title, items, color) in enumerate(columns):
            x = 0.5 + i * 3
            
            # Add header box
            header = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.5), Inches(2.8), Inches(0.6)
            )
            fill = header.fill
            fill.solid()
            fill.fore_color.rgb = color
            
            text_frame = header.text_frame
            text_frame.text = title
            text_frame.paragraphs[0].font.size = Pt(14)
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].font.color.rgb = self.colors['white']
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Add items box
            items_box = slide.shapes.add_textbox(Inches(x), Inches(2.2), Inches(2.8), Inches(4))
            items_frame = items_box.text_frame
            items_frame.clear()
            
            for j, item in enumerate(items):
                p = items_frame.paragraphs[0] if j == 0 else items_frame.add_paragraph()
                p.text = f"□ {item}"
                p.font.size = Pt(11)
                p.font.color.rgb = self.colors['dark']
                p.line_spacing = 1.5
        
        return slide
    
    def save(self, filename):
        """Save the presentation"""
        self.prs.save(filename)
        print(f"✓ Presentation saved as: {filename}")
        return filename


# Example usage
def create_sample_presentation():
    """Create a sample case competition presentation"""
    ppt = CaseCompetitionPPT()
    
    # Title slide
    ppt.add_title_slide(
        title="Digital Transformation Strategy",
        subtitle="Accelerating Growth Through Technology",
        team_members=["John Doe", "Jane Smith", "Mike Johnson", "Sarah Williams"]
    )
    
    # Executive Summary
    ppt.add_executive_summary(
        problem="Company facing 20% YoY decline in market share due to digital disruption",
        solution="Implement comprehensive digital transformation across all channels",
        impact="Expected 35% revenue growth and 25% cost reduction within 18 months",
        timeline="Phase 1: Q1-Q2 | Phase 2: Q3-Q4 | Full Implementation: Year 2"
    )
    
    # Market Share Pie Chart
    market_data = {
        "Our Company": 15,
        "Competitor A": 25,
        "Competitor B": 20,
        "Competitor C": 18,
        "Others": 22
    }
    ppt.add_pie_chart("Current Market Share Distribution", market_data)
    
    # Revenue Projection Bar Chart
    years = ['2023', '2024', '2025', '2026']
    revenues = [45, 52, 68, 85]
    ppt.add_bar_chart("Revenue Projection ($M)", years, revenues, "Revenue ($M)")
    
    # Implementation Timeline
    milestones = [
        ("Q1 2024", "Platform Selection"),
        ("Q2 2024", "Pilot Launch"),
        ("Q3 2024", "Scale to 50%"),
        ("Q4 2024", "Full Rollout"),
        ("Q1 2025", "Optimization")
    ]
    ppt.add_timeline("Implementation Timeline", milestones)
    
    # Financial Analysis Table
    headers = ["Metric", "Current", "Year 1", "Year 2", "Year 3"]
    data = [
        ["Revenue ($M)", "45", "52", "68", "85"],
        ["Operating Margin (%)", "12%", "15%", "20%", "25%"],
        ["EBITDA ($M)", "5.4", "7.8", "13.6", "21.3"],
        ["ROI (%)", "-", "125%", "180%", "250%"]
    ]
    ppt.add_table("Financial Projections", headers, data)
    
    # SWOT Analysis
    ppt.add_swot_analysis(
        strengths=["Strong brand recognition", "Experienced team", "Financial stability"],
        weaknesses=["Legacy systems", "Limited digital expertise", "Resistance to change"],
        opportunities=["Growing digital market", "Customer demand for online", "Cost reduction potential"],
        threats=["Aggressive competitors", "Technology risks", "Implementation complexity"]
    )
    
    # Next Steps
    ppt.add_next_steps(
        immediate_actions=["Form transformation team", "Conduct technology audit", "Define success metrics"],
        milestones_30day=["Complete vendor selection", "Launch pilot program", "Train core team"],
        long_term=["Achieve 50% digital revenue", "Reduce costs by 25%", "Market leader position"]
    )
    
    # Save presentation
    ppt.save("Case_Competition_Presentation.pptx")
    
    return ppt


if __name__ == "__main__":
    print("Creating advanced case competition presentation...")
    create_sample_presentation()
    print("\n✓ All visualization capabilities ready!")
    print("✓ Can create: Pie charts, Bar charts, Tables, Timelines, SWOT, Financial Analysis")
    print("✓ Professional McKinsey-style formatting applied")