#!/usr/bin/env python3
"""
DisruptX Round 1 - MediChain Presentation VERSION 1
Professional BCG/McKinsey-style 3-slide deck with perfect spacing and formatting
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import numpy as np
import io

class DisruptXVersion1:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9
        self.prs.slide_height = Inches(7.5)
        
        # Professional color palette
        self.colors = {
            'primary_blue': RGBColor(0, 83, 159),
            'teal': RGBColor(0, 150, 170), 
            'light_teal': RGBColor(179, 229, 234),
            'dark_gray': RGBColor(64, 64, 64),
            'medium_gray': RGBColor(128, 128, 128),
            'light_gray': RGBColor(191, 191, 191),
            'very_light_gray': RGBColor(242, 242, 242),
            'white': RGBColor(255, 255, 255),
            'green': RGBColor(119, 187, 65),
            'orange': RGBColor(255, 138, 0),
            'red': RGBColor(220, 38, 127),
        }
    
    def _add_title_with_margin(self, slide, text, y_position=0.5):
        """Add title with proper margins"""
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(y_position), 
            Inches(11.8), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.clear()
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Calibri'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
        p.alignment = PP_ALIGN.LEFT
        
    def _add_text_box_with_padding(self, slide, text, left, top, width, height, font_size=12, bold=False, color=None):
        """Add text box with proper padding"""
        text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = text_box.text_frame
        tf.clear()
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Calibri'
        p.font.size = Pt(font_size)
        p.font.bold = bold
        if color:
            p.font.color.rgb = color
        else:
            p.font.color.rgb = self.colors['dark_gray']
        
        return text_box
        
    def _add_bottom_banner(self, slide, text):
        """Add bottom banner"""
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(6.8),
            Inches(13.333), Inches(0.7)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = self.colors['primary_blue']
        banner.line.fill.background()
        
        tf = banner.text_frame
        tf.clear()
        tf.margin_left = Inches(0.5)
        tf.margin_right = Inches(0.5)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Calibri'
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
    def create_slide1_market_opportunity(self):
        """Slide 1: The 600M Opportunity"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_title_with_margin(slide, "Tier 2/3 India: The Untapped Healthcare Goldmine")
        
        # Left section - Market size visual
        left_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.75), Inches(1.5),
            Inches(5.5), Inches(4.8)
        )
        left_box.fill.solid()
        left_box.fill.fore_color.rgb = self.colors['very_light_gray']
        left_box.line.color.rgb = self.colors['light_gray']
        
        # Market Size Chart
        fig, ax = plt.subplots(figsize=(5, 4))
        
        # Create layered circles for population representation
        circle1 = plt.Circle((0.5, 0.5), 0.45, color='#0053A0', alpha=0.8, label='Total: 600M')
        circle2 = plt.Circle((0.5, 0.5), 0.35, color='#0096AA', alpha=0.8, label='Underserved: 450M')
        circle3 = plt.Circle((0.5, 0.5), 0.25, color='#FF8A00', alpha=0.8, label='Critical Need: 300M')
        
        ax.add_patch(circle1)
        ax.add_patch(circle2)
        ax.add_patch(circle3)
        
        # Add text annotations with proper spacing
        ax.text(0.5, 0.9, '600M', fontsize=28, fontweight='bold', ha='center', color='white')
        ax.text(0.5, 0.8, 'Total Population', fontsize=12, ha='center', color='white')
        ax.text(0.5, 0.5, '450M', fontsize=22, fontweight='bold', ha='center', color='white')
        ax.text(0.5, 0.4, 'Underserved', fontsize=10, ha='center', color='white')
        ax.text(0.5, 0.15, '300M', fontsize=18, fontweight='bold', ha='center', color='white')
        ax.text(0.5, 0.05, 'Critical', fontsize=9, ha='center', color='white')
        
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        ax.axis('off')
        ax.set_title('Healthcare Access Gap', fontsize=16, fontweight='bold', pad=20)
        
        # Save and add to slide
        buf = io.BytesIO()
        plt.tight_layout()
        plt.savefig(buf, format='png', dpi=300, transparent=True)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(1), Inches(2), height=Inches(3.5))
        plt.close()
        
        # Right section - Key challenges with icons
        challenges = [
            ("🏥", "Doctor Availability", "1 doctor per 11,000 people (vs 1:1,000 in cities)"),
            ("🚗", "Distance Barrier", "Avg 47km to nearest hospital, 3+ hours travel"),
            ("💰", "Affordability Crisis", "60% earn <₹10,000/month, 85% lack insurance"),
            ("📱", "Digital Readiness", "67% own smartphones, 4G coverage expanding")
        ]
        
        y_pos = 1.8
        for icon, title, desc in challenges:
            # Icon box
            icon_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(6.5), Inches(y_pos),
                Inches(0.6), Inches(0.6)
            )
            icon_box.fill.solid()
            icon_box.fill.fore_color.rgb = self.colors['teal']
            icon_box.line.fill.background()
            
            # Add icon text
            tf = icon_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = icon
            p.font.size = Pt(20)
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            
            # Title and description with proper spacing
            self._add_text_box_with_padding(
                slide, title, 7.3, y_pos, 5, 0.3,
                font_size=14, bold=True, color=self.colors['dark_gray']
            )
            self._add_text_box_with_padding(
                slide, desc, 7.3, y_pos + 0.3, 5, 0.4,
                font_size=11, color=self.colors['medium_gray']
            )
            
            y_pos += 1.1
        
        self._add_bottom_banner(slide, "Healthcare gap in Tier 2/3 cities represents India's largest underserved market opportunity")
        
        # Add slide number
        num_box = slide.shapes.add_textbox(Inches(12.8), Inches(7.1), Inches(0.4), Inches(0.3))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = "1"
        p.font.size = Pt(10)
        p.font.color.rgb = self.colors['medium_gray']
        p.alignment = PP_ALIGN.RIGHT
        
    def create_slide2_solution(self):
        """Slide 2: MediChain - The Tech-Enabled Solution"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_title_with_margin(slide, "MediChain: AI-Powered Healthcare Ecosystem for Bharat")
        
        # Solution architecture diagram
        fig, ax = plt.subplots(figsize=(11, 4.5))
        
        # Central platform
        central_rect = patches.FancyBboxPatch(
            (0.4, 0.4), 0.2, 0.2, boxstyle="round,pad=0.02",
            facecolor='#0053A0', edgecolor='#0053A0', linewidth=2
        )
        ax.add_patch(central_rect)
        ax.text(0.5, 0.5, 'MediChain\nPlatform', fontsize=14, fontweight='bold', 
                ha='center', va='center', color='white')
        
        # Components with proper spacing
        components = [
            (0.15, 0.7, "AI Triage\nSystem", '#0096AA'),
            (0.15, 0.3, "Vernacular\nInterface", '#77BB41'),
            (0.5, 0.8, "Doctor\nNetwork", '#FF8A00'),
            (0.85, 0.7, "Medicine\nDelivery", '#DC267F'),
            (0.85, 0.3, "Diagnostic\nIntegration", '#664D91'),
            (0.5, 0.2, "ASHA Worker\nApp", '#00B050')
        ]
        
        for x, y, label, color in components:
            # Component box
            comp_rect = patches.FancyBboxPatch(
                (x-0.08, y-0.08), 0.16, 0.16, boxstyle="round,pad=0.02",
                facecolor=color, edgecolor=color, linewidth=2, alpha=0.9
            )
            ax.add_patch(comp_rect)
            ax.text(x, y, label, fontsize=11, fontweight='bold', 
                    ha='center', va='center', color='white')
            
            # Connect to center
            if x != 0.5 or y != 0.5:
                ax.plot([x, 0.5], [y, 0.5], 'k--', alpha=0.3, linewidth=1)
        
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        ax.axis('off')
        
        # Save and add to slide
        buf = io.BytesIO()
        plt.tight_layout()
        plt.savefig(buf, format='png', dpi=300, transparent=True)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(0.75), Inches(1.5), width=Inches(11.8))
        plt.close()
        
        # Competitive Landscape - Fixed version
        comp_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.75), Inches(4.2),
            Inches(5.5), Inches(2.2)
        )
        comp_box.fill.solid()
        comp_box.fill.fore_color.rgb = self.colors['very_light_gray']
        comp_box.line.color.rgb = self.colors['light_gray']
        
        self._add_text_box_with_padding(
            slide, "Competitive Positioning", 0.9, 4.3, 5, 0.4,
            font_size=14, bold=True, color=self.colors['primary_blue']
        )
        
        # Create competitive matrix as a table
        table_left = 1
        table_top = 4.8
        table_width = 5
        table_height = 1.3
        
        table = slide.shapes.add_table(3, 3, Inches(table_left), Inches(table_top), 
                                      Inches(table_width), Inches(table_height)).table
        
        # Set column widths
        table.columns[0].width = Inches(1.8)
        table.columns[1].width = Inches(1.6)
        table.columns[2].width = Inches(1.6)
        
        # Headers
        headers = ['', 'Urban', 'Rural']
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Data
        data = [
            ['Broad', 'Practo, 1mg', 'Gap: Vernacular\n+ Low-cost'],
            ['Narrow', 'PharmEasy', 'Gap: Trust\n+ ASHA']
        ]
        
        for i, row in enumerate(data):
            table.cell(i+1, 0).text = row[0]
            table.cell(i+1, 0).text_frame.paragraphs[0].font.bold = True
            table.cell(i+1, 1).text = row[1]
            table.cell(i+1, 2).text = row[2]
            
            # Color the gap cells
            table.cell(i+1, 2).fill.solid()
            table.cell(i+1, 2).fill.fore_color.rgb = self.colors['light_teal']
        
        # Patient Journey Funnel - Fixed
        funnel_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.5), Inches(4.2),
            Inches(6.3), Inches(2.2)
        )
        funnel_box.fill.solid()
        funnel_box.fill.fore_color.rgb = self.colors['very_light_gray']
        funnel_box.line.color.rgb = self.colors['light_gray']
        
        self._add_text_box_with_padding(
            slide, "Patient Journey & Conversion", 6.65, 4.3, 6, 0.4,
            font_size=14, bold=True, color=self.colors['primary_blue']
        )
        
        # Create improved funnel visualization
        fig, ax = plt.subplots(figsize=(5.5, 1.5))
        
        stages = ['Awareness', 'Registration', 'Consultation', 'Treatment', 'Follow-up']
        percentages = [100, 75, 60, 45, 35]
        colors = ['#0053A0', '#0096AA', '#77BB41', '#FF8A00', '#DC267F']
        
        y_pos = 0.5
        bar_height = 0.3
        
        for i, (stage, pct, color) in enumerate(zip(stages, percentages, colors)):
            # Draw percentage bar
            bar_width = pct / 100 * 0.9
            bar = patches.Rectangle((0.05, y_pos - bar_height/2), bar_width, bar_height,
                                   facecolor=color, edgecolor='none', alpha=0.8)
            ax.add_patch(bar)
            
            # Add percentage text inside bar
            ax.text(0.05 + bar_width/2, y_pos, f'{pct}%', fontsize=11, fontweight='bold',
                   ha='center', va='center', color='white')
            
            # Add stage label below
            ax.text(0.05 + bar_width/2, y_pos - bar_height/2 - 0.1, stage, 
                   fontsize=9, ha='center', va='top')
        
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        ax.axis('off')
        
        buf = io.BytesIO()
        plt.tight_layout()
        plt.savefig(buf, format='png', dpi=300, transparent=True)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(6.7), Inches(4.8), height=Inches(1.4))
        plt.close()
        
        self._add_bottom_banner(slide, "MediChain's integrated approach addresses every barrier in rural healthcare delivery")
        
        # Add slide number
        num_box = slide.shapes.add_textbox(Inches(12.8), Inches(7.1), Inches(0.4), Inches(0.3))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = "2"
        p.font.size = Pt(10)
        p.font.color.rgb = self.colors['medium_gray']
        p.alignment = PP_ALIGN.RIGHT
        
    def create_slide3_impact(self):
        """Slide 3: Implementation Roadmap & Financial Impact"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_title_with_margin(slide, "Path to ₹5,000 Cr: Scale, Impact & Returns")
        
        # Timeline roadmap
        timeline_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.75), Inches(1.5),
            Inches(11.8), Inches(2)
        )
        timeline_box.fill.solid()
        timeline_box.fill.fore_color.rgb = self.colors['very_light_gray']
        timeline_box.line.color.rgb = self.colors['light_gray']
        
        # Create timeline graphic
        fig, ax = plt.subplots(figsize=(11, 1.8))
        
        # Timeline data
        phases = [
            ("Phase 1\n6 months", "MVP Launch\n5 districts\n10K users", 0.15),
            ("Phase 2\n12 months", "State Expansion\n50 districts\n500K users", 0.35),
            ("Phase 3\n18 months", "Multi-State\n200 districts\n5M users", 0.55),
            ("Phase 4\n24 months", "National Scale\n500+ districts\n50M users", 0.75)
        ]
        
        # Draw timeline
        ax.plot([0.1, 0.8], [0.5, 0.5], 'k-', linewidth=3)
        
        for phase, desc, x_pos in phases:
            # Milestone circle
            circle = plt.Circle((x_pos, 0.5), 0.04, color='#0053A0', zorder=10)
            ax.add_patch(circle)
            
            # Phase text above
            ax.text(x_pos, 0.75, phase, fontsize=10, fontweight='bold',
                   ha='center', va='center')
            
            # Description below
            ax.text(x_pos, 0.25, desc, fontsize=8,
                   ha='center', va='center', style='italic')
        
        ax.set_xlim(0, 0.9)
        ax.set_ylim(0, 1)
        ax.axis('off')
        
        buf = io.BytesIO()
        plt.tight_layout()
        plt.savefig(buf, format='png', dpi=300, transparent=True)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(0.9), Inches(1.7), width=Inches(11.5))
        plt.close()
        
        # Financial metrics (left)
        fin_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.75), Inches(3.8),
            Inches(5.5), Inches(2.7)
        )
        fin_box.fill.solid()
        fin_box.fill.fore_color.rgb = self.colors['very_light_gray']
        fin_box.line.color.rgb = self.colors['light_gray']
        
        self._add_text_box_with_padding(
            slide, "Financial Projections (Year 5)", 0.9, 3.9, 5, 0.4,
            font_size=14, bold=True, color=self.colors['primary_blue']
        )
        
        # Financial metrics with proper spacing
        metrics = [
            ("Revenue", "₹5,000 Cr", "#0053A0"),
            ("EBITDA Margin", "22%", "#0096AA"),
            ("Break-even", "Month 30", "#77BB41"),
            ("Funding Need", "₹500 Cr", "#FF8A00")
        ]
        
        y_start = 4.4
        for i, (label, value, color) in enumerate(metrics):
            # Metric box
            metric_rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1 + (i % 2) * 2.5), Inches(y_start + (i // 2) * 0.9),
                Inches(2.3), Inches(0.7)
            )
            metric_rect.fill.solid()
            metric_rect.fill.fore_color.rgb = RGBColor(*[int(color[i:i+2], 16) for i in (1, 3, 5)])
            metric_rect.line.fill.background()
            
            tf = metric_rect.text_frame
            tf.clear()
            tf.margin_all = Inches(0.1)
            p = tf.paragraphs[0]
            p.text = f"{label}\n{value}"
            p.font.name = 'Calibri'
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # Impact metrics (right)
        impact_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.5), Inches(3.8),
            Inches(6.3), Inches(2.7)
        )
        impact_box.fill.solid()
        impact_box.fill.fore_color.rgb = self.colors['very_light_gray']
        impact_box.line.color.rgb = self.colors['light_gray']
        
        self._add_text_box_with_padding(
            slide, "Social Impact by Year 5", 6.65, 3.9, 6, 0.4,
            font_size=14, bold=True, color=self.colors['primary_blue']
        )
        
        # Create impact visualization
        fig, ax = plt.subplots(figsize=(5.5, 2))
        
        impacts = [
            ("Lives Impacted", "50M", "Direct healthcare access"),
            ("Jobs Created", "100K", "ASHA workers, delivery partners"),
            ("Cost Reduction", "70%", "vs traditional healthcare"),
            ("Time Saved", "3 hrs", "Per consultation on average")
        ]
        
        y_positions = np.linspace(0.8, 0.2, len(impacts))
        
        for (metric, value, desc), y_pos in zip(impacts, y_positions):
            # Value circle
            circle = plt.Circle((0.15, y_pos), 0.08, color='#0053A0', alpha=0.8)
            ax.add_patch(circle)
            ax.text(0.15, y_pos, value, fontsize=10, fontweight='bold',
                   ha='center', va='center', color='white')
            
            # Metric and description
            ax.text(0.28, y_pos + 0.04, metric, fontsize=11, fontweight='bold')
            ax.text(0.28, y_pos - 0.04, desc, fontsize=9, style='italic', color='#666666')
        
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        ax.axis('off')
        
        buf = io.BytesIO()
        plt.tight_layout()
        plt.savefig(buf, format='png', dpi=300, transparent=True)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(6.7), Inches(4.4), width=Inches(5.8))
        plt.close()
        
        self._add_bottom_banner(slide, "MediChain: Transforming rural healthcare while building a sustainable, profitable business")
        
        # Add slide number
        num_box = slide.shapes.add_textbox(Inches(12.8), Inches(7.1), Inches(0.4), Inches(0.3))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = "3"
        p.font.size = Pt(10)
        p.font.color.rgb = self.colors['medium_gray']
        p.alignment = PP_ALIGN.RIGHT
        
    def generate_presentation(self):
        """Generate the complete presentation"""
        self.create_slide1_market_opportunity()
        self.create_slide2_solution()
        self.create_slide3_impact()
        
        # Save to PPT Generated folder
        output_path = "/mnt/e/AI and Projects/Case Comp PPT/PPT Generated/DisruptX_MediChain_Version1.pptx"
        self.prs.save(output_path)
        print(f"Presentation saved to: {output_path}")
        return output_path

if __name__ == "__main__":
    # Create the presentation
    creator = DisruptXVersion1()
    output_file = creator.generate_presentation()
    print(f"\nDisruptX MediChain presentation Version 1 created successfully!")
    print(f"Location: {output_file}")