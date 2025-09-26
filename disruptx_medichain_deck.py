#!/usr/bin/env python3
"""
DisruptX Round 1 - MediChain Presentation
BCG/McKinsey-style 3-slide deck for Tier-2/3 India healthcare disruption
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, Circle, Rectangle, Arrow, FancyArrowPatch
import seaborn as sns
import pandas as pd
import numpy as np
import io
from datetime import datetime

class DisruptXPresentation:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 widescreen
        self.prs.slide_height = Inches(7.5)
        
        # BCG/McKinsey color scheme
        self.colors = {
            'primary_blue': RGBColor(0, 83, 159),      # BCG Blue
            'teal': RGBColor(0, 150, 170),            # Teal
            'light_blue': RGBColor(102, 204, 238),    # Light Blue
            'dark_gray': RGBColor(64, 64, 64),        # Dark Gray
            'medium_gray': RGBColor(128, 128, 128),   # Medium Gray
            'light_gray': RGBColor(191, 191, 191),    # Light Gray
            'very_light_gray': RGBColor(242, 242, 242), # Background gray
            'white': RGBColor(255, 255, 255),         # White
            'accent_green': RGBColor(119, 187, 65),    # Success green
            'accent_orange': RGBColor(255, 138, 0),    # Warning orange
            'accent_red': RGBColor(220, 38, 127),      # Alert red
        }
        
        # Icon dictionary (using text symbols)
        self.icons = {
            'healthcare': '🏥',
            'education': '🎓',
            'finance': '💰',
            'agriculture': '🌾',
            'smartphone': '📱',
            'internet': '🌐',
            'location': '📍',
            'doctor': '👨‍⚕️',
            'patient': '👥',
            'ai': '🤖',
            'blockchain': '🔗',
            'pharmacy': '💊',
            'diagnostic': '🔬',
            'arrow_up': '↗',
            'arrow_right': '→',
            'checkmark': '✓'
        }
    
    def _add_slide_number(self, slide, number):
        """Add slide number to bottom right"""
        textbox = slide.shapes.add_textbox(
            Inches(12.5), Inches(7.1), Inches(0.5), Inches(0.3)
        )
        tf = textbox.text_frame
        tf.text = str(number)
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = self.colors['medium_gray']
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
    def _add_bottom_banner(self, slide, text):
        """Add bottom banner with key takeaway"""
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(6.8),
            self.prs.slide_width, Inches(0.7)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = self.colors['primary_blue']
        banner.line.fill.background()
        
        tf = banner.text_frame
        tf.margin_left = Inches(0.5)
        tf.margin_right = Inches(0.5)
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Calibri'
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
    def _add_speaker_notes(self, slide, notes):
        """Add speaker notes to slide"""
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
        
    def _inches_to_float(self, inches_obj):
        """Convert Inches object to float value"""
        return inches_obj / 914400  # EMU to inches conversion
    
    def create_slide_1_opportunity_landscape(self):
        """Create Slide 1: Opportunity Landscape"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6))
        tf = title_box.text_frame
        tf.text = "Why Tier-2 & Tier-3 India is Ripe for Disruption"
        p = tf.paragraphs[0]
        p.font.name = 'Calibri'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
        
        # Left section: Macro Trends (40%)
        self._add_macro_trends(slide, Inches(0.5), Inches(1.2), Inches(5.5), Inches(2.8))
        
        # Right section: India Map (30%)
        self._add_india_map_visual(slide, Inches(6.2), Inches(1.2), Inches(4), Inches(2.8))
        
        # Bottom left: Underserved Sectors (35%)
        self._add_underserved_sectors(slide, Inches(0.5), Inches(4.2), Inches(5.5), Inches(2.3))
        
        # Bottom right: Infrastructure Access Chart (35%)
        self._add_infrastructure_chart(slide, Inches(6.2), Inches(4.2), Inches(6.6), Inches(2.3))
        
        # Bottom banner
        self._add_bottom_banner(slide, 
            "Digital readiness + structural gaps = fertile ground for tech-enabled disruption in Bharat.")
        
        # Speaker notes
        self._add_speaker_notes(slide,
            "Tier-2 and Tier-3 India are digitally enabled — UPI, smartphones, cheap data — "
            "but structurally underserved across healthcare, education, finance and agriculture. "
            "This duality creates massive disruption potential.")
        
        # Slide number
        self._add_slide_number(slide, 1)
        
        return slide
    
    def _add_macro_trends(self, slide, x, y, width, height):
        """Add macro trends section with icons"""
        # Container
        trends_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height
        )
        trends_box.fill.solid()
        trends_box.fill.fore_color.rgb = self.colors['very_light_gray']
        trends_box.line.color.rgb = self.colors['light_gray']
        trends_box.line.width = Pt(1)
        
        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, Inches(0.4)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.colors['teal']
        header.line.fill.background()
        
        tf = header.text_frame
        tf.text = "MACRO TRENDS DRIVING DISRUPTION"
        p = tf.paragraphs[0]
        p.font.name = 'Calibri'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # Trends data
        trends = [
            ('💹', 'GDP Growth', 'Tier-2/3 cities → ~45% of India\'s GDP by 2025'),
            ('👥', 'Population', '~650M population base'),
            ('📱', 'Smartphone', '~60% penetration in Tier-2/3 (vs 78% urban)'),
            ('💳', 'UPI Adoption', '>12B transactions/month'),
            ('🌐', 'Data Cost', 'Among lowest globally (~$0.17/GB)')
        ]
        
        # Add trends in 2 columns
        col_width = width / 2 - Inches(0.2)
        y_start = y + Inches(0.6)
        
        for i, (icon, title, desc) in enumerate(trends):
            col = i % 2
            row = i // 2
            x_pos = x + Inches(0.1) + (col_width + Inches(0.2)) * col
            y_pos = y_start + Inches(0.5) * row
            
            # Icon and text
            text_box = slide.shapes.add_textbox(x_pos, y_pos, col_width, Inches(0.5))
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            
            # Icon
            run = p.add_run()
            run.text = f"{icon} "
            run.font.size = Pt(16)
            
            # Title
            run = p.add_run()
            run.text = f"{title}: "
            run.font.name = 'Calibri'
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = self.colors['primary_blue']
            
            # Description
            run = p.add_run()
            run.text = desc
            run.font.name = 'Calibri'
            run.font.size = Pt(10)
            run.font.color.rgb = self.colors['dark_gray']
    
    def _add_india_map_visual(self, slide, x, y, width, height):
        """Add India map highlighting Tier-2/3 cities"""
        # Create matplotlib figure for map
        fig, ax = plt.subplots(figsize=(self._inches_to_float(width), self._inches_to_float(height)))
        
        # Simple India outline (stylized)
        india_x = [0.5, 0.3, 0.2, 0.3, 0.5, 0.7, 0.8, 0.7, 0.5]
        india_y = [0.9, 0.7, 0.5, 0.3, 0.1, 0.2, 0.5, 0.8, 0.9]
        
        ax.fill(india_x, india_y, color='#E8F4F8', edgecolor='#0053A0', linewidth=2)
        
        # Highlight regions
        regions = {
            'North': (0.5, 0.7, '#0096AA'),
            'East': (0.7, 0.5, '#0096AA'),
            'West': (0.3, 0.5, '#0096AA'),
            'South': (0.5, 0.3, '#0096AA')
        }
        
        for region, (rx, ry, color) in regions.items():
            circle = plt.Circle((rx, ry), 0.08, color=color, alpha=0.5)
            ax.add_patch(circle)
            ax.text(rx, ry, region[0], ha='center', va='center', 
                   fontsize=10, fontweight='bold', color='white')
        
        # Sample cities
        cities = [
            ('Lucknow', 0.5, 0.65),
            ('Jaipur', 0.35, 0.6),
            ('Nagpur', 0.5, 0.5),
            ('Coimbatore', 0.45, 0.25),
            ('Bhubaneswar', 0.75, 0.45)
        ]
        
        for city, cx, cy in cities:
            ax.plot(cx, cy, 'ro', markersize=6)
            ax.text(cx + 0.05, cy, city, fontsize=8, style='italic')
        
        # Title
        ax.text(0.5, 0.95, 'Tier-2/3 India: Focus Regions', 
               ha='center', fontsize=12, fontweight='bold')
        
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        ax.axis('off')
        
        plt.tight_layout()
        
        # Save and add to slide
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=300, transparent=True, bbox_inches='tight')
        plt.close()
        
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, x, y, width=width)
    
    def _add_underserved_sectors(self, slide, x, y, width, height):
        """Add underserved sectors with icons"""
        # Container
        sectors_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height
        )
        sectors_box.fill.solid()
        sectors_box.fill.fore_color.rgb = self.colors['white']
        sectors_box.line.color.rgb = self.colors['light_gray']
        sectors_box.line.width = Pt(1)
        
        # Header
        header_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.1), width - Inches(0.4), Inches(0.3))
        tf = header_box.text_frame
        tf.text = "UNDERSERVED SECTORS"
        p = tf.paragraphs[0]
        p.font.name = 'Calibri'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
        
        # Sectors
        sectors = [
            ('🏥', 'Healthcare', '600M underserved\nAvg travel 50+ km'),
            ('🎓', 'Education', 'Teacher ratio 1:60\nvs 1:30 ideal'),
            ('💰', 'Finance', '190M unbanked\n<5% insurance'),
            ('🌾', 'Agriculture', 'Post-harvest loss\n₹90,000 Cr/yr')
        ]
        
        # Add sectors in grid
        sector_width = (width - Inches(0.4)) / 4
        y_pos = y + Inches(0.6)
        
        for i, (icon, name, stats) in enumerate(sectors):
            x_pos = x + Inches(0.2) + sector_width * i
            
            # Icon box
            icon_box = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x_pos + sector_width/2 - Inches(0.3), y_pos,
                Inches(0.6), Inches(0.6)
            )
            icon_box.fill.solid()
            icon_box.fill.fore_color.rgb = self.colors['teal']
            icon_box.line.fill.background()
            
            # Icon
            tf = icon_box.text_frame
            tf.text = icon
            p = tf.paragraphs[0]
            p.font.size = Pt(20)
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            
            # Sector name
            name_box = slide.shapes.add_textbox(
                x_pos, y_pos + Inches(0.7), sector_width, Inches(0.3)
            )
            tf = name_box.text_frame
            tf.text = name
            p = tf.paragraphs[0]
            p.font.name = 'Calibri'
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.colors['dark_gray']
            p.alignment = PP_ALIGN.CENTER
            
            # Stats
            stats_box = slide.shapes.add_textbox(
                x_pos, y_pos + Inches(1), sector_width, Inches(0.6)
            )
            tf = stats_box.text_frame
            tf.text = stats
            p = tf.paragraphs[0]
            p.font.name = 'Calibri'
            p.font.size = Pt(9)
            p.font.color.rgb = self.colors['medium_gray']
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.2
    
    def _add_infrastructure_chart(self, slide, x, y, width, height):
        """Add infrastructure access comparison chart"""
        # Create chart data
        categories = ['Doctors per\n1,000 people', 'Bank branches\nper 100k', 
                      'Internet\npenetration (%)', 'Smartphone\npenetration (%)']
        urban = [1.8, 18, 78, 85]
        tier23 = [0.5, 6, 62, 60]
        
        # Create matplotlib figure
        fig, ax = plt.subplots(figsize=(self._inches_to_float(width), self._inches_to_float(height)))
        
        x_pos = np.arange(len(categories))
        bar_width = 0.35
        
        # Create bars
        bars1 = ax.bar(x_pos - bar_width/2, urban, bar_width, 
                       label='Urban', color='#0053A0', alpha=0.8)
        bars2 = ax.bar(x_pos + bar_width/2, tier23, bar_width,
                       label='Tier-2/3', color='#0096AA', alpha=0.8)
        
        # Add value labels
        for bars in [bars1, bars2]:
            for bar in bars:
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2., height,
                       f'{height:g}', ha='center', va='bottom', fontsize=9)
        
        # Styling
        ax.set_xlabel('')
        ax.set_title('Infrastructure Access: Urban vs Tier-2/3', 
                    fontsize=12, fontweight='bold', pad=10)
        ax.set_xticks(x_pos)
        ax.set_xticklabels(categories, fontsize=9)
        ax.legend(loc='upper right', fontsize=10)
        ax.grid(axis='y', alpha=0.3)
        ax.set_axisbelow(True)
        
        # Remove top and right spines
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        
        plt.tight_layout()
        
        # Save and add to slide
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=300, bbox_inches='tight')
        plt.close()
        
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, x, y, width=width)
    
    def create_slide_2_sector_focus(self):
        """Create Slide 2: Healthcare Sector Focus"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6))
        tf = title_box.text_frame
        tf.text = "Healthcare in Tier-2/3 India: Urgent Problem, Large Market"
        p = tf.paragraphs[0]
        p.font.name = 'Calibri'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
        
        # Top left: Problem Pillars (40%)
        self._add_problem_pillars(slide, Inches(0.5), Inches(1.2), Inches(5), Inches(2.5))
        
        # Top right: Market Potential (30%)
        self._add_market_potential(slide, Inches(5.7), Inches(1.2), Inches(3.8), Inches(2.5))
        
        # Bottom left: Healthcare Financing Pie (30%)
        self._add_healthcare_financing_pie(slide, Inches(9.7), Inches(1.2), Inches(3.1), Inches(2.5))
        
        # Bottom: Competitive Landscape Matrix (100%)
        self._add_competitive_matrix(slide, Inches(0.5), Inches(3.9), Inches(12.3), Inches(2.6))
        
        # Bottom banner
        self._add_bottom_banner(slide,
            "Healthcare = burning platform: unmet need, adoption signals, and policy tailwinds make it the top disruption sector.")
        
        # Speaker notes
        self._add_speaker_notes(slide,
            "Healthcare in Tier-2/3 has the sharpest pain points: doctor shortages, catastrophic OOP costs, "
            "and proven adoption of telemedicine. This sector is where tech disruption has the highest impact "
            "and fastest scalability.")
        
        # Slide number
        self._add_slide_number(slide, 2)
        
        return slide
    
    def _add_problem_pillars(self, slide, x, y, width, height):
        """Add three problem pillars"""
        pillars = [
            {
                'title': 'Accessibility',
                'icon': '🏥',
                'stats': ['75% of doctors in urban areas', '600M underserved', 'Travel >50 km'],
                'color': self.colors['accent_red']
            },
            {
                'title': 'Affordability', 
                'icon': '💰',
                'stats': ['OOP spend = 62%', '~60M pushed to poverty', 'Catastrophic costs'],
                'color': self.colors['accent_orange']
            },
            {
                'title': 'Awareness & Trust',
                'icon': '🤝',
                'stats': ['Preventive care stigma', 'Reliance on unlicensed', 'Low health literacy'],
                'color': self.colors['primary_blue']
            }
        ]
        
        pillar_width = width / 3 - Inches(0.1)
        
        for i, pillar in enumerate(pillars):
            x_pos = x + (pillar_width + Inches(0.1)) * i
            
            # Pillar container
            pillar_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_pos, y, pillar_width, height
            )
            pillar_box.fill.solid()
            pillar_box.fill.fore_color.rgb = self.colors['very_light_gray']
            pillar_box.line.color.rgb = pillar['color']
            pillar_box.line.width = Pt(2)
            
            # Header
            header = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x_pos, y, pillar_width, Inches(0.5)
            )
            header.fill.solid()
            header.fill.fore_color.rgb = pillar['color']
            header.line.fill.background()
            
            tf = header.text_frame
            tf.text = f"{pillar['icon']} {pillar['title']}"
            p = tf.paragraphs[0]
            p.font.name = 'Calibri'
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            
            # Stats
            stats_box = slide.shapes.add_textbox(
                x_pos + Inches(0.1), y + Inches(0.6),
                pillar_width - Inches(0.2), height - Inches(0.7)
            )
            tf = stats_box.text_frame
            
            for stat in pillar['stats']:
                p = tf.paragraphs[0] if stat == pillar['stats'][0] else tf.add_paragraph()
                p.text = f"• {stat}"
                p.font.name = 'Calibri'
                p.font.size = Pt(10)
                p.font.color.rgb = self.colors['dark_gray']
                p.line_spacing = 1.3
    
    def _add_market_potential(self, slide, x, y, width, height):
        """Add market potential data box"""
        # Container
        market_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height
        )
        market_box.fill.solid()
        market_box.fill.fore_color.rgb = self.colors['teal']
        market_box.line.fill.background()
        
        tf = market_box.text_frame
        tf.margin_all = Inches(0.3)
        
        # Header
        p = tf.paragraphs[0]
        p.text = "📈 MARKET POTENTIAL"
        p.font.name = 'Calibri'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        # Market data
        market_data = [
            ('Healthcare market', 'USD 372B by 2025', 'CAGR 22%'),
            ('Telemedicine', 'USD 5.4B by 2025', 'Massive growth'),
            ('Diagnostics', 'CAGR ~20%', 'Underserved'),
            ('eSanjeevani', '160M+ teleconsults', 'Adoption proof')
        ]
        
        for data in market_data:
            p = tf.add_paragraph()
            p.text = ""  # Spacing
            
            p = tf.add_paragraph()
            p.text = f"▸ {data[0]}"
            p.font.name = 'Calibri'
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            
            p = tf.add_paragraph()
            p.text = f"  {data[1]} | {data[2]}"
            p.font.name = 'Calibri'
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors['white']
            p.line_spacing = 1.2
    
    def _add_healthcare_financing_pie(self, slide, x, y, width, height):
        """Add healthcare financing pie chart"""
        # Create pie chart
        fig, ax = plt.subplots(figsize=(self._inches_to_float(width), self._inches_to_float(height)))
        
        # Data
        sizes = [62, 30, 8]
        labels = ['Out-of-Pocket\n(62%)', 'Government\n(30%)', 'Private Insurance\n(8%)']
        colors = ['#DC267F', '#0096AA', '#77BB41']
        
        # Create pie
        wedges, texts, autotexts = ax.pie(sizes, labels=labels, colors=colors,
                                          autopct='', startangle=90,
                                          wedgeprops=dict(width=0.7))
        
        # Enhance text
        for text in texts:
            text.set_fontsize(10)
            text.set_fontweight('bold')
        
        # Title
        ax.text(0, 0, 'Healthcare\nFinancing', ha='center', va='center',
               fontsize=11, fontweight='bold')
        
        plt.tight_layout()
        
        # Save and add to slide
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=300, transparent=True, bbox_inches='tight')
        plt.close()
        
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, x, y, width=width)
    
    def _add_competitive_matrix(self, slide, x, y, width, height):
        """Add competitive landscape 2x2 matrix"""
        # Create matplotlib figure
        fig, ax = plt.subplots(figsize=(self._inches_to_float(width), self._inches_to_float(height)))
        
        # Set up the plot
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 10)
        
        # Draw quadrant lines
        ax.axhline(y=5, color='gray', linestyle='-', alpha=0.3, linewidth=2)
        ax.axvline(x=5, color='gray', linestyle='-', alpha=0.3, linewidth=2)
        
        # Add axis labels
        ax.text(5, -0.5, 'Geography', ha='center', fontsize=12, fontweight='bold')
        ax.text(0.5, 10.5, '← Urban', ha='left', fontsize=10)
        ax.text(9.5, 10.5, 'Rural →', ha='right', fontsize=10)
        
        ax.text(-0.5, 5, 'Service\nBreadth', ha='center', va='center', fontsize=12, 
                fontweight='bold', rotation=90)
        ax.text(-0.5, 9.5, 'Broad ↑', ha='center', fontsize=10, rotation=90)
        ax.text(-0.5, 0.5, '↓ Narrow', ha='center', fontsize=10, rotation=90)
        
        # Add competitors
        competitors = [
            ('Practo', 2, 7, '#0053A0'),
            ('Tata Health', 3, 8, '#0053A0'),
            ('1mg/PharmEasy', 2.5, 6, '#0096AA'),
            ('eSanjeevani\n(Govt)', 6, 3, '#77BB41')
        ]
        
        for name, x_pos, y_pos, color in competitors:
            # Add bubble
            circle = plt.Circle((x_pos, y_pos), 0.8, color=color, alpha=0.7)
            ax.add_patch(circle)
            ax.text(x_pos, y_pos, name, ha='center', va='center',
                   fontsize=9, fontweight='bold', color='white')
        
        # Highlight gap area
        gap_rect = plt.Rectangle((6, 6), 3, 3, fill=False, 
                                edgecolor='#DC267F', linewidth=3, linestyle='--')
        ax.add_patch(gap_rect)
        ax.text(7.5, 7.5, 'Gap:\nVernacular,\nlow-cost,\ndiagnostic\n+ trust', 
               ha='center', va='center', fontsize=10, fontweight='bold',
               color='#DC267F', style='italic')
        
        # Title
        ax.text(5, 11, 'Competitive Landscape', ha='center', fontsize=14, fontweight='bold')
        
        # Remove axes
        ax.set_xticks([])
        ax.set_yticks([])
        for spine in ax.spines.values():
            spine.set_visible(False)
        
        plt.tight_layout()
        
        # Save and add to slide
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=300, transparent=True, bbox_inches='tight')
        plt.close()
        
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, x, y, width=width)
    
    def create_slide_3_medichain_model(self):
        """Create Slide 3: MediChain Solution"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6))
        tf = title_box.text_frame
        tf.text = "MediChain — Tech-enabled Primary Care & Diagnostics"
        p = tf.paragraphs[0]
        p.font.name = 'Calibri'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
        
        # Top: 4-part solution infographic (100%)
        self._add_solution_infographic(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.5))
        
        # Middle left: Differentiators (30%)
        self._add_differentiators(slide, Inches(0.5), Inches(2.9), Inches(4), Inches(1.8))
        
        # Middle center: Funnel diagram (35%)
        self._add_funnel_diagram(slide, Inches(4.7), Inches(2.9), Inches(4), Inches(1.8))
        
        # Middle right: Impact Framework (35%)
        self._add_impact_framework(slide, Inches(8.9), Inches(2.9), Inches(3.9), Inches(1.8))
        
        # Bottom: 5-year roadmap (100%)
        self._add_roadmap(slide, Inches(0.5), Inches(4.9), Inches(12.3), Inches(1.6))
        
        # Bottom banner
        self._add_bottom_banner(slide,
            "MediChain = A vernacular, trust-first, low-cost healthcare pathway for Bharat's Tier-2/3 markets.")
        
        # Speaker notes
        self._add_speaker_notes(slide,
            "MediChain combines vernacular AI triage, kiosk-based diagnostics, secure records and pharmacy tie-ups. "
            "It pilots small (100 kiosks, 50K users) but scales to 100M lives by year five — "
            "balancing profitability and social impact.")
        
        # Slide number
        self._add_slide_number(slide, 3)
        
        return slide
    
    def _add_solution_infographic(self, slide, x, y, width, height):
        """Add 4-part solution infographic"""
        solutions = [
            {
                'icon': '🤖',
                'title': 'AI Symptom Triage',
                'desc': 'Vernacular chatbot\nTriage cost <₹20',
                'color': self.colors['teal']
            },
            {
                'icon': '🔬',
                'title': 'IoT Diagnostic Kiosks',
                'desc': 'BP, sugar, ECG, SPO₂\nKiosk cost ~₹1L',
                'color': self.colors['primary_blue']
            },
            {
                'icon': '🔗',
                'title': 'Blockchain Health Records',
                'desc': 'NDHM-aligned\nPortable, consented',
                'color': self.colors['accent_green']
            },
            {
                'icon': '💊',
                'title': 'Phygital Linkages',
                'desc': 'Local pharmacies\nLast-mile meds',
                'color': self.colors['accent_orange']
            }
        ]
        
        # Container
        solution_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, height
        )
        solution_box.fill.solid()
        solution_box.fill.fore_color.rgb = self.colors['very_light_gray']
        solution_box.line.color.rgb = self.colors['light_gray']
        solution_box.line.width = Pt(1)
        
        # Add each solution component
        comp_width = width / 4 - Inches(0.1)
        
        for i, solution in enumerate(solutions):
            x_pos = x + Inches(0.05) + (comp_width + Inches(0.1)) * i
            
            # Component box
            comp_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_pos, y + Inches(0.1),
                comp_width, height - Inches(0.2)
            )
            comp_box.fill.solid()
            comp_box.fill.fore_color.rgb = solution['color']
            comp_box.line.fill.background()
            
            tf = comp_box.text_frame
            tf.margin_all = Inches(0.15)
            
            # Icon
            p = tf.paragraphs[0]
            p.text = solution['icon']
            p.font.size = Pt(24)
            p.alignment = PP_ALIGN.CENTER
            
            # Title
            p = tf.add_paragraph()
            p.text = solution['title']
            p.font.name = 'Calibri'
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            
            # Description
            p = tf.add_paragraph()
            p.text = solution['desc']
            p.font.name = 'Calibri'
            p.font.size = Pt(9)
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.1
            
            # Add arrow between components
            if i < 3:
                arrow_x = x_pos + comp_width
                arrow = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT, 
                    arrow_x, y + height/2,
                    arrow_x + Inches(0.1), y + height/2
                )
                arrow.line.color.rgb = self.colors['medium_gray']
                arrow.line.width = Pt(2)
    
    def _add_differentiators(self, slide, x, y, width, height):
        """Add differentiators section"""
        # Container
        diff_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height
        )
        diff_box.fill.solid()
        diff_box.fill.fore_color.rgb = self.colors['primary_blue']
        diff_box.line.fill.background()
        
        tf = diff_box.text_frame
        tf.margin_all = Inches(0.2)
        
        # Header
        p = tf.paragraphs[0]
        p.text = "✨ KEY DIFFERENTIATORS"
        p.font.name = 'Calibri'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        # Differentiators
        diffs = [
            ('🗣️', 'Vernacular-first UX for Tier-2/3'),
            ('💸', 'Micro-pricing: <₹100 consults, ₹499 family plan'),
            ('🤝', 'Trust: kiosks in pharmacies, NGO/govt tie-ups')
        ]
        
        for icon, diff in diffs:
            p = tf.add_paragraph()
            p.text = ""  # Spacing
            
            p = tf.add_paragraph()
            p.text = f"{icon} {diff}"
            p.font.name = 'Calibri'
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors['white']
            p.line_spacing = 1.2
    
    def _add_funnel_diagram(self, slide, x, y, width, height):
        """Add conversion funnel diagram"""
        # Create matplotlib figure
        fig, ax = plt.subplots(figsize=(self._inches_to_float(width), self._inches_to_float(height)))
        
        # Funnel data
        stages = ['Addressable\npopulation', 'Reach via\ndigital+kiosks', 'Active users\n(6 months)', 
                  'Teleconsults', 'Paying\nsubscribers']
        values = [650000, 75000, 50000, 20000, 5000]
        
        # Create funnel
        y_positions = np.linspace(1, 0, len(stages))
        
        for i, (stage, value, y_pos) in enumerate(zip(stages, values, y_positions)):
            # Calculate width based on value
            max_width = 0.8
            width = max_width * (value / values[0])
            
            # Create trapezoid (funnel segment)
            if i < len(stages) - 1:
                next_width = max_width * (values[i+1] / values[0])
                x_coords = [(0.5 - width/2), (0.5 + width/2), 
                           (0.5 + next_width/2), (0.5 - next_width/2)]
                y_coords = [y_pos, y_pos, y_positions[i+1], y_positions[i+1]]
            else:
                x_coords = [(0.5 - width/2), (0.5 + width/2), 0.5, 0.5]
                y_coords = [y_pos, y_pos, y_pos - 0.1, y_pos - 0.1]
            
            # Draw segment
            color = plt.cm.Blues(0.3 + 0.5 * (i / len(stages)))
            ax.fill(x_coords, y_coords, color=color, alpha=0.8, edgecolor='darkblue', linewidth=2)
            
            # Add text
            ax.text(0.5, y_pos - 0.05, f"{stage}\n{value:,}", 
                   ha='center', va='center', fontsize=9, fontweight='bold')
        
        # Title
        ax.text(0.5, 1.1, 'Pilot Conversion Funnel', ha='center', fontsize=12, fontweight='bold')
        
        # Remove axes
        ax.set_xlim(0, 1)
        ax.set_ylim(-0.1, 1.2)
        ax.axis('off')
        
        plt.tight_layout()
        
        # Save and add to slide
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=300, transparent=True, bbox_inches='tight')
        plt.close()
        
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, x, y, width=width)
    
    def _add_impact_framework(self, slide, x, y, width, height):
        """Add impact framework"""
        # Container
        impact_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height
        )
        impact_box.fill.solid()
        impact_box.fill.fore_color.rgb = self.colors['white']
        impact_box.line.color.rgb = self.colors['accent_green']
        impact_box.line.width = Pt(2)
        
        # Header
        header_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, Inches(0.4)
        )
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = self.colors['accent_green']
        header_box.line.fill.background()
        
        tf = header_box.text_frame
        tf.text = "IMPACT & SCALE"
        p = tf.paragraphs[0]
        p.font.name = 'Calibri'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # Two columns
        impacts = [
            {
                'title': '💰 Economic Impact',
                'items': ['Affordable care', 'Lower OOP burden', 'Scalable subscription']
            },
            {
                'title': '🌍 Social Impact',
                'items': ['100M underserved', 'Preventive adoption', 'SDG-3 aligned']
            }
        ]
        
        col_width = width / 2 - Inches(0.2)
        y_start = y + Inches(0.5)
        
        for i, impact in enumerate(impacts):
            x_pos = x + Inches(0.1) + (col_width + Inches(0.2)) * i
            
            # Column content
            col_box = slide.shapes.add_textbox(x_pos, y_start, col_width, height - Inches(0.6))
            tf = col_box.text_frame
            
            # Title
            p = tf.paragraphs[0]
            p.text = impact['title']
            p.font.name = 'Calibri'
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.colors['dark_gray']
            
            # Items
            for item in impact['items']:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.name = 'Calibri'
                p.font.size = Pt(9)
                p.font.color.rgb = self.colors['medium_gray']
                p.line_spacing = 1.2
    
    def _add_roadmap(self, slide, x, y, width, height):
        """Add 5-year roadmap with arrows"""
        # Create matplotlib figure
        fig, ax = plt.subplots(figsize=(self._inches_to_float(width), self._inches_to_float(height)))
        
        # Roadmap phases
        phases = [
            {'year': 'Y1', 'target': '100 kiosks\n50K users', 'desc': 'Pilot\n5 districts'},
            {'year': 'Y3', 'target': '25 cities\n1M users', 'desc': 'Scale'},
            {'year': 'Y5', 'target': 'Pan-India\n100M lives', 'desc': 'National rollout'}
        ]
        
        # Timeline base
        timeline_y = 0.5
        ax.arrow(0.1, timeline_y, 0.85, 0, head_width=0.05, head_length=0.02, 
                fc='#0053A0', ec='#0053A0', linewidth=3)
        
        # Add phases
        x_positions = [0.2, 0.5, 0.8]
        
        for i, (phase, x_pos) in enumerate(zip(phases, x_positions)):
            # Phase marker
            circle = plt.Circle((x_pos, timeline_y), 0.06, color='#0096AA', zorder=5)
            ax.add_patch(circle)
            ax.text(x_pos, timeline_y, phase['year'], ha='center', va='center',
                   fontsize=10, fontweight='bold', color='white', zorder=6)
            
            # Phase details box
            box_y = 0.7 if i % 2 == 0 else 0.3
            
            # Connection line
            ax.plot([x_pos, x_pos], [timeline_y, box_y], 'k--', alpha=0.5, linewidth=1)
            
            # Details box
            bbox = dict(boxstyle="round,pad=0.3", facecolor='#E8F4F8', edgecolor='#0096AA', linewidth=2)
            ax.text(x_pos, box_y, f"{phase['target']}\n{phase['desc']}", 
                   ha='center', va='center', fontsize=9, bbox=bbox)
        
        # Title
        ax.text(0.5, 0.9, '5-Year Scale Roadmap', ha='center', fontsize=12, fontweight='bold')
        
        # Remove axes
        ax.set_xlim(0, 1)
        ax.set_ylim(0.1, 1)
        ax.axis('off')
        
        plt.tight_layout()
        
        # Save and add to slide
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=300, transparent=True, bbox_inches='tight')
        plt.close()
        
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, x, y, width=width)
    
    def save_presentation(self, filename="DisruptX_Round1_MediChain.pptx"):
        """Save the presentation"""
        # Create all slides
        self.create_slide_1_opportunity_landscape()
        self.create_slide_2_sector_focus()
        self.create_slide_3_medichain_model()
        
        # Save
        self.prs.save(filename)
        return filename


def create_disruptx_presentation():
    """Create the DisruptX presentation"""
    ppt = DisruptXPresentation()
    filename = ppt.save_presentation()
    print(f"✓ Created DisruptX presentation: {filename}")
    print("✓ 3 slides with dense visuals")
    print("✓ BCG/McKinsey consulting style")
    print("✓ Speaker notes included")
    print("✓ Professional blue/teal/grey theme")
    return filename


if __name__ == "__main__":
    print("Creating DisruptX Round 1 - MediChain presentation...")
    create_disruptx_presentation()
    print("\n✓ Presentation ready for competition!")