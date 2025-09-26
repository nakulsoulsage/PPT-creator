#!/usr/bin/env python3
"""
Professional Case Competition PPT System
Creates ultra-clean, informative, McKinsey-quality presentations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd
import numpy as np
from datetime import datetime
import io
import os

class ProfessionalPPTSystem:
    def __init__(self, style='mckinsey'):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 widescreen
        self.prs.slide_height = Inches(7.5)
        
        # Professional color schemes
        self.styles = {
            'mckinsey': {
                'primary': RGBColor(0, 32, 96),      # Navy Blue
                'secondary': RGBColor(0, 102, 204),  # McKinsey Blue
                'accent': RGBColor(0, 145, 220),     # Light Blue
                'success': RGBColor(0, 164, 153),    # Teal
                'warning': RGBColor(255, 186, 8),    # Gold
                'danger': RGBColor(211, 47, 47),     # Red
                'text': RGBColor(51, 51, 51),        # Dark Gray
                'subtext': RGBColor(117, 117, 117),  # Medium Gray
                'background': RGBColor(248, 248, 248), # Light Gray
                'white': RGBColor(255, 255, 255)
            },
            'bcg': {
                'primary': RGBColor(0, 128, 0),      # BCG Green
                'secondary': RGBColor(0, 155, 119),  # Light Green
                'accent': RGBColor(0, 176, 80),      # Bright Green
                'success': RGBColor(76, 175, 80),    # Success Green
                'warning': RGBColor(255, 152, 0),    # Orange
                'danger': RGBColor(244, 67, 54),     # Red
                'text': RGBColor(33, 33, 33),        # Almost Black
                'subtext': RGBColor(97, 97, 97),     # Dark Gray
                'background': RGBColor(250, 250, 250), # Off White
                'white': RGBColor(255, 255, 255)
            },
            'bain': {
                'primary': RGBColor(237, 28, 36),    # Bain Red
                'secondary': RGBColor(255, 102, 102), # Light Red
                'accent': RGBColor(255, 138, 128),    # Coral
                'success': RGBColor(102, 187, 106),   # Green
                'warning': RGBColor(255, 167, 38),    # Orange
                'danger': RGBColor(229, 57, 53),      # Dark Red
                'text': RGBColor(66, 66, 66),         # Dark Gray
                'subtext': RGBColor(117, 117, 117),   # Medium Gray
                'background': RGBColor(253, 253, 253), # Near White
                'white': RGBColor(255, 255, 255)
            }
        }
        
        self.colors = self.styles[style]
        self.style = style
        
    def _add_slide_number(self, slide, number=None):
        """Add slide number to bottom right"""
        if number is None:
            number = len(self.prs.slides)
            
        textbox = slide.shapes.add_textbox(
            Inches(12.5), Inches(7), Inches(0.5), Inches(0.3)
        )
        tf = textbox.text_frame
        tf.text = str(number)
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = self.colors['subtext']
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
    def _add_logo_placeholder(self, slide):
        """Add logo/branding area"""
        logo_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(11.5), Inches(0.2), Inches(1.5), Inches(0.5)
        )
        logo_box.fill.background()
        logo_box.line.fill.background()
        text_frame = logo_box.text_frame
        text_frame.text = "[Logo]"
        text_frame.paragraphs[0].font.size = Pt(8)
        text_frame.paragraphs[0].font.color.rgb = self.colors['subtext']
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
    def add_title_slide(self, title, subtitle, date=None, presenters=None):
        """Create minimalist professional title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.text = title.upper()
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        p.alignment = PP_ALIGN.LEFT
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(11), Inches(1)
        )
        tf = subtitle_box.text_frame
        tf.text = subtitle
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(24)
        p.font.color.rgb = self.colors['text']
        p.alignment = PP_ALIGN.LEFT
        
        # Date
        if date is None:
            date = datetime.now().strftime("%B %Y")
        date_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5), Inches(4), Inches(0.5)
        )
        tf = date_box.text_frame
        tf.text = date
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.color.rgb = self.colors['subtext']
        
        # Presenters
        if presenters:
            presenters_box = slide.shapes.add_textbox(
                Inches(1), Inches(5.5), Inches(11), Inches(0.5)
            )
            tf = presenters_box.text_frame
            tf.text = " | ".join(presenters)
            p = tf.paragraphs[0]
            p.font.name = 'Arial'
            p.font.size = Pt(12)
            p.font.color.rgb = self.colors['subtext']
            
        # Add subtle accent line
        line = slide.shapes.add_connector(
            1, Inches(1), Inches(5.2), Inches(7), Inches(5.2)
        )
        line.line.color.rgb = self.colors['accent']
        line.line.width = Pt(2)
        
        return slide
    
    def add_agenda_slide(self, sections):
        """Create clean agenda slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(11), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.text = "AGENDA"
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        # Sections
        y_position = 2
        for i, section in enumerate(sections, 1):
            # Number circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, 
                Inches(1), Inches(y_position - 0.15), 
                Inches(0.5), Inches(0.5)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.colors['accent']
            circle.line.fill.background()
            
            tf = circle.text_frame
            tf.text = str(i)
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            
            # Section text
            text_box = slide.shapes.add_textbox(
                Inches(1.8), Inches(y_position - 0.1), Inches(9), Inches(0.5)
            )
            tf = text_box.text_frame
            tf.text = section['title']
            p = tf.paragraphs[0]
            p.font.name = 'Arial'
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['text']
            
            # Description
            if 'description' in section:
                desc_box = slide.shapes.add_textbox(
                    Inches(1.8), Inches(y_position + 0.4), Inches(9), Inches(0.4)
                )
                tf = desc_box.text_frame
                tf.text = section['description']
                p = tf.paragraphs[0]
                p.font.name = 'Arial'
                p.font.size = Pt(12)
                p.font.color.rgb = self.colors['subtext']
                
            y_position += 1.2
            
        self._add_slide_number(slide)
        return slide
    
    def add_key_message_slide(self, title, key_message, supporting_points):
        """Create slide with key message and supporting points"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(11), Inches(0.6)
        )
        tf = title_box.text_frame
        tf.text = title.upper()
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        # Key message box
        key_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0.5), Inches(1.5), 
            Inches(12.3), Inches(1.2)
        )
        key_box.fill.solid()
        key_box.fill.fore_color.rgb = self.colors['background']
        key_box.line.color.rgb = self.colors['accent']
        key_box.line.width = Pt(2)
        
        tf = key_box.text_frame
        tf.margin_left = Inches(0.5)
        tf.margin_right = Inches(0.5)
        tf.margin_top = Inches(0.25)
        tf.margin_bottom = Inches(0.25)
        tf.text = key_message
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # Supporting points
        y_position = 3.5
        for i, point in enumerate(supporting_points):
            # Bullet
            bullet = slide.shapes.add_shape(
                MSO_SHAPE.DIAMOND, 
                Inches(1.5), Inches(y_position + 0.1), 
                Inches(0.15), Inches(0.15)
            )
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = self.colors['accent']
            bullet.line.fill.background()
            
            # Text
            text_box = slide.shapes.add_textbox(
                Inches(2), Inches(y_position), Inches(10), Inches(0.6)
            )
            tf = text_box.text_frame
            tf.text = point
            p = tf.paragraphs[0]
            p.font.name = 'Arial'
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors['text']
            p.line_spacing = 1.2
            
            y_position += 0.8
            
        self._add_slide_number(slide)
        return slide
    
    def add_data_slide(self, title, chart_type, data, insights=None):
        """Create slide with data visualization and insights"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(11), Inches(0.6)
        )
        tf = title_box.text_frame
        tf.text = title.upper()
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        # Create visualization
        if chart_type == 'bar':
            self._add_bar_visualization(slide, data)
        elif chart_type == 'line':
            self._add_line_visualization(slide, data)
        elif chart_type == 'pie':
            self._add_pie_visualization(slide, data)
        elif chart_type == 'waterfall':
            self._add_waterfall_visualization(slide, data)
        elif chart_type == 'matrix':
            self._add_matrix_visualization(slide, data)
            
        # Add insights box
        if insights:
            insights_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                Inches(8.5), Inches(2), 
                Inches(4), Inches(3)
            )
            insights_box.fill.solid()
            insights_box.fill.fore_color.rgb = self.colors['background']
            insights_box.line.color.rgb = self.colors['subtext']
            insights_box.line.width = Pt(0.5)
            
            tf = insights_box.text_frame
            tf.margin_left = Inches(0.3)
            tf.margin_right = Inches(0.3)
            tf.margin_top = Inches(0.3)
            
            p = tf.paragraphs[0]
            p.text = "KEY INSIGHTS"
            p.font.name = 'Arial'
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.colors['primary']
            
            for insight in insights:
                p = tf.add_paragraph()
                p.text = f"• {insight}"
                p.font.name = 'Arial'
                p.font.size = Pt(11)
                p.font.color.rgb = self.colors['text']
                p.line_spacing = 1.3
                
        self._add_slide_number(slide)
        return slide
    
    def _add_bar_visualization(self, slide, data):
        """Add clean bar chart"""
        # Create matplotlib figure
        fig, ax = plt.subplots(figsize=(7, 4))
        
        # Data
        categories = data['categories']
        values = data['values']
        
        # Create bars
        bars = ax.bar(categories, values, color='#0066CC', width=0.6)
        
        # Add value labels on bars
        for bar in bars:
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height,
                   f'{height:,.0f}',
                   ha='center', va='bottom', fontsize=10, fontweight='bold')
        
        # Styling
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_visible(False)
        ax.yaxis.set_visible(False)
        ax.xaxis.set_ticks_position('none')
        
        # Grid
        ax.yaxis.grid(True, linestyle='-', alpha=0.2)
        ax.set_axisbelow(True)
        
        # Labels
        ax.set_xlabel(data.get('x_label', ''), fontsize=12, fontweight='bold')
        plt.xticks(rotation=0)
        plt.tight_layout()
        
        # Save and add to slide
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=300, bbox_inches='tight', facecolor='white')
        plt.close()
        
        img_stream.seek(0)
        pic = slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), width=Inches(7))
        
    def _add_line_visualization(self, slide, data):
        """Add clean line chart"""
        fig, ax = plt.subplots(figsize=(7, 4))
        
        # Plot lines
        for series in data['series']:
            ax.plot(data['x_values'], series['values'], 
                   marker='o', linewidth=2.5, markersize=8,
                   label=series['name'], color=series.get('color', '#0066CC'))
        
        # Styling
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.grid(True, linestyle='-', alpha=0.2)
        ax.set_axisbelow(True)
        
        # Labels
        ax.set_xlabel(data.get('x_label', ''), fontsize=12, fontweight='bold')
        ax.set_ylabel(data.get('y_label', ''), fontsize=12, fontweight='bold')
        
        if len(data['series']) > 1:
            ax.legend(frameon=False, loc='best')
            
        plt.tight_layout()
        
        # Save and add to slide
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=300, bbox_inches='tight', facecolor='white')
        plt.close()
        
        img_stream.seek(0)
        pic = slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), width=Inches(7))
        
    def _add_pie_visualization(self, slide, data):
        """Add clean pie chart"""
        fig, ax = plt.subplots(figsize=(6, 6))
        
        # Create pie
        colors = ['#0066CC', '#0099FF', '#66B2FF', '#99CCFF', '#CCE5FF']
        wedges, texts, autotexts = ax.pie(
            data['values'], 
            labels=data['labels'],
            autopct='%1.1f%%',
            colors=colors[:len(data['values'])],
            startangle=90,
            pctdistance=0.85
        )
        
        # Styling
        for text in texts:
            text.set_fontsize(12)
            text.set_fontweight('bold')
            
        for autotext in autotexts:
            autotext.set_color('white')
            autotext.set_fontsize(11)
            autotext.set_fontweight('bold')
            
        # Add center circle for donut effect
        centre_circle = plt.Circle((0, 0), 0.70, fc='white')
        ax.add_artist(centre_circle)
        
        plt.tight_layout()
        
        # Save and add to slide
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=300, bbox_inches='tight', facecolor='white')
        plt.close()
        
        img_stream.seek(0)
        pic = slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), width=Inches(6))
        
    def add_comparison_slide(self, title, items):
        """Create comparison slide with multiple columns"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(11), Inches(0.6)
        )
        tf = title_box.text_frame
        tf.text = title.upper()
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        # Calculate column width
        num_items = len(items)
        total_width = 11.5
        gap = 0.3
        col_width = (total_width - (gap * (num_items - 1))) / num_items
        
        # Create comparison columns
        for i, item in enumerate(items):
            x_position = 1 + (col_width + gap) * i
            
            # Header box
            header = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x_position), Inches(1.5),
                Inches(col_width), Inches(0.8)
            )
            header.fill.solid()
            header.fill.fore_color.rgb = self.colors['primary']
            header.line.fill.background()
            
            tf = header.text_frame
            tf.text = item['name']
            p = tf.paragraphs[0]
            p.font.name = 'Arial'
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            
            # Content box
            content = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x_position), Inches(2.3),
                Inches(col_width), Inches(4.5)
            )
            content.fill.solid()
            content.fill.fore_color.rgb = self.colors['background']
            content.line.color.rgb = self.colors['subtext']
            content.line.width = Pt(0.5)
            
            tf = content.text_frame
            tf.margin_left = Inches(0.2)
            tf.margin_right = Inches(0.2)
            tf.margin_top = Inches(0.2)
            
            for j, point in enumerate(item['points']):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = f"• {point}"
                p.font.name = 'Arial'
                p.font.size = Pt(11)
                p.font.color.rgb = self.colors['text']
                p.line_spacing = 1.3
                
        self._add_slide_number(slide)
        return slide
    
    def add_framework_slide(self, title, framework_type, data):
        """Create framework slide (2x2 matrix, Porter's 5 Forces, etc.)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(11), Inches(0.6)
        )
        tf = title_box.text_frame
        tf.text = title.upper()
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        if framework_type == '2x2_matrix':
            self._add_2x2_matrix(slide, data)
        elif framework_type == 'porter5':
            self._add_porter5_forces(slide, data)
        elif framework_type == 'value_chain':
            self._add_value_chain(slide, data)
            
        self._add_slide_number(slide)
        return slide
    
    def _add_2x2_matrix(self, slide, data):
        """Add 2x2 matrix framework"""
        # Matrix center
        center_x = 6.7
        center_y = 4
        box_width = 3.5
        box_height = 2
        
        # Draw axes
        h_line = slide.shapes.add_connector(
            1, Inches(2), Inches(center_y), Inches(11.3), Inches(center_y)
        )
        h_line.line.color.rgb = self.colors['text']
        h_line.line.width = Pt(1)
        
        v_line = slide.shapes.add_connector(
            1, Inches(center_x), Inches(2), Inches(center_x), Inches(6)
        )
        v_line.line.color.rgb = self.colors['text']
        v_line.line.width = Pt(1)
        
        # Add axis labels
        # X-axis label
        x_label = slide.shapes.add_textbox(
            Inches(5), Inches(6.2), Inches(3.5), Inches(0.3)
        )
        tf = x_label.text_frame
        tf.text = data['x_axis']
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors['text']
        p.alignment = PP_ALIGN.CENTER
        
        # Y-axis label (rotated effect using positioning)
        y_label = slide.shapes.add_textbox(
            Inches(1.5), Inches(3.8), Inches(0.5), Inches(0.4)
        )
        tf = y_label.text_frame
        tf.text = data['y_axis']
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors['text']
        
        # Add quadrants
        quadrants = [
            {'x': center_x - box_width - 0.2, 'y': center_y - box_height - 0.2, 'data': data['quadrants'][0]},  # Top Left
            {'x': center_x + 0.2, 'y': center_y - box_height - 0.2, 'data': data['quadrants'][1]},             # Top Right
            {'x': center_x - box_width - 0.2, 'y': center_y + 0.2, 'data': data['quadrants'][2]},              # Bottom Left
            {'x': center_x + 0.2, 'y': center_y + 0.2, 'data': data['quadrants'][3]}                           # Bottom Right
        ]
        
        colors = [self.colors['accent'], self.colors['success'], 
                  self.colors['warning'], self.colors['secondary']]
        
        for i, quad in enumerate(quadrants):
            box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(quad['x']), Inches(quad['y']),
                Inches(box_width), Inches(box_height)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = colors[i]
            box.fill.transparency = 0.8  # Make it lighter
            box.line.color.rgb = colors[i]
            box.line.width = Pt(1)
            
            tf = box.text_frame
            tf.margin_left = Inches(0.2)
            tf.margin_right = Inches(0.2)
            tf.margin_top = Inches(0.2)
            
            # Title
            p = tf.paragraphs[0]
            p.text = quad['data']['title']
            p.font.name = 'Arial'
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.colors['text']
            p.alignment = PP_ALIGN.CENTER
            
            # Items
            for item in quad['data']['items'][:3]:  # Limit to 3 items
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.name = 'Arial'
                p.font.size = Pt(10)
                p.font.color.rgb = self.colors['text']
    
    def add_recommendation_slide(self, title, recommendations, implementation_timeline=None):
        """Create recommendation slide with clear action items"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(11), Inches(0.6)
        )
        tf = title_box.text_frame
        tf.text = title.upper()
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        # Recommendations
        y_position = 1.8
        for i, rec in enumerate(recommendations, 1):
            # Number box
            num_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1), Inches(y_position),
                Inches(0.6), Inches(0.6)
            )
            num_box.fill.solid()
            num_box.fill.fore_color.rgb = self.colors['primary']
            num_box.line.fill.background()
            
            tf = num_box.text_frame
            tf.text = str(i)
            p = tf.paragraphs[0]
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            
            # Recommendation text
            rec_box = slide.shapes.add_textbox(
                Inches(1.8), Inches(y_position), Inches(10), Inches(1.2)
            )
            tf = rec_box.text_frame
            
            # Title
            p = tf.paragraphs[0]
            p.text = rec['title']
            p.font.name = 'Arial'
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.colors['text']
            
            # Description
            if 'description' in rec:
                p = tf.add_paragraph()
                p.text = rec['description']
                p.font.name = 'Arial'
                p.font.size = Pt(12)
                p.font.color.rgb = self.colors['subtext']
                p.line_spacing = 1.2
                
            # Impact
            if 'impact' in rec:
                p = tf.add_paragraph()
                p.text = f"Impact: {rec['impact']}"
                p.font.name = 'Arial'
                p.font.size = Pt(11)
                p.font.color.rgb = self.colors['success']
                p.font.italic = True
                
            y_position += 1.5
            
        # Timeline
        if implementation_timeline:
            timeline_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1), Inches(5.5),
                Inches(11.3), Inches(1.5)
            )
            timeline_box.fill.solid()
            timeline_box.fill.fore_color.rgb = self.colors['background']
            timeline_box.line.color.rgb = self.colors['accent']
            timeline_box.line.width = Pt(1)
            
            tf = timeline_box.text_frame
            tf.margin_left = Inches(0.3)
            tf.margin_top = Inches(0.2)
            
            p = tf.paragraphs[0]
            p.text = "IMPLEMENTATION TIMELINE"
            p.font.name = 'Arial'
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.colors['primary']
            
            p = tf.add_paragraph()
            p.text = implementation_timeline
            p.font.name = 'Arial'
            p.font.size = Pt(11)
            p.font.color.rgb = self.colors['text']
            
        self._add_slide_number(slide)
        return slide
    
    def add_summary_slide(self, key_takeaways):
        """Create summary slide with key takeaways"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(11), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.text = "KEY TAKEAWAYS"
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        # Create takeaway boxes
        y_position = 1.8
        colors = [self.colors['primary'], self.colors['secondary'], self.colors['accent']]
        
        for i, takeaway in enumerate(key_takeaways[:3]):  # Limit to 3
            # Box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1), Inches(y_position),
                Inches(11.3), Inches(1.5)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = colors[i % len(colors)]
            box.line.fill.background()
            
            tf = box.text_frame
            tf.margin_left = Inches(0.5)
            tf.margin_right = Inches(0.5)
            tf.margin_top = Inches(0.3)
            tf.margin_bottom = Inches(0.3)
            
            # Number
            p = tf.paragraphs[0]
            p.text = f"{i + 1}. {takeaway}"
            p.font.name = 'Arial'
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            
            y_position += 1.8
            
        self._add_slide_number(slide)
        return slide
    
    def add_appendix_divider(self):
        """Create appendix divider slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['primary']
        bg.line.fill.background()
        
        # Text
        text_box = slide.shapes.add_textbox(
            Inches(1), Inches(3), Inches(11.3), Inches(1.5)
        )
        tf = text_box.text_frame
        tf.text = "APPENDIX"
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def save(self, filename):
        """Save the presentation"""
        self.prs.save(filename)
        return filename


def create_complete_case_presentation():
    """Create a complete, professional case competition presentation"""
    
    # Initialize presentation
    ppt = ProfessionalPPTSystem(style='mckinsey')
    
    # 1. Title Slide
    ppt.add_title_slide(
        title="Digital Banking Transformation",
        subtitle="Capturing $2.5B opportunity through customer-centric innovation",
        presenters=["SIMSREE Team Alpha", "Mumbai"]
    )
    
    # 2. Agenda
    ppt.add_agenda_slide([
        {"title": "Executive Summary", "description": "Problem statement and proposed solution"},
        {"title": "Market Analysis", "description": "Industry trends and competitive landscape"},
        {"title": "Strategic Framework", "description": "Our approach to digital transformation"},
        {"title": "Implementation Roadmap", "description": "Phased approach with quick wins"},
        {"title": "Financial Impact", "description": "ROI analysis and value creation"},
        {"title": "Recommendations", "description": "Next steps and success factors"}
    ])
    
    # 3. Executive Summary
    ppt.add_key_message_slide(
        title="Executive Summary",
        key_message="Traditional banks must transform digitally or lose 40% market share to fintech competitors by 2026",
        supporting_points=[
            "Customer expectations have fundamentally shifted - 78% prefer digital-first banking",
            "Fintech disruption accelerating with $150B invested globally in 2023",
            "Early movers capturing 3x revenue growth vs. traditional peers",
            "Implementation window closing - first-mover advantage critical"
        ]
    )
    
    # 4. Market Analysis
    market_data = {
        'categories': ['Traditional Banks', 'Neo Banks', 'Fintech Apps', 'Big Tech', 'Others'],
        'values': [45, 15, 20, 12, 8],
        'x_label': 'Market Segments'
    }
    ppt.add_data_slide(
        title="Market Share Evolution",
        chart_type='bar',
        data=market_data,
        insights=[
            "Traditional banks losing 2% share annually",
            "Neo banks growing at 45% CAGR",
            "Big Tech entry accelerating disruption"
        ]
    )
    
    # 5. Customer Journey Analysis
    ppt.add_comparison_slide(
        title="Customer Experience Comparison",
        items=[
            {
                "name": "Traditional Banking",
                "points": [
                    "Branch-centric model",
                    "5-7 days account opening",
                    "Limited digital features",
                    "9am-5pm availability",
                    "High operational costs"
                ]
            },
            {
                "name": "Digital-First Banking",
                "points": [
                    "Mobile-native experience",
                    "5-minute account opening",
                    "AI-powered insights",
                    "24/7 instant support",
                    "70% lower cost-to-serve"
                ]
            },
            {
                "name": "Our Transformation",
                "points": [
                    "Hybrid optimal model",
                    "Instant digital onboarding",
                    "Personalized offerings",
                    "Omnichannel excellence",
                    "50% cost reduction"
                ]
            }
        ]
    )
    
    # 6. Strategic Framework
    matrix_data = {
        'x_axis': 'Digital Maturity →',
        'y_axis': 'Customer Value ↑',
        'quadrants': [
            {
                'title': 'TRANSFORM',
                'items': ['Core Banking', 'Lending Platform', 'Risk Systems']
            },
            {
                'title': 'ACCELERATE',
                'items': ['Mobile Banking', 'Payment Systems', 'Analytics']
            },
            {
                'title': 'OPTIMIZE',
                'items': ['Branch Network', 'Call Centers', 'Back Office']
            },
            {
                'title': 'INNOVATE',
                'items': ['AI Advisors', 'Blockchain', 'Open Banking']
            }
        ]
    }
    ppt.add_framework_slide(
        title="Digital Transformation Framework",
        framework_type='2x2_matrix',
        data=matrix_data
    )
    
    # 7. Implementation Roadmap
    timeline_data = {
        'categories': ['Q1 2024', 'Q2 2024', 'Q3 2024', 'Q4 2024', 'Q1 2025', 'Q2 2025'],
        'values': [100, 250, 450, 750, 1200, 1800],
        'x_label': 'Timeline'
    }
    ppt.add_data_slide(
        title="Digital Adoption Trajectory",
        chart_type='line',
        data={
            'x_values': timeline_data['categories'],
            'series': [
                {
                    'name': 'Active Digital Users (000s)',
                    'values': timeline_data['values'],
                    'color': '#0066CC'
                }
            ],
            'x_label': 'Implementation Timeline',
            'y_label': 'Users (thousands)'
        },
        insights=[
            "Month 1-3: Foundation & pilot launch",
            "Month 4-6: Scale to 25% customers",
            "Month 7-12: Full rollout & optimization"
        ]
    )
    
    # 8. Financial Analysis
    roi_data = {
        'labels': ['Technology Investment', 'Process Optimization', 'Revenue Growth', 'Cost Savings'],
        'values': [35, 20, 30, 15]
    }
    ppt.add_data_slide(
        title="Value Creation Breakdown",
        chart_type='pie',
        data=roi_data,
        insights=[
            "$2.5B total value creation over 3 years",
            "280% ROI with 14-month payback",
            "Break-even at Month 9"
        ]
    )
    
    # 9. Recommendations
    ppt.add_recommendation_slide(
        title="Strategic Recommendations",
        recommendations=[
            {
                'title': 'Launch Digital Transformation Office',
                'description': 'Establish dedicated team with C-suite sponsorship to drive end-to-end transformation',
                'impact': 'Accelerate delivery by 40%'
            },
            {
                'title': 'Partner with Leading Fintech Platforms',
                'description': 'Strategic partnerships for payments, lending, and wealth management capabilities',
                'impact': 'Reduce time-to-market by 18 months'
            },
            {
                'title': 'Invest in Data & AI Capabilities',
                'description': 'Build advanced analytics platform for personalization and risk management',
                'impact': 'Increase revenue per customer by 35%'
            }
        ],
        implementation_timeline="Phase 1 (0-6 months): Foundation | Phase 2 (6-12 months): Scale | Phase 3 (12-24 months): Optimize"
    )
    
    # 10. Key Takeaways
    ppt.add_summary_slide([
        "Digital transformation is existential - act now or lose relevance",
        "$2.5B value creation opportunity with proven 280% ROI",
        "Customer-centric approach with phased implementation ensures success"
    ])
    
    # Save presentation
    filename = "Professional_Case_Competition_Presentation.pptx"
    ppt.save(filename)
    print(f"\n✓ Created professional presentation: {filename}")
    print("✓ Total slides: 10")
    print("✓ Style: McKinsey professional")
    
    return filename


if __name__ == "__main__":
    print("Creating ultra-professional case competition presentation...")
    create_complete_case_presentation()
    print("\n✓ Presentation ready for competition!")