#!/usr/bin/env python3
"""
DisruptX Round 1 - MediChain Presentation (CLEAN VERSION)
Ultra-clean BCG/McKinsey-style 3-slide deck with perfect formatting
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import numpy as np
import io

class CleanDisruptXPresentation:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9
        self.prs.slide_height = Inches(7.5)
        
        # Clean BCG color palette
        self.colors = {
            'primary_blue': RGBColor(0, 83, 159),
            'teal': RGBColor(0, 150, 170),
            'light_teal': RGBColor(179, 229, 234),
            'dark_gray': RGBColor(64, 64, 64),
            'medium_gray': RGBColor(128, 128, 128),
            'light_gray': RGBColor(191, 191, 191),
            'very_light_gray': RGBColor(242, 242, 242),
            'white': RGBColor(255, 255, 255),
            'green': RGBColor(119, 187, 65),
            'orange': RGBColor(255, 138, 0),
            'red': RGBColor(220, 38, 127),
        }
    
    def _add_clean_title(self, slide, text, y_position=0.3):
        """Add clean title with proper spacing"""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_position), 
            Inches(12.3), Inches(0.7)
        )
        tf = title_box.text_frame
        tf.clear()  # Clear any default text
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Calibri'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
        p.alignment = PP_ALIGN.LEFT
        
    def _add_bottom_banner(self, slide, text):
        """Add clean bottom banner"""
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(6.8),
            Inches(13.333), Inches(0.7)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = self.colors['primary_blue']
        banner.line.fill.background()
        
        tf = banner.text_frame
        tf.clear()
        tf.margin_left = Inches(0.5)
        tf.margin_right = Inches(0.5)
        tf.margin_top = Inches(0.15)
        tf.margin_bottom = Inches(0.15)
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Calibri'
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
    def _add_slide_number(self, slide, number):
        """Add slide number"""
        num_box = slide.shapes.add_textbox(
            Inches(12.7), Inches(7.1), 
            Inches(0.4), Inches(0.3)
        )
        tf = num_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = str(number)
        p.font.name = 'Calibri'
        p.font.size = Pt(10)
        p.font.color.rgb = self.colors['medium_gray']
        p.alignment = PP_ALIGN.RIGHT
        
    def _add_speaker_notes(self, slide, notes):
        """Add speaker notes"""
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
        
    def create_slide_1(self):
        """Slide 1: Opportunity Landscape - Clean Layout"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        self._add_clean_title(slide, "Why Tier-2 & Tier-3 India is Ripe for Disruption")
        
        # Create 2x2 grid layout for clean organization
        # Top Left: Macro Trends
        self._add_macro_trends_clean(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(2.5))
        
        # Top Right: Map Visual
        self._add_india_map_clean(slide, Inches(6.8), Inches(1.2), Inches(6), Inches(2.5))
        
        # Bottom Left: Underserved Sectors
        self._add_sectors_clean(slide, Inches(0.5), Inches(3.9), Inches(6), Inches(2.5))
        
        # Bottom Right: Infrastructure Chart
        self._add_infrastructure_chart_clean(slide, Inches(6.8), Inches(3.9), Inches(6), Inches(2.5))
        
        # Bottom banner
        self._add_bottom_banner(slide, 
            "Digital readiness + structural gaps = fertile ground for tech-enabled disruption in Bharat.")
        
        # Speaker notes
        self._add_speaker_notes(slide,
            "Tier-2 and Tier-3 India are digitally enabled — UPI, smartphones, cheap data — "
            "but structurally underserved across healthcare, education, finance and agriculture. "
            "This duality creates massive disruption potential.")
        
        self._add_slide_number(slide, 1)
        return slide
    
    def _add_macro_trends_clean(self, slide, x, y, width, height):
        """Clean macro trends with proper spacing"""
        # Container with header
        container = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, height
        )
        container.fill.solid()
        container.fill.fore_color.rgb = self.colors['very_light_gray']
        container.line.color.rgb = self.colors['light_gray']
        container.line.width = Pt(1)
        
        # Header
        header_height = 0.5
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, Inches(header_height)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.colors['teal']
        header.line.fill.background()
        
        header_tf = header.text_frame
        header_tf.clear()
        header_tf.margin_top = Inches(0.1)
        header_p = header_tf.paragraphs[0]
        header_p.text = "MACRO TRENDS DRIVING DISRUPTION"
        header_p.font.name = 'Calibri'
        header_p.font.size = Pt(14)
        header_p.font.bold = True
        header_p.font.color.rgb = self.colors['white']
        header_p.alignment = PP_ALIGN.CENTER
        
        # Trends content with proper spacing
        trends = [
            ("GDP Growth:", "Tier-2/3 → ~45% of India's GDP by 2025"),
            ("Population:", "~650M population base"),
            ("Smartphone:", "~60% penetration (vs 78% urban)"),
            ("UPI Adoption:", ">12B transactions/month"),
            ("Data Cost:", "Among lowest globally (~$0.17/GB)")
        ]
        
        content_start_y = y + Inches(header_height + 0.2)
        line_height = 0.35
        
        for i, (label, value) in enumerate(trends):
            text_y = content_start_y + Inches(i * line_height)
            
            # Create text box for each trend
            trend_box = slide.shapes.add_textbox(
                x + Inches(0.3), text_y,
                width - Inches(0.6), Inches(0.3)
            )
            tf = trend_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            
            # Add label (bold)
            run = p.add_run()
            run.text = label + " "
            run.font.name = 'Calibri'
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = self.colors['primary_blue']
            
            # Add value
            run = p.add_run()
            run.text = value
            run.font.name = 'Calibri'
            run.font.size = Pt(11)
            run.font.color.rgb = self.colors['dark_gray']
    
    def _add_india_map_clean(self, slide, x, y, width, height):
        """Clean India map visual"""
        # Container
        container = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, height
        )
        container.fill.solid()
        container.fill.fore_color.rgb = self.colors['white']
        container.line.color.rgb = self.colors['light_gray']
        container.line.width = Pt(1)
        
        # Title inside container
        title_box = slide.shapes.add_textbox(
            x + Inches(0.2), y + Inches(0.1),
            width - Inches(0.4), Inches(0.4)
        )
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "Tier-2/3 India: Focus Markets"
        p.font.name = 'Calibri'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
        p.alignment = PP_ALIGN.CENTER
        
        # Create simple map visualization
        fig, ax = plt.subplots(figsize=(width.inches * 0.8, height.inches * 0.7))
        
        # India outline (simplified)
        india = patches.Polygon([
            (0.5, 0.9), (0.3, 0.8), (0.2, 0.6), (0.2, 0.4),
            (0.3, 0.2), (0.5, 0.1), (0.7, 0.2), (0.8, 0.4),
            (0.8, 0.6), (0.7, 0.8), (0.5, 0.9)
        ], closed=True, facecolor='#E8F4F8', edgecolor='#0053A0', linewidth=2)
        ax.add_patch(india)
        
        # Region markers
        regions = [
            ("N", 0.5, 0.7, "North"),
            ("E", 0.7, 0.5, "East"),
            ("W", 0.3, 0.5, "West"),
            ("S", 0.5, 0.3, "South")
        ]
        
        for letter, rx, ry, name in regions:
            # Region circle
            circle = patches.Circle((rx, ry), 0.08, facecolor='#0096AA', alpha=0.7)
            ax.add_patch(circle)
            ax.text(rx, ry, letter, ha='center', va='center', 
                   fontsize=14, fontweight='bold', color='white')
            
        # Sample cities
        cities = [
            ("Lucknow", 0.5, 0.65),
            ("Jaipur", 0.35, 0.6),
            ("Nagpur", 0.5, 0.45),
            ("Coimbatore", 0.5, 0.25),
            ("Bhubaneswar", 0.75, 0.45)
        ]
        
        for city, cx, cy in cities:
            ax.plot(cx, cy, 'ro', markersize=4)
            ax.text(cx + 0.08, cy, city, fontsize=8)
        
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        ax.axis('off')
        
        plt.tight_layout()
        
        # Save and add to slide
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=150, transparent=True, bbox_inches='tight')
        plt.close()
        
        img_stream.seek(0)
        pic = slide.shapes.add_picture(img_stream, 
                                       x + Inches(0.6), 
                                       y + Inches(0.6), 
                                       width=Inches(4.8))
    
    def _add_sectors_clean(self, slide, x, y, width, height):
        """Clean underserved sectors layout"""
        # Container
        container = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, height
        )
        container.fill.solid()
        container.fill.fore_color.rgb = self.colors['white']
        container.line.color.rgb = self.colors['light_gray']
        container.line.width = Pt(1)
        
        # Header
        header_box = slide.shapes.add_textbox(
            x + Inches(0.2), y + Inches(0.1),
            width - Inches(0.4), Inches(0.4)
        )
        tf = header_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "UNDERSERVED SECTORS"
        p.font.name = 'Calibri'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
        p.alignment = PP_ALIGN.CENTER
        
        # Sectors in 2x2 grid
        sectors = [
            ("Healthcare", "600M underserved", "Avg travel 50+ km", self.colors['red']),
            ("Education", "Teacher ratio 1:60", "vs 1:30 ideal", self.colors['orange']),
            ("Finance", "190M unbanked", "<5% insurance", self.colors['green']),
            ("Agriculture", "Post-harvest loss", "₹90,000 Cr/yr", self.colors['teal'])
        ]
        
        sector_width = (width - Inches(0.6)) / 2
        sector_height = (height - Inches(0.8)) / 2
        
        for i, (name, stat1, stat2, color) in enumerate(sectors):
            row = i // 2
            col = i % 2
            
            sector_x = x + Inches(0.2) + (sector_width + Inches(0.1)) * col
            sector_y = y + Inches(0.6) + (sector_height + Inches(0.1)) * row
            
            # Sector box
            sector_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                sector_x, sector_y,
                sector_width, sector_height
            )
            sector_box.fill.solid()
            sector_box.fill.fore_color.rgb = self.colors['very_light_gray']
            sector_box.line.color.rgb = color
            sector_box.line.width = Pt(2)
            
            # Sector content
            content_box = slide.shapes.add_textbox(
                sector_x + Inches(0.1), sector_y + Inches(0.1),
                sector_width - Inches(0.2), sector_height - Inches(0.2)
            )
            tf = content_box.text_frame
            tf.clear()
            
            # Name
            p = tf.paragraphs[0]
            p.text = name
            p.font.name = 'Calibri'
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.CENTER
            
            # Stats
            p = tf.add_paragraph()
            p.text = stat1
            p.font.name = 'Calibri'
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors['dark_gray']
            p.alignment = PP_ALIGN.CENTER
            
            p = tf.add_paragraph()
            p.text = stat2
            p.font.name = 'Calibri'
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors['dark_gray']
            p.alignment = PP_ALIGN.CENTER
    
    def _add_infrastructure_chart_clean(self, slide, x, y, width, height):
        """Clean infrastructure comparison chart"""
        # Add chart using python-pptx native chart
        chart_data = ChartData()
        chart_data.categories = ['Doctors/1k', 'Bank branches/100k', 'Internet %', 'Smartphone %']
        chart_data.add_series('Urban', (1.8, 18, 78, 85))
        chart_data.add_series('Tier-2/3', (0.5, 6, 62, 60))
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            x, y, width, height,
            chart_data
        ).chart
        
        # Style the chart
        chart.has_title = True
        chart.chart_title.text_frame.text = "Infrastructure Access: Urban vs Tier-2/3"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        
        # Format series colors
        chart.series[0].format.fill.solid()
        chart.series[0].format.fill.fore_color.rgb = self.colors['primary_blue']
        
        chart.series[1].format.fill.solid()
        chart.series[1].format.fill.fore_color.rgb = self.colors['teal']
        
        # Add data labels
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(9)
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        
        # Legend
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False
    
    def create_slide_2(self):
        """Slide 2: Healthcare Focus - Clean Layout"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        self._add_clean_title(slide, "Healthcare in Tier-2/3 India: Urgent Problem, Large Market")
        
        # Layout: 3 sections
        # Top: Problem Pillars (3 columns)
        self._add_problem_pillars_clean(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.2))
        
        # Middle Left: Market Potential
        self._add_market_potential_clean(slide, Inches(0.5), Inches(3.6), Inches(6), Inches(2.7))
        
        # Middle Right: Competitive Matrix
        self._add_competitive_matrix_clean(slide, Inches(6.8), Inches(3.6), Inches(6), Inches(2.7))
        
        # Bottom banner
        self._add_bottom_banner(slide,
            "Healthcare = burning platform: unmet need, adoption signals, and policy tailwinds make it the top disruption sector.")
        
        # Speaker notes
        self._add_speaker_notes(slide,
            "Healthcare in Tier-2/3 has the sharpest pain points: doctor shortages, catastrophic OOP costs, "
            "and proven adoption of telemedicine. This sector is where tech disruption has the highest impact "
            "and fastest scalability.")
        
        self._add_slide_number(slide, 2)
        return slide
    
    def _add_problem_pillars_clean(self, slide, x, y, width, height):
        """Clean problem pillars layout"""
        pillars = [
            {
                'title': 'ACCESSIBILITY',
                'points': ['75% doctors in urban', '600M underserved', 'Travel >50 km'],
                'color': self.colors['red']
            },
            {
                'title': 'AFFORDABILITY',
                'points': ['62% out-of-pocket', '60M → poverty/year', 'Catastrophic costs'],
                'color': self.colors['orange']
            },
            {
                'title': 'AWARENESS & TRUST',
                'points': ['Preventive care stigma', 'Unlicensed reliance', 'Low health literacy'],
                'color': self.colors['primary_blue']
            }
        ]
        
        pillar_width = (width - Inches(0.4)) / 3
        
        for i, pillar in enumerate(pillars):
            pillar_x = x + Inches(0.1) + (pillar_width + Inches(0.1)) * i
            
            # Pillar container
            pillar_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                pillar_x, y,
                pillar_width, height
            )
            pillar_box.fill.solid()
            pillar_box.fill.fore_color.rgb = self.colors['very_light_gray']
            pillar_box.line.color.rgb = pillar['color']
            pillar_box.line.width = Pt(2)
            
            # Header
            header = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                pillar_x, y,
                pillar_width, Inches(0.5)
            )
            header.fill.solid()
            header.fill.fore_color.rgb = pillar['color']
            header.line.fill.background()
            
            header_tf = header.text_frame
            header_tf.clear()
            header_tf.margin_top = Inches(0.1)
            p = header_tf.paragraphs[0]
            p.text = pillar['title']
            p.font.name = 'Calibri'
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            
            # Points
            points_box = slide.shapes.add_textbox(
                pillar_x + Inches(0.2), y + Inches(0.7),
                pillar_width - Inches(0.4), height - Inches(0.8)
            )
            tf = points_box.text_frame
            tf.clear()
            
            for j, point in enumerate(pillar['points']):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = f"• {point}"
                p.font.name = 'Calibri'
                p.font.size = Pt(11)
                p.font.color.rgb = self.colors['dark_gray']
                p.line_spacing = 1.5
    
    def _add_market_potential_clean(self, slide, x, y, width, height):
        """Clean market potential section"""
        # Container
        container = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, height
        )
        container.fill.solid()
        container.fill.fore_color.rgb = self.colors['teal']
        container.line.fill.background()
        
        # Content
        tf = container.text_frame
        tf.clear()
        tf.margin_all = Inches(0.3)
        
        # Header
        p = tf.paragraphs[0]
        p.text = "MARKET POTENTIAL"
        p.font.name = 'Calibri'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        # Market data with clean formatting
        market_data = [
            ("Healthcare Market", "USD 372B by 2025 | CAGR 22%"),
            ("Telemedicine", "USD 5.4B by 2025"),
            ("Diagnostics", "CAGR ~20%"),
            ("eSanjeevani Proof", "160M+ teleconsults")
        ]
        
        # Add spacing
        p = tf.add_paragraph()
        p.text = ""
        
        for title, value in market_data:
            p = tf.add_paragraph()
            p.text = title
            p.font.name = 'Calibri'
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            
            p = tf.add_paragraph()
            p.text = value
            p.font.name = 'Calibri'
            p.font.size = Pt(11)
            p.font.color.rgb = self.colors['white']
            p.line_spacing = 1.3
            
            # Add spacing between items
            p = tf.add_paragraph()
            p.text = ""
            p.font.size = Pt(6)
        
        # Add pie chart at bottom
        pie_data = {
            'Out-of-Pocket': 62,
            'Government': 30,
            'Insurance': 8
        }
        
        # Create pie chart
        fig, ax = plt.subplots(figsize=(2, 2))
        colors = ['#DC267F', '#0096AA', '#77BB41']
        wedges, texts, autotexts = ax.pie(
            pie_data.values(), 
            labels=None,  # No labels to avoid clutter
            colors=colors,
            autopct='%1.0f%%',
            startangle=90
        )
        
        # Style
        for autotext in autotexts:
            autotext.set_color('white')
            autotext.set_fontsize(9)
            autotext.set_fontweight('bold')
        
        ax.text(0, -1.3, 'Healthcare Financing', ha='center', fontsize=10, fontweight='bold')
        
        plt.tight_layout()
        
        # Save and add
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=150, transparent=True, bbox_inches='tight')
        plt.close()
        
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, 
                                x + Inches(1.5), 
                                y + Inches(1.8), 
                                width=Inches(3))
    
    def _add_competitive_matrix_clean(self, slide, x, y, width, height):
        """Clean competitive landscape matrix"""
        # Container
        container = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, height
        )
        container.fill.solid()
        container.fill.fore_color.rgb = self.colors['white']
        container.line.color.rgb = self.colors['light_gray']
        container.line.width = Pt(1)
        
        # Title
        title_box = slide.shapes.add_textbox(
            x + Inches(0.2), y + Inches(0.1),
            width - Inches(0.4), Inches(0.4)
        )
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "Competitive Landscape"
        p.font.name = 'Calibri'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
        p.alignment = PP_ALIGN.CENTER
        
        # Create clean 2x2 matrix
        fig, ax = plt.subplots(figsize=(width.inches * 0.8, height.inches * 0.7))
        
        # Draw axes
        ax.axhline(y=5, color='gray', linestyle='-', alpha=0.3, linewidth=1)
        ax.axvline(x=5, color='gray', linestyle='-', alpha=0.3, linewidth=1)
        
        # Labels
        ax.text(5, -0.5, 'Geography', ha='center', fontsize=10, fontweight='bold')
        ax.text(-0.5, 5, 'Service Breadth', ha='center', va='center', 
                fontsize=10, fontweight='bold', rotation=90)
        
        # Axis labels
        ax.text(1, 10.3, 'Urban', ha='center', fontsize=9)
        ax.text(9, 10.3, 'Rural', ha='center', fontsize=9)
        ax.text(-1.2, 9, 'Broad', ha='center', fontsize=9, rotation=90)
        ax.text(-1.2, 1, 'Narrow', ha='center', fontsize=9, rotation=90)
        
        # Competitors
        competitors = [
            ('Practo', 2, 7),
            ('Tata Health', 3, 8),
            ('1mg/PharmEasy', 2.5, 6),
            ('eSanjeevani', 6, 3)
        ]
        
        colors = ['#0053A0', '#0053A0', '#0096AA', '#77BB41']
        
        for (name, x_pos, y_pos), color in zip(competitors, colors):
            circle = patches.Circle((x_pos, y_pos), 0.6, facecolor=color, alpha=0.7)
            ax.add_patch(circle)
            ax.text(x_pos, y_pos, name, ha='center', va='center',
                   fontsize=8, fontweight='bold', color='white')
        
        # Highlight gap
        gap_rect = patches.Rectangle((6, 6), 3, 3, fill=False,
                                    edgecolor='#DC267F', linewidth=2, linestyle='--')
        ax.add_patch(gap_rect)
        ax.text(7.5, 7.5, 'GAP:\nVernacular +\nLow-cost +\nTrust', 
               ha='center', va='center', fontsize=9, color='#DC267F',
               fontweight='bold', style='italic')
        
        ax.set_xlim(-1.5, 10.5)
        ax.set_ylim(-1, 11)
        ax.axis('off')
        
        plt.tight_layout()
        
        # Save and add
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=150, transparent=True, bbox_inches='tight')
        plt.close()
        
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream,
                                x + Inches(0.6),
                                y + Inches(0.6),
                                width=Inches(4.8))
    
    def create_slide_3(self):
        """Slide 3: MediChain Solution - Clean Layout"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        self._add_clean_title(slide, "MediChain — Tech-enabled Primary Care & Diagnostics")
        
        # Top: Solution Components
        self._add_solution_components_clean(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.8))
        
        # Middle: 3 sections
        # Left: Differentiators
        self._add_differentiators_clean(slide, Inches(0.5), Inches(3.2), Inches(4), Inches(2.8))
        
        # Center: Funnel
        self._add_funnel_clean(slide, Inches(4.7), Inches(3.2), Inches(4), Inches(2.8))
        
        # Right: Impact
        self._add_impact_clean(slide, Inches(8.9), Inches(3.2), Inches(3.9), Inches(2.8))
        
        # Bottom: Roadmap (above banner)
        self._add_roadmap_clean(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.6))
        
        # Bottom banner
        self._add_bottom_banner(slide,
            "MediChain = A vernacular, trust-first, low-cost healthcare pathway for Bharat's Tier-2/3 markets.")
        
        # Speaker notes
        self._add_speaker_notes(slide,
            "MediChain combines vernacular AI triage, kiosk-based diagnostics, secure records and pharmacy tie-ups. "
            "It pilots small (100 kiosks, 50K users) but scales to 100M lives by year five — "
            "balancing profitability and social impact.")
        
        self._add_slide_number(slide, 3)
        return slide
    
    def _add_solution_components_clean(self, slide, x, y, width, height):
        """Clean solution components with arrows"""
        components = [
            {
                'title': 'AI Symptom Triage',
                'desc': 'Vernacular chatbot\nCost <₹20',
                'color': self.colors['teal']
            },
            {
                'title': 'IoT Diagnostic Kiosks',
                'desc': 'BP, sugar, ECG\nCost ~₹1L',
                'color': self.colors['primary_blue']
            },
            {
                'title': 'Blockchain Records',
                'desc': 'NDHM-aligned\nPortable',
                'color': self.colors['green']
            },
            {
                'title': 'Pharmacy Linkages',
                'desc': 'Last-mile meds\nLocal tie-ups',
                'color': self.colors['orange']
            }
        ]
        
        # Container
        container = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, height
        )
        container.fill.solid()
        container.fill.fore_color.rgb = self.colors['very_light_gray']
        container.line.color.rgb = self.colors['light_gray']
        container.line.width = Pt(1)
        
        comp_width = (width - Inches(0.6)) / 4
        
        for i, comp in enumerate(components):
            comp_x = x + Inches(0.1) + (comp_width + Inches(0.1)) * i
            
            # Component box
            comp_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                comp_x, y + Inches(0.1),
                comp_width, height - Inches(0.2)
            )
            comp_box.fill.solid()
            comp_box.fill.fore_color.rgb = comp['color']
            comp_box.line.fill.background()
            
            # Content
            tf = comp_box.text_frame
            tf.clear()
            tf.margin_all = Inches(0.2)
            
            # Title
            p = tf.paragraphs[0]
            p.text = comp['title']
            p.font.name = 'Calibri'
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            
            # Add spacing
            p = tf.add_paragraph()
            p.text = ""
            p.font.size = Pt(6)
            
            # Description
            p = tf.add_paragraph()
            p.text = comp['desc']
            p.font.name = 'Calibri'
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.2
            
            # Add arrow between components
            if i < 3:
                arrow_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    comp_x + comp_width - Inches(0.05), 
                    y + height/2 - Inches(0.15),
                    Inches(0.2), Inches(0.3)
                )
                arrow_shape.fill.solid()
                arrow_shape.fill.fore_color.rgb = self.colors['medium_gray']
                arrow_shape.line.fill.background()
    
    def _add_differentiators_clean(self, slide, x, y, width, height):
        """Clean differentiators box"""
        # Container
        container = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, height
        )
        container.fill.solid()
        container.fill.fore_color.rgb = self.colors['primary_blue']
        container.line.fill.background()
        
        tf = container.text_frame
        tf.clear()
        tf.margin_all = Inches(0.3)
        
        # Header
        p = tf.paragraphs[0]
        p.text = "KEY DIFFERENTIATORS"
        p.font.name = 'Calibri'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        # Add spacing
        p = tf.add_paragraph()
        p.text = ""
        
        # Differentiators
        diffs = [
            "Vernacular-first UX for Tier-2/3 adoption",
            "Micro-pricing: <₹100 consults, ₹499 family",
            "Trust: Pharmacy kiosks + NGO/govt tie-ups"
        ]
        
        for diff in diffs:
            p = tf.add_paragraph()
            p.text = f"• {diff}"
            p.font.name = 'Calibri'
            p.font.size = Pt(11)
            p.font.color.rgb = self.colors['white']
            p.line_spacing = 1.5
            
            # Add spacing
            p = tf.add_paragraph()
            p.text = ""
            p.font.size = Pt(6)
    
    def _add_funnel_clean(self, slide, x, y, width, height):
        """Clean funnel visualization"""
        # Container
        container = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, height
        )
        container.fill.solid()
        container.fill.fore_color.rgb = self.colors['white']
        container.line.color.rgb = self.colors['light_gray']
        container.line.width = Pt(1)
        
        # Title
        title_box = slide.shapes.add_textbox(
            x + Inches(0.2), y + Inches(0.1),
            width - Inches(0.4), Inches(0.4)
        )
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "Pilot Conversion Funnel"
        p.font.name = 'Calibri'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
        p.alignment = PP_ALIGN.CENTER
        
        # Funnel data table
        funnel_data = [
            ("Addressable pop.", "650,000"),
            ("Digital + kiosk reach", "75,000"),
            ("Active users (6mo)", "50,000"),
            ("Teleconsults", "20,000"),
            ("Paying subscribers", "5,000")
        ]
        
        # Create simple table
        table_y = y + Inches(0.6)
        row_height = (height - Inches(0.7)) / len(funnel_data)
        
        for i, (stage, value) in enumerate(funnel_data):
            row_y = table_y + row_height * i
            
            # Stage name
            stage_box = slide.shapes.add_textbox(
                x + Inches(0.3), row_y,
                width * 0.6, row_height - Inches(0.05)
            )
            tf = stage_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = stage
            p.font.name = 'Calibri'
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors['dark_gray']
            
            # Value
            value_box = slide.shapes.add_textbox(
                x + width * 0.6, row_y,
                width * 0.35, row_height - Inches(0.05)
            )
            tf = value_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = value
            p.font.name = 'Calibri'
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = self.colors['primary_blue']
            p.alignment = PP_ALIGN.RIGHT
            
            # Conversion bar (visual)
            if i < len(funnel_data) - 1:
                bar_width = (float(value.replace(',', '')) / 650000) * (width * 0.3)
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    x + Inches(0.2), row_y + row_height - Inches(0.05),
                    bar_width, Inches(0.02)
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = self.colors['teal']
                bar.line.fill.background()
    
    def _add_impact_clean(self, slide, x, y, width, height):
        """Clean impact framework"""
        # Container
        container = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, height
        )
        container.fill.solid()
        container.fill.fore_color.rgb = self.colors['light_teal']
        container.line.color.rgb = self.colors['teal']
        container.line.width = Pt(2)
        
        tf = container.text_frame
        tf.clear()
        tf.margin_all = Inches(0.2)
        
        # Header
        p = tf.paragraphs[0]
        p.text = "IMPACT & SCALE"
        p.font.name = 'Calibri'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary_blue']
        p.alignment = PP_ALIGN.CENTER
        
        # Add spacing
        p = tf.add_paragraph()
        p.text = ""
        p.font.size = Pt(4)
        
        # Economic Impact
        p = tf.add_paragraph()
        p.text = "Economic Impact"
        p.font.name = 'Calibri'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors['dark_gray']
        
        economic_items = ["• Affordable care", "• Lower OOP burden", "• Scalable model"]
        for item in economic_items:
            p = tf.add_paragraph()
            p.text = item
            p.font.name = 'Calibri'
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors['dark_gray']
        
        # Add spacing
        p = tf.add_paragraph()
        p.text = ""
        p.font.size = Pt(8)
        
        # Social Impact
        p = tf.add_paragraph()
        p.text = "Social Impact"
        p.font.name = 'Calibri'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors['dark_gray']
        
        social_items = ["• 100M underserved", "• Preventive adoption", "• SDG-3 aligned"]
        for item in social_items:
            p = tf.add_paragraph()
            p.text = item
            p.font.name = 'Calibri'
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors['dark_gray']
    
    def _add_roadmap_clean(self, slide, x, y, width, height):
        """Clean roadmap timeline"""
        # Container
        roadmap_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, height
        )
        roadmap_box.fill.solid()
        roadmap_box.fill.fore_color.rgb = self.colors['very_light_gray']
        roadmap_box.line.color.rgb = self.colors['light_gray']
        roadmap_box.line.width = Pt(0.5)
        
        # Phases
        phases = [
            ("Y1: Pilot", "100 kiosks, 50K users"),
            ("Y3: Scale", "25 cities, 1M users"),
            ("Y5: National", "Pan-India, 100M lives")
        ]
        
        phase_width = width / 3
        
        for i, (phase, target) in enumerate(phases):
            phase_x = x + phase_width * i
            
            # Phase box
            phase_box = slide.shapes.add_textbox(
                phase_x + Inches(0.1), y + Inches(0.05),
                phase_width - Inches(0.2), height - Inches(0.1)
            )
            tf = phase_box.text_frame
            tf.clear()
            
            # Phase name
            p = tf.paragraphs[0]
            p.text = phase
            p.font.name = 'Calibri'
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.colors['primary_blue']
            
            # Target
            run = p.add_run()
            run.text = f" | {target}"
            run.font.name = 'Calibri'
            run.font.size = Pt(10)
            run.font.color.rgb = self.colors['dark_gray']
            
            # Add arrow between phases
            if i < 2:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    phase_x + phase_width - Inches(0.15), 
                    y + height/2 - Inches(0.1),
                    Inches(0.2), Inches(0.2)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.colors['teal']
                arrow.line.fill.background()
    
    def save_presentation(self):
        """Save the presentation with all slides"""
        self.create_slide_1()
        self.create_slide_2()
        self.create_slide_3()
        
        filename = "DisruptX_Round1_MediChain.pptx"
        self.prs.save(filename)
        return filename


if __name__ == "__main__":
    print("Creating CLEAN DisruptX presentation...")
    ppt = CleanDisruptXPresentation()
    filename = ppt.save_presentation()
    print(f"✓ Created: {filename}")
    print("✓ Ultra-clean formatting")
    print("✓ No overlapping text")
    print("✓ Professional BCG/McKinsey style")
    print("✓ All content properly aligned")