#!/usr/bin/env python3
"""
DisruptX Round 1 - MediChain Presentation VERSION 2
Ultra-dense BCG/McKinsey-style 3-slide deck with maximum content density
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import FancyBboxPatch, Circle, Rectangle, FancyArrowPatch
import matplotlib.patches as mpatches
import numpy as np
import io

class DisruptXVersion2:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9
        self.prs.slide_height = Inches(7.5)
        
        # BCG color palette
        self.colors = {
            'bcg_blue': RGBColor(0, 83, 159),
            'bcg_green': RGBColor(0, 155, 119),
            'bcg_orange': RGBColor(242, 148, 0),
            'bcg_purple': RGBColor(138, 53, 168),
            'bcg_pink': RGBColor(237, 0, 140),
            'bcg_light_blue': RGBColor(0, 181, 226),
            'dark_gray': RGBColor(51, 51, 51),
            'medium_gray': RGBColor(128, 128, 128),
            'light_gray': RGBColor(217, 217, 217),
            'very_light_gray': RGBColor(242, 242, 242),
            'white': RGBColor(255, 255, 255),
        }
    
    def _add_dense_title(self, slide, text, subtitle=None):
        """Add title with subtitle - minimal spacing"""
        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.15), 
            Inches(12.7), Inches(0.6)
        )
        tf = title_box.text_frame
        tf.clear()
        tf.margin_all = Inches(0.05)
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Calibri'
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.colors['bcg_blue']
        p.alignment = PP_ALIGN.LEFT
        
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.name = 'Calibri'
            p2.font.size = Pt(14)
            p2.font.color.rgb = self.colors['medium_gray']
            p2.alignment = PP_ALIGN.LEFT
            p2.space_before = Pt(2)
            
    def _add_bottom_insight(self, slide, text):
        """Add bottom insight banner"""
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(7.15),
            Inches(13.333), Inches(0.35)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = self.colors['bcg_blue']
        banner.line.fill.background()
        
        tf = banner.text_frame
        tf.clear()
        tf.margin_all = Inches(0.05)
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Calibri'
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
    def create_slide1_market_opportunity(self):
        """Slide 1: The 600M Opportunity - Ultra Dense"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_dense_title(slide, "Tier 2/3 India: The ₹1.2 Trillion Healthcare Opportunity", 
                             "600M underserved population with growing digital adoption")
        
        # Create 6-panel grid layout
        # Top row: Market Size, Healthcare Gap, Digital Penetration
        # Bottom row: Income Analysis, Infrastructure Challenge, Growth Drivers
        
        # Panel 1: Market Size Visualization
        panel1 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.2), Inches(0.9),
            Inches(4.2), Inches(3)
        )
        panel1.fill.solid()
        panel1.fill.fore_color.rgb = self.colors['very_light_gray']
        panel1.line.color.rgb = self.colors['light_gray']
        
        # Market size chart with multiple data points
        fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(4, 2.8))
        
        # Population breakdown
        sizes = [300, 150, 150]
        labels = ['Critical Need\n300M', 'Basic Access\n150M', 'Moderate Access\n150M']
        colors = ['#ED008C', '#F29400', '#009B77']
        explode = (0.1, 0, 0)
        
        ax1.pie(sizes, explode=explode, labels=labels, colors=colors, autopct='%1.0f%%',
                shadow=True, startangle=90)
        ax1.set_title('Healthcare Access Segments', fontsize=10, fontweight='bold')
        
        # Growth projection
        years = [2023, 2024, 2025, 2026, 2027]
        market_size = [120, 180, 280, 420, 640]
        
        ax2.plot(years, market_size, marker='o', linewidth=2, markersize=8, color='#005395')
        ax2.fill_between(years, market_size, alpha=0.3, color='#00B5E2')
        ax2.set_title('Market Size (₹ Billion)', fontsize=10, fontweight='bold')
        ax2.set_xlabel('Year', fontsize=8)
        ax2.grid(True, alpha=0.3)
        
        for i, (year, size) in enumerate(zip(years, market_size)):
            ax2.text(year, size + 20, f'₹{size}B', ha='center', fontsize=8)
        
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=300, transparent=True)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(0.3), Inches(1.1), height=Inches(2.7))
        plt.close()
        
        # Panel 2: Healthcare Infrastructure Gap
        panel2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4.6), Inches(0.9),
            Inches(4.2), Inches(3)
        )
        panel2.fill.solid()
        panel2.fill.fore_color.rgb = self.colors['very_light_gray']
        panel2.line.color.rgb = self.colors['light_gray']
        
        # Infrastructure comparison chart
        fig, ax = plt.subplots(figsize=(4, 2.8))
        
        categories = ['Doctors\nper 10K', 'Hospital\nBeds/1K', 'Diagnostic\nCenters', 'Pharmacies\nper 10K']
        urban = [14.5, 5.8, 12.3, 8.9]
        rural = [0.9, 0.3, 0.8, 1.2]
        
        x = np.arange(len(categories))
        width = 0.35
        
        bars1 = ax.bar(x - width/2, urban, width, label='Urban', color='#005395', alpha=0.8)
        bars2 = ax.bar(x + width/2, rural, width, label='Rural', color='#F29400', alpha=0.8)
        
        # Add value labels
        for bars in [bars1, bars2]:
            for bar in bars:
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2., height,
                       f'{height:.1f}', ha='center', va='bottom', fontsize=8)
        
        ax.set_ylabel('Density', fontsize=10)
        ax.set_title('Healthcare Infrastructure: Urban vs Rural Gap', fontsize=11, fontweight='bold', pad=10)
        ax.set_xticks(x)
        ax.set_xticklabels(categories, fontsize=8)
        ax.legend(loc='upper right', fontsize=8)
        ax.grid(True, axis='y', alpha=0.3)
        
        # Add gap percentages
        for i, (u, r) in enumerate(zip(urban, rural)):
            gap_pct = ((u - r) / u) * 100
            ax.text(i, max(u, r) + 0.5, f'-{gap_pct:.0f}%', ha='center', fontsize=8, 
                   color='red', fontweight='bold')
        
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=300, transparent=True)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(4.7), Inches(1.1), height=Inches(2.7))
        plt.close()
        
        # Panel 3: Digital Adoption Metrics
        panel3 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(9), Inches(0.9),
            Inches(4.1), Inches(3)
        )
        panel3.fill.solid()
        panel3.fill.fore_color.rgb = self.colors['very_light_gray']
        panel3.line.color.rgb = self.colors['light_gray']
        
        # Digital metrics visualization
        fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(3.9, 2.7))
        
        # Smartphone penetration trend
        years = ['2020', '2021', '2022', '2023', '2024E']
        penetration = [32, 41, 52, 67, 78]
        
        ax1.bar(years, penetration, color='#009B77', alpha=0.8)
        ax1.set_title('Smartphone Penetration (%)', fontsize=10, fontweight='bold')
        ax1.set_ylim(0, 100)
        
        for i, v in enumerate(penetration):
            ax1.text(i, v + 1, f'{v}%', ha='center', fontsize=8, fontweight='bold')
        
        # Internet usage by purpose
        purposes = ['Social\nMedia', 'E-comm', 'Education', 'Healthcare']
        usage = [85, 62, 48, 12]
        
        bars = ax2.barh(purposes, usage, color=['#8A35A8', '#F29400', '#00B5E2', '#ED008C'])
        ax2.set_title('Internet Usage by Purpose (%)', fontsize=10, fontweight='bold')
        ax2.set_xlim(0, 100)
        
        for i, (bar, val) in enumerate(zip(bars, usage)):
            ax2.text(val + 1, bar.get_y() + bar.get_height()/2, 
                    f'{val}%', va='center', fontsize=8)
        
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=300, transparent=True)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(9.1), Inches(1.1), height=Inches(2.7))
        plt.close()
        
        # Panel 4: Income Distribution Analysis
        panel4 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.2), Inches(4),
            Inches(4.2), Inches(2.9)
        )
        panel4.fill.solid()
        panel4.fill.fore_color.rgb = self.colors['very_light_gray']
        panel4.line.color.rgb = self.colors['light_gray']
        
        # Income analysis chart
        fig, ax = plt.subplots(figsize=(4, 2.7))
        
        income_brackets = ['<5K', '5-10K', '10-20K', '20-30K', '>30K']
        population_pct = [28, 35, 22, 10, 5]
        healthcare_spend = [2, 4, 7, 12, 18]  # % of income
        
        x = np.arange(len(income_brackets))
        width = 0.4
        
        ax2 = ax.twinx()
        
        bars1 = ax.bar(x - width/2, population_pct, width, label='Population %', 
                       color='#005395', alpha=0.7)
        bars2 = ax2.bar(x + width/2, healthcare_spend, width, label='Healthcare Spend %', 
                        color='#ED008C', alpha=0.7)
        
        ax.set_xlabel('Monthly Income (₹)', fontsize=9)
        ax.set_ylabel('Population %', fontsize=9, color='#005395')
        ax2.set_ylabel('Healthcare Spend %', fontsize=9, color='#ED008C')
        ax.set_title('Income Distribution & Healthcare Affordability', fontsize=11, fontweight='bold')
        
        ax.set_xticks(x)
        ax.set_xticklabels(income_brackets, fontsize=8)
        ax.tick_params(axis='y', labelcolor='#005395')
        ax2.tick_params(axis='y', labelcolor='#ED008C')
        
        # Combined legend
        lines1, labels1 = ax.get_legend_handles_labels()
        lines2, labels2 = ax2.get_legend_handles_labels()
        ax.legend(lines1 + lines2, labels1 + labels2, loc='upper right', fontsize=8)
        
        # Add affordability index
        for i, (pop, spend) in enumerate(zip(population_pct, healthcare_spend)):
            affordability = 100 - (spend * 5)  # Simple affordability metric
            ax.text(i, -5, f'{affordability:.0f}', ha='center', fontsize=7, 
                   bbox=dict(boxstyle='round,pad=0.3', facecolor='yellow', alpha=0.5))
        
        ax.text(len(income_brackets)/2, -8, 'Affordability Index ↑', 
               ha='center', fontsize=8, style='italic')
        
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=300, transparent=True)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(0.3), Inches(4.1), height=Inches(2.8))
        plt.close()
        
        # Panel 5: Distance & Access Challenge
        panel5 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4.6), Inches(4),
            Inches(4.2), Inches(2.9)
        )
        panel5.fill.solid()
        panel5.fill.fore_color.rgb = self.colors['very_light_gray']
        panel5.line.color.rgb = self.colors['light_gray']
        
        # Access challenge visualization
        fig, ax = plt.subplots(figsize=(4, 2.7))
        
        # Create a heatmap-style visualization
        states = ['UP', 'Bihar', 'MP', 'Rajasthan', 'Jharkhand']
        metrics = ['Avg Distance\nto Hospital', 'Travel Time\n(hours)', 'Transport\nAvailable', 'Road\nQuality']
        
        data = np.array([
            [47, 3.2, 35, 28],  # UP
            [52, 3.8, 28, 22],  # Bihar
            [43, 3.0, 40, 35],  # MP
            [58, 4.1, 32, 30],  # Rajasthan
            [49, 3.5, 30, 25]   # Jharkhand
        ])
        
        # Normalize data for color mapping
        normalized_data = (data - data.min(axis=0)) / (data.max(axis=0) - data.min(axis=0))
        
        im = ax.imshow(normalized_data.T, cmap='RdYlGn_r', aspect='auto')
        
        ax.set_xticks(np.arange(len(states)))
        ax.set_yticks(np.arange(len(metrics)))
        ax.set_xticklabels(states, fontsize=9)
        ax.set_yticklabels(metrics, fontsize=8)
        ax.set_title('Healthcare Access Challenges by State', fontsize=11, fontweight='bold', pad=10)
        
        # Add text annotations
        for i in range(len(states)):
            for j in range(len(metrics)):
                if j == 0:
                    text = f'{data[i, j]:.0f}km'
                elif j == 1:
                    text = f'{data[i, j]:.1f}h'
                else:
                    text = f'{data[i, j]:.0f}%'
                ax.text(i, j, text, ha='center', va='center', fontsize=8,
                       color='white' if normalized_data[i, j] > 0.5 else 'black')
        
        # Add colorbar
        cbar = plt.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
        cbar.ax.set_ylabel('Challenge Level', fontsize=8)
        cbar.ax.tick_params(labelsize=7)
        
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=300, transparent=True)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(4.7), Inches(4.1), height=Inches(2.8))
        plt.close()
        
        # Panel 6: Growth Drivers & Opportunities
        panel6 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(9), Inches(4),
            Inches(4.1), Inches(2.9)
        )
        panel6.fill.solid()
        panel6.fill.fore_color.rgb = self.colors['very_light_gray']
        panel6.line.color.rgb = self.colors['light_gray']
        
        # Growth drivers radar chart
        fig, ax = plt.subplots(figsize=(3.9, 2.7), subplot_kw=dict(projection='polar'))
        
        categories = ['Govt\nSupport', 'Digital\nAdoption', 'Income\nGrowth', 
                     'Health\nAwareness', 'Insurance\nPenetration', 'Infrastructure']
        values = [85, 78, 65, 72, 45, 58]
        
        angles = np.linspace(0, 2 * np.pi, len(categories), endpoint=False).tolist()
        values += values[:1]
        angles += angles[:1]
        
        ax.plot(angles, values, 'o-', linewidth=2, color='#009B77')
        ax.fill(angles, values, alpha=0.25, color='#009B77')
        
        ax.set_theta_offset(np.pi / 2)
        ax.set_theta_direction(-1)
        ax.set_xticks(angles[:-1])
        ax.set_xticklabels(categories, fontsize=8)
        ax.set_ylim(0, 100)
        ax.set_title('Growth Enablers Score (0-100)', fontsize=11, fontweight='bold', pad=20)
        ax.grid(True)
        
        # Add value labels
        for angle, value, cat in zip(angles[:-1], values[:-1], categories):
            ax.text(angle, value + 5, str(value), ha='center', fontsize=8, fontweight='bold')
        
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=300, transparent=True)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(9.1), Inches(4.1), height=Inches(2.8))
        plt.close()
        
        # Add key stats callouts
        stats = [
            ("₹1.2T", "Market Size\nby 2027", self.colors['bcg_blue']),
            ("93%", "Unmet Healthcare\nDemand", self.colors['bcg_pink']),
            ("3.5hr", "Avg Travel Time\nto Hospital", self.colors['bcg_orange']),
            ("67%", "Smartphone\nPenetration", self.colors['bcg_green'])
        ]
        
        for i, (value, label, color) in enumerate(stats):
            x_pos = 0.2 + (i * 3.3)
            
            stat_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(0.75),
                Inches(0.8), Inches(0.6)
            )
            stat_box.fill.solid()
            stat_box.fill.fore_color.rgb = color
            stat_box.line.fill.background()
            
            tf = stat_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = value
            p.font.name = 'Calibri'
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            
            label_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.85), Inches(0.75),
                Inches(1.2), Inches(0.6)
            )
            tf = label_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = label
            p.font.name = 'Calibri'
            p.font.size = Pt(9)
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.LEFT
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        self._add_bottom_insight(slide, "600M lives | ₹1.2T opportunity | 93% unmet demand = India's largest untapped market")
        
    def create_slide2_solution(self):
        """Slide 2: MediChain Solution - Ultra Dense"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_dense_title(slide, "MediChain: End-to-End Digital Health Platform", 
                             "AI-first approach | Vernacular-native | ASHA-integrated | Asset-light model")
        
        # Create comprehensive solution architecture
        fig, ax = plt.subplots(figsize=(12.8, 5.5))
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 10)
        ax.axis('off')
        
        # Central MediChain Core
        core_rect = FancyBboxPatch((4, 4), 2, 2, boxstyle="round,pad=0.1",
                                  facecolor='#005395', edgecolor='#005395', linewidth=3)
        ax.add_patch(core_rect)
        ax.text(5, 5, 'MediChain\nIntelligent Core', fontsize=14, fontweight='bold',
               ha='center', va='center', color='white')
        
        # Layer 1: User Touchpoints (Top)
        touchpoints = [
            (1, 8.5, 'WhatsApp\nBot', '#25D366'),
            (2.5, 8.5, 'Voice\nAssistant', '#F29400'),
            (4, 8.5, 'Mobile\nApp', '#00B5E2'),
            (5.5, 8.5, 'ASHA\nInterface', '#ED008C'),
            (7, 8.5, 'Web\nPortal', '#8A35A8'),
            (8.5, 8.5, 'IVRS\nSystem', '#009B77')
        ]
        
        for x, y, label, color in touchpoints:
            rect = FancyBboxPatch((x-0.4, y-0.4), 0.8, 0.8, boxstyle="round,pad=0.05",
                                 facecolor=color, edgecolor=color, alpha=0.9)
            ax.add_patch(rect)
            ax.text(x, y, label, fontsize=9, ha='center', va='center', color='white', fontweight='bold')
            # Connect to core
            ax.plot([x, 5], [y-0.4, 6], 'k--', alpha=0.3, linewidth=1)
        
        # Layer 2: Core Services (Surrounding)
        services = [
            (1.5, 5, 'AI Triage\n& Diagnosis', '#00B5E2', 'left'),
            (1.5, 3.5, 'Appointment\nScheduling', '#009B77', 'left'),
            (8.5, 5, 'Medicine\nDelivery', '#F29400', 'right'),
            (8.5, 3.5, 'Lab Test\nBooking', '#ED008C', 'right'),
            (5, 7.5, 'Doctor\nConsultation', '#8A35A8', 'center'),
            (5, 2, 'Payment\n& Insurance', '#005395', 'center')
        ]
        
        for x, y, label, color, align in services:
            rect = FancyBboxPatch((x-0.6, y-0.4), 1.2, 0.8, boxstyle="round,pad=0.05",
                                 facecolor=color, edgecolor=color, alpha=0.8)
            ax.add_patch(rect)
            ax.text(x, y, label, fontsize=10, ha='center', va='center', color='white', fontweight='bold')
            # Connect to core
            if align == 'left':
                ax.plot([x+0.6, 4], [y, 5], 'k-', alpha=0.3, linewidth=2)
            elif align == 'right':
                ax.plot([x-0.6, 6], [y, 5], 'k-', alpha=0.3, linewidth=2)
            else:
                ax.plot([x, 5], [y+0.4 if y > 5 else y-0.4, 4 if y < 5 else 6], 'k-', alpha=0.3, linewidth=2)
        
        # Layer 3: Backend Infrastructure (Bottom)
        infrastructure = [
            (1, 0.5, 'Cloud\nInfra', '#808080'),
            (2.5, 0.5, 'AI/ML\nEngine', '#606060'),
            (4, 0.5, 'Data\nLake', '#404040'),
            (5.5, 0.5, 'Security\n& Privacy', '#202020'),
            (7, 0.5, 'Analytics\nEngine', '#303030'),
            (8.5, 0.5, 'Integration\nAPIs', '#505050')
        ]
        
        for x, y, label, color in infrastructure:
            rect = FancyBboxPatch((x-0.4, y-0.3), 0.8, 0.6, boxstyle="round,pad=0.05",
                                 facecolor=color, edgecolor=color)
            ax.add_patch(rect)
            ax.text(x, y, label, fontsize=8, ha='center', va='center', color='white')
            # Connect to core
            ax.plot([x, 5], [y+0.3, 4], 'k:', alpha=0.3, linewidth=1)
        
        # Add data flow indicators
        # User journey flow
        arrow1 = FancyArrowPatch((4, 8.1), (4, 6.4), connectionstyle="arc3,rad=0.3",
                                arrowstyle='->', mutation_scale=20, linewidth=2, color='green')
        ax.add_patch(arrow1)
        ax.text(3.5, 7.2, 'User\nJourney', fontsize=8, color='green', fontweight='bold')
        
        # Data flow
        arrow2 = FancyArrowPatch((6, 4.5), (7.5, 1), connectionstyle="arc3,rad=-0.3",
                                arrowstyle='->', mutation_scale=20, linewidth=2, color='blue')
        ax.add_patch(arrow2)
        ax.text(7, 2.5, 'Data\nFlow', fontsize=8, color='blue', fontweight='bold')
        
        # Service integration
        arrow3 = FancyArrowPatch((2, 4.5), (4, 5), connectionstyle="arc3,rad=0.2",
                                arrowstyle='->', mutation_scale=20, linewidth=2, color='orange')
        ax.add_patch(arrow3)
        ax.text(2.5, 4.8, 'Service\nIntegration', fontsize=8, color='orange', fontweight='bold')
        
        # Add key differentiators
        diff_box = FancyBboxPatch((0.2, 1.5), 2.3, 2, boxstyle="round,pad=0.1",
                                 facecolor='#F0F0F0', edgecolor='#005395', linewidth=2)
        ax.add_patch(diff_box)
        ax.text(1.35, 3.2, 'Key Differentiators', fontsize=10, fontweight='bold', ha='center')
        ax.text(0.4, 2.8, '✓ 11 Vernacular languages', fontsize=8, va='top')
        ax.text(0.4, 2.5, '✓ Offline-first architecture', fontsize=8, va='top')
        ax.text(0.4, 2.2, '✓ ASHA worker network', fontsize=8, va='top')
        ax.text(0.4, 1.9, '✓ AI-powered triage', fontsize=8, va='top')
        
        # Add metrics box
        metrics_box = FancyBboxPatch((7.5, 1.5), 2.3, 2, boxstyle="round,pad=0.1",
                                    facecolor='#F0F0F0', edgecolor='#009B77', linewidth=2)
        ax.add_patch(metrics_box)
        ax.text(8.65, 3.2, 'Platform Metrics', fontsize=10, fontweight='bold', ha='center')
        ax.text(7.7, 2.8, '• <30s response time', fontsize=8, va='top')
        ax.text(7.7, 2.5, '• 99.9% uptime SLA', fontsize=8, va='top')
        ax.text(7.7, 2.2, '• 15 min avg consultation', fontsize=8, va='top')
        ax.text(7.7, 1.9, '• 4.8/5 user rating', fontsize=8, va='top')
        
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=300, transparent=True)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(0.2), Inches(0.85), width=Inches(12.9))
        plt.close()
        
        # Bottom section: Competitive Analysis + Value Chain + Revenue Model
        
        # Competitive Positioning Matrix
        comp_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.2), Inches(6.2),
            Inches(4.2), Inches(1.8)
        )
        comp_box.fill.solid()
        comp_box.fill.fore_color.rgb = self.colors['very_light_gray']
        comp_box.line.color.rgb = self.colors['light_gray']
        
        # Create detailed competitive analysis
        comp_table = slide.shapes.add_table(5, 4, Inches(0.3), Inches(6.35), 
                                           Inches(4), Inches(1.6)).table
        
        # Headers
        headers = ['', 'MediChain', 'Practo/1mg', 'PharmEasy']
        for i, header in enumerate(headers):
            cell = comp_table.cell(0, i)
            cell.text = header
            tf = cell.text_frame
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.size = Pt(9)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            if i == 1:  # MediChain column
                cell.fill.fore_color.rgb = self.colors['bcg_green']
                tf.paragraphs[0].font.color.rgb = self.colors['white']
            else:
                cell.fill.fore_color.rgb = self.colors['light_gray']
        
        # Comparison data
        criteria = [
            ('Rural Focus', '✓✓✓', '✗', '✓'),
            ('Vernacular', '11 langs', 'English', '2 langs'),
            ('ASHA Network', '✓', '✗', '✗'),
            ('Affordability', '₹99/mo', '₹499/mo', '₹299/mo')
        ]
        
        for i, (criterion, mc, pr, pe) in enumerate(criteria):
            comp_table.cell(i+1, 0).text = criterion
            comp_table.cell(i+1, 0).text_frame.paragraphs[0].font.bold = True
            comp_table.cell(i+1, 0).text_frame.paragraphs[0].font.size = Pt(8)
            
            comp_table.cell(i+1, 1).text = mc
            comp_table.cell(i+1, 2).text = pr
            comp_table.cell(i+1, 3).text = pe
            
            for j in range(1, 4):
                comp_table.cell(i+1, j).text_frame.paragraphs[0].font.size = Pt(8)
                comp_table.cell(i+1, j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Value Chain Analysis
        value_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4.6), Inches(6.2),
            Inches(4.2), Inches(1.8)
        )
        value_box.fill.solid()
        value_box.fill.fore_color.rgb = self.colors['very_light_gray']
        value_box.line.color.rgb = self.colors['light_gray']
        
        # Value chain visualization
        fig, ax = plt.subplots(figsize=(4, 1.6))
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 2)
        ax.axis('off')
        
        # Value chain stages
        stages = [
            (1, 'User\nAcquisition', '₹50', '#005395'),
            (3, 'Service\nDelivery', '₹200', '#009B77'),
            (5, 'Retention', '₹150', '#F29400'),
            (7, 'Upsell', '₹300', '#ED008C'),
            (9, 'Network\nEffect', '₹500', '#8A35A8')
        ]
        
        prev_x = 0
        for x, stage, value, color in stages:
            # Stage box
            rect = FancyBboxPatch((x-0.7, 0.5), 1.4, 1, boxstyle="round,pad=0.05",
                                 facecolor=color, edgecolor=color, alpha=0.8)
            ax.add_patch(rect)
            ax.text(x, 1.2, stage, fontsize=9, ha='center', va='center', 
                   color='white', fontweight='bold')
            ax.text(x, 0.8, value, fontsize=8, ha='center', va='center', color='white')
            ax.text(x, 0.3, 'LTV', fontsize=7, ha='center', va='center', style='italic')
            
            # Arrow
            if prev_x > 0:
                arrow = FancyArrowPatch((prev_x+0.7, 1), (x-0.7, 1),
                                       arrowstyle='->', mutation_scale=15, linewidth=2)
                ax.add_patch(arrow)
            prev_x = x
        
        ax.text(5, 1.7, 'Customer Value Chain', fontsize=10, fontweight='bold', ha='center')
        
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=300, transparent=True)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(4.7), Inches(6.25), height=Inches(1.7))
        plt.close()
        
        # Revenue Model
        rev_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(9), Inches(6.2),
            Inches(4.1), Inches(1.8)
        )
        rev_box.fill.solid()
        rev_box.fill.fore_color.rgb = self.colors['very_light_gray']
        rev_box.line.color.rgb = self.colors['light_gray']
        
        # Revenue streams pie chart
        fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(3.9, 1.6))
        
        # Revenue mix
        revenues = [35, 25, 20, 15, 5]
        labels = ['Consultation', 'Pharmacy', 'Diagnostics', 'Insurance', 'Data']
        colors = ['#005395', '#009B77', '#F29400', '#ED008C', '#8A35A8']
        
        ax1.pie(revenues, labels=labels, colors=colors, autopct='%1.0f%%',
                textprops={'fontsize': 7})
        ax1.set_title('Revenue Mix', fontsize=9, fontweight='bold')
        
        # Unit economics
        metrics = ['CAC', 'LTV', 'Payback', 'Margin']
        values = ['₹150', '₹2,500', '3 mo', '22%']
        
        ax2.axis('off')
        for i, (metric, value) in enumerate(zip(metrics, values)):
            y_pos = 1.4 - (i * 0.35)
            ax2.text(0.1, y_pos, f'{metric}:', fontsize=8, fontweight='bold')
            ax2.text(0.6, y_pos, value, fontsize=8, color='green' if metric == 'Margin' else 'black')
        
        ax2.set_title('Unit Economics', fontsize=9, fontweight='bold')
        ax2.set_xlim(0, 1)
        ax2.set_ylim(0, 1.6)
        
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=300, transparent=True)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(9.1), Inches(6.25), height=Inches(1.7))
        plt.close()
        
        self._add_bottom_insight(slide, "Platform approach + Network effects + Multi-revenue streams = Sustainable competitive advantage")
        
    def create_slide3_impact(self):
        """Slide 3: Implementation & Impact - Ultra Dense"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_dense_title(slide, "₹5,000 Cr Revenue | 50M Lives | Path to Profitability",
                             "24-month national rollout | Break-even by Month 30 | 22% EBITDA margin")
        
        # Create 4-quadrant layout
        # Top-left: Detailed Implementation Roadmap
        roadmap_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.2), Inches(0.85),
            Inches(6.4), Inches(3)
        )
        roadmap_box.fill.solid()
        roadmap_box.fill.fore_color.rgb = self.colors['very_light_gray']
        roadmap_box.line.color.rgb = self.colors['light_gray']
        
        # Gantt chart style roadmap
        fig, ax = plt.subplots(figsize=(6.2, 2.8))
        
        # Define phases and activities
        phases = [
            ('Phase 1: MVP', 0, 6, '#005395'),
            ('Phase 2: Scale', 6, 12, '#009B77'),
            ('Phase 3: Expand', 12, 18, '#F29400'),
            ('Phase 4: National', 18, 24, '#ED008C')
        ]
        
        activities = [
            ('Tech Platform Build', 0, 4, 0.8),
            ('ASHA Network Setup', 1, 6, 0.6),
            ('Doctor Onboarding', 2, 24, 0.4),
            ('Partnership Development', 3, 18, 0.2),
            ('Marketing Campaign', 4, 24, 0),
            ('Regulatory Compliance', 0, 24, -0.2),
            ('Funding Rounds', 0, 18, -0.4)
        ]
        
        # Draw phases
        for phase, start, end, color in phases:
            rect = Rectangle((start, 1.2), end-start, 0.6, 
                           facecolor=color, alpha=0.3, edgecolor=color)
            ax.add_patch(rect)
            ax.text((start+end)/2, 1.5, phase, ha='center', fontsize=8, fontweight='bold')
        
        # Draw activities
        for activity, start, end, y_pos in activities:
            rect = Rectangle((start, y_pos-0.15), end-start, 0.3,
                           facecolor='#333333', alpha=0.7)
            ax.add_patch(rect)
            ax.text(-0.5, y_pos, activity, ha='right', va='center', fontsize=7)
        
        # Add milestones
        milestones = [
            (3, 'MVP Launch'),
            (9, '100K Users'),
            (15, '1M Users'),
            (21, '10M Users')
        ]
        
        for month, milestone in milestones:
            ax.plot([month, month], [-0.6, 2], 'r--', alpha=0.5)
            ax.text(month, 2.1, milestone, ha='center', fontsize=7, 
                   rotation=45, color='red', fontweight='bold')
        
        ax.set_xlim(-4, 25)
        ax.set_ylim(-0.8, 2.5)
        ax.set_xlabel('Months', fontsize=9)
        ax.set_title('24-Month Implementation Roadmap', fontsize=10, fontweight='bold')
        ax.grid(True, axis='x', alpha=0.3)
        
        # Remove y-axis
        ax.set_yticks([])
        
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=300, transparent=True)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(0.3), Inches(1), height=Inches(2.8))
        plt.close()
        
        # Top-right: Financial Projections
        fin_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.8), Inches(0.85),
            Inches(6.3), Inches(3)
        )
        fin_box.fill.solid()
        fin_box.fill.fore_color.rgb = self.colors['very_light_gray']
        fin_box.line.color.rgb = self.colors['light_gray']
        
        # Financial charts
        fig, ((ax1, ax2), (ax3, ax4)) = plt.subplots(2, 2, figsize=(6.1, 2.8))
        
        # Revenue growth
        years = ['Y1', 'Y2', 'Y3', 'Y4', 'Y5']
        revenue = [50, 250, 800, 2000, 5000]
        
        bars = ax1.bar(years, revenue, color='#005395', alpha=0.8)
        ax1.set_title('Revenue (₹ Cr)', fontsize=9, fontweight='bold')
        ax1.set_ylim(0, 6000)
        
        for bar, val in zip(bars, revenue):
            ax1.text(bar.get_x() + bar.get_width()/2, val + 100, 
                    f'₹{val}', ha='center', fontsize=7, fontweight='bold')
        
        # User growth
        users = [0.1, 0.5, 2.5, 10, 50]
        ax2.plot(years, users, marker='o', linewidth=2, markersize=8, color='#009B77')
        ax2.fill_between(range(len(years)), users, alpha=0.3, color='#009B77')
        ax2.set_title('Users (Million)', fontsize=9, fontweight='bold')
        ax2.set_ylim(0, 60)
        
        for i, (year, user) in enumerate(zip(years, users)):
            ax2.text(i, user + 2, f'{user}M', ha='center', fontsize=7)
        
        # Margin evolution
        margins = [-45, -20, 5, 18, 22]
        colors = ['red' if m < 0 else 'green' for m in margins]
        
        bars = ax3.bar(years, margins, color=colors, alpha=0.7)
        ax3.set_title('EBITDA Margin (%)', fontsize=9, fontweight='bold')
        ax3.axhline(y=0, color='black', linestyle='-', linewidth=0.5)
        ax3.set_ylim(-50, 30)
        
        for bar, val in zip(bars, margins):
            y_pos = val + 2 if val > 0 else val - 2
            ax3.text(bar.get_x() + bar.get_width()/2, y_pos, 
                    f'{val}%', ha='center', fontsize=7, fontweight='bold')
        
        # Funding & valuation
        funding_rounds = ['Seed', 'Series A', 'Series B', 'Series C']
        amounts = [10, 50, 150, 300]
        valuations = [50, 300, 1000, 3000]
        
        ax4_twin = ax4.twinx()
        
        bars = ax4.bar(range(len(funding_rounds)), amounts, alpha=0.7, 
                       color='#F29400', label='Funding')
        line = ax4_twin.plot(range(len(funding_rounds)), valuations, 
                            marker='o', linewidth=2, markersize=8, 
                            color='#ED008C', label='Valuation')
        
        ax4.set_xticks(range(len(funding_rounds)))
        ax4.set_xticklabels(funding_rounds, fontsize=7)
        ax4.set_ylabel('Funding (₹ Cr)', fontsize=8)
        ax4_twin.set_ylabel('Valuation (₹ Cr)', fontsize=8)
        ax4.set_title('Funding History', fontsize=9, fontweight='bold')
        
        # Add value labels
        for i, (amt, val) in enumerate(zip(amounts, valuations)):
            ax4.text(i, amt + 5, f'₹{amt}Cr', ha='center', fontsize=6)
            ax4_twin.text(i, val + 100, f'₹{val}Cr', ha='center', fontsize=6, color='#ED008C')
        
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=300, transparent=True)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(6.9), Inches(1), height=Inches(2.8))
        plt.close()
        
        # Bottom-left: Social Impact Metrics
        impact_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.2), Inches(3.95),
            Inches(6.4), Inches(3)
        )
        impact_box.fill.solid()
        impact_box.fill.fore_color.rgb = self.colors['very_light_gray']
        impact_box.line.color.rgb = self.colors['light_gray']
        
        # Impact visualization
        fig, ax = plt.subplots(figsize=(6.2, 2.8))
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 10)
        ax.axis('off')
        
        # Central impact circle
        main_circle = Circle((5, 5), 2, facecolor='#009B77', alpha=0.3, edgecolor='#009B77', linewidth=3)
        ax.add_patch(main_circle)
        ax.text(5, 5, '50M\nLives\nImpacted', fontsize=16, fontweight='bold', 
               ha='center', va='center', color='#005395')
        
        # Surrounding impact metrics
        impacts = [
            (2, 8, '100K', 'Direct\nJobs', '#005395'),
            (8, 8, '500K', 'Indirect\nJobs', '#F29400'),
            (1, 5, '3 hrs', 'Time Saved\nper Visit', '#ED008C'),
            (9, 5, '70%', 'Cost\nReduction', '#8A35A8'),
            (2, 2, '₹2,000', 'Avg Savings\nper Family', '#00B5E2'),
            (8, 2, '11', 'Languages\nSupported', '#009B77')
        ]
        
        for x, y, value, label, color in impacts:
            # Metric circle
            circle = Circle((x, y), 0.8, facecolor=color, alpha=0.8)
            ax.add_patch(circle)
            ax.text(x, y+0.1, value, fontsize=12, fontweight='bold', 
                   ha='center', va='center', color='white')
            ax.text(x, y-0.3, label, fontsize=8, ha='center', va='center', 
                   color='white', style='italic')
            
            # Connect to center
            ax.plot([x + (0.8 if x < 5 else -0.8) * np.cos(np.arctan2(5-y, 5-x)),
                    5 + 2 * np.cos(np.arctan2(y-5, x-5))],
                   [y + (0.8 if x < 5 else -0.8) * np.sin(np.arctan2(5-y, 5-x)),
                    5 + 2 * np.sin(np.arctan2(y-5, x-5))],
                   'k--', alpha=0.3, linewidth=1)
        
        # Add SDG alignment
        sdg_text = 'Aligned with SDG 3: Good Health & Well-being'
        ax.text(5, 0.5, sdg_text, fontsize=10, ha='center', fontweight='bold', 
               bbox=dict(boxstyle='round,pad=0.5', facecolor='yellow', alpha=0.5))
        
        ax.set_title('Social Impact Dashboard', fontsize=12, fontweight='bold', pad=20)
        
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=300, transparent=True)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(0.3), Inches(4.05), height=Inches(2.85))
        plt.close()
        
        # Bottom-right: Risk Mitigation & Next Steps
        risk_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.8), Inches(3.95),
            Inches(6.3), Inches(3)
        )
        risk_box.fill.solid()
        risk_box.fill.fore_color.rgb = self.colors['very_light_gray']
        risk_box.line.color.rgb = self.colors['light_gray']
        
        # Risk matrix and next steps
        fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(6.1, 2.8))
        
        # Risk matrix
        risks = [
            ('Regulatory', 3, 4, 'High'),
            ('Competition', 4, 3, 'Medium'),
            ('Technology', 2, 4, 'Medium'),
            ('Adoption', 3, 3, 'Medium'),
            ('Funding', 2, 3, 'Low'),
            ('Operations', 3, 2, 'Low')
        ]
        
        for risk, prob, impact, level in risks:
            color = '#ED008C' if level == 'High' else '#F29400' if level == 'Medium' else '#009B77'
            ax1.scatter(prob, impact, s=300, alpha=0.7, color=color)
            ax1.text(prob, impact, risk[:3], fontsize=8, ha='center', va='center', 
                    fontweight='bold', color='white')
        
        ax1.set_xlim(0, 5)
        ax1.set_ylim(0, 5)
        ax1.set_xlabel('Probability', fontsize=9)
        ax1.set_ylabel('Impact', fontsize=9)
        ax1.set_title('Risk Assessment Matrix', fontsize=10, fontweight='bold')
        ax1.grid(True, alpha=0.3)
        
        # Add quadrant labels
        ax1.text(1.5, 4.5, 'Monitor', fontsize=8, ha='center', style='italic', alpha=0.7)
        ax1.text(3.5, 4.5, 'Critical', fontsize=8, ha='center', style='italic', alpha=0.7)
        ax1.text(1.5, 1.5, 'Low', fontsize=8, ha='center', style='italic', alpha=0.7)
        ax1.text(3.5, 1.5, 'Manage', fontsize=8, ha='center', style='italic', alpha=0.7)
        
        # Next steps
        ax2.axis('off')
        ax2.set_title('Immediate Next Steps', fontsize=10, fontweight='bold')
        
        steps = [
            '1. Secure Series A funding (₹50 Cr)',
            '2. Launch MVP in 5 districts of UP',
            '3. Onboard 100 doctors & 500 ASHA workers',
            '4. Establish pharma partnerships',
            '5. Complete regulatory approvals',
            '6. Build core tech team (20 engineers)'
        ]
        
        for i, step in enumerate(steps):
            y_pos = 0.9 - (i * 0.15)
            ax2.text(0.05, y_pos, step, fontsize=9, va='top', 
                    bbox=dict(boxstyle='round,pad=0.3', 
                             facecolor='lightblue' if i < 3 else 'lightgray', 
                             alpha=0.7))
        
        # Add timeline
        ax2.text(0.5, 0.05, 'Timeline: 90 days', fontsize=10, ha='center', 
                fontweight='bold', color='red',
                bbox=dict(boxstyle='round,pad=0.5', facecolor='yellow', alpha=0.5))
        
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=300, transparent=True)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(6.9), Inches(4.05), height=Inches(2.85))
        plt.close()
        
        # Add key metrics callouts at the top
        metrics = [
            ("₹5,000 Cr", "Year 5 Revenue", self.colors['bcg_blue']),
            ("50M", "Lives Impacted", self.colors['bcg_green']),
            ("22%", "EBITDA Margin", self.colors['bcg_orange']),
            ("Month 30", "Break-even", self.colors['bcg_purple'])
        ]
        
        for i, (value, label, color) in enumerate(metrics):
            x_pos = 0.2 + (i * 3.3)
            
            metric_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(0.68),
                Inches(0.7), Inches(0.5)
            )
            metric_box.fill.solid()
            metric_box.fill.fore_color.rgb = color
            metric_box.line.fill.background()
            
            tf = metric_box.text_frame
            tf.clear()
            tf.margin_all = Inches(0.02)
            p = tf.paragraphs[0]
            p.text = value
            p.font.name = 'Calibri'
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            
            label_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.75), Inches(0.68),
                Inches(1.1), Inches(0.5)
            )
            tf = label_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = label
            p.font.name = 'Calibri'
            p.font.size = Pt(8)
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.LEFT
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        self._add_bottom_insight(slide, "MediChain: Where social impact meets financial returns - transforming healthcare for Bharat")
        
    def generate_presentation(self):
        """Generate the complete presentation"""
        self.create_slide1_market_opportunity()
        self.create_slide2_solution()
        self.create_slide3_impact()
        
        # Save to PPT Generated folder
        output_path = "/mnt/e/AI and Projects/Case Comp PPT/PPT Generated/DisruptX_MediChain_Version2.pptx"
        self.prs.save(output_path)
        print(f"Presentation saved to: {output_path}")
        return output_path

if __name__ == "__main__":
    # Create the presentation
    creator = DisruptXVersion2()
    output_file = creator.generate_presentation()
    print(f"\nDisruptX MediChain presentation Version 2 created successfully!")
    print(f"Ultra-dense BCG/McKinsey style with maximum content")
    print(f"Location: {output_file}")