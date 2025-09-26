from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import numpy as np
from io import BytesIO
import seaborn as sns

# Configure matplotlib for clean infographics
plt.style.use('seaborn-v0_8-darkgrid')
plt.rcParams['font.family'] = 'sans-serif'
plt.rcParams['font.size'] = 12
plt.rcParams['figure.dpi'] = 150
plt.rcParams['savefig.bbox'] = 'tight'
plt.rcParams['savefig.pad_inches'] = 0.1

class MediChainInfographicPresentation:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)

        # Professional color palette for infographics
        self.colors = {
            'primary': '#0096AA',    # Teal
            'secondary': '#00539F',  # Navy Blue
            'accent1': '#00B482',    # Green
            'accent2': '#FF8C00',    # Orange
            'accent3': '#ED1C24',    # Red
            'dark': '#2D2D2D',       # Dark Gray
            'light': '#F5F5F5',      # Light Gray
            'white': '#FFFFFF'       # White
        }

        # RGB colors for shapes
        self.rgb_colors = {
            'primary': RGBColor(0, 150, 170),
            'secondary': RGBColor(0, 83, 159),
            'accent1': RGBColor(0, 180, 130),
            'accent2': RGBColor(255, 140, 0),
            'accent3': RGBColor(237, 28, 36),
            'dark': RGBColor(45, 45, 45),
            'light': RGBColor(245, 245, 245),
            'white': RGBColor(255, 255, 255)
        }

    def add_title_slide_infographic(self, slide):
        """Add title with minimal text, maximum visual"""
        # Create gradient background effect
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Inches(9)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.rgb_colors['white']
        bg.line.fill.background()

        # Title section
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = "MediChain: Tier-2/3 Healthcare Disruption"
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.rgb_colors['secondary']
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def create_slide1_infographic(self):
        """Slide 1: Opportunity Landscape - Pure Infographic"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_title_slide_infographic(slide)

        # Create comprehensive infographic
        fig = plt.figure(figsize=(14, 6.5))

        # Main grid layout
        gs = fig.add_gridspec(2, 3, height_ratios=[1, 1], width_ratios=[1, 1, 1],
                             hspace=0.3, wspace=0.25)

        # 1. Population & GDP Pie Chart
        ax1 = fig.add_subplot(gs[0, 0])
        sizes = [45, 55]
        labels = ['Tier-2/3\n(45% GDP)', 'Others\n(55% GDP)']
        colors = [self.colors['primary'], self.colors['light']]
        explode = (0.1, 0)

        wedges, texts, autotexts = ax1.pie(sizes, explode=explode, labels=labels,
                                            colors=colors, autopct='%1.0f%%',
                                            shadow=True, startangle=90)
        ax1.set_title('GDP Contribution by 2025', fontsize=14, fontweight='bold', pad=10)

        # Make text more visible
        for text in texts:
            text.set_fontsize(11)
            text.set_fontweight('bold')
        for autotext in autotexts:
            autotext.set_color('white')
            autotext.set_fontsize(12)
            autotext.set_fontweight('bold')

        # 2. Digital Penetration Bar Chart
        ax2 = fig.add_subplot(gs[0, 1])
        categories = ['Smart\nphones', 'Internet', 'UPI\nUsers', 'Digital\nLiteracy']
        urban = [78, 82, 85, 75]
        tier23 = [60, 60, 65, 45]

        x = np.arange(len(categories))
        width = 0.35

        bars1 = ax2.bar(x - width/2, urban, width, label='Urban',
                       color=self.colors['secondary'], edgecolor='white', linewidth=2)
        bars2 = ax2.bar(x + width/2, tier23, width, label='Tier-2/3',
                       color=self.colors['primary'], edgecolor='white', linewidth=2)

        ax2.set_ylabel('Penetration %', fontsize=12, fontweight='bold')
        ax2.set_title('Digital Adoption Gap', fontsize=14, fontweight='bold')
        ax2.set_xticks(x)
        ax2.set_xticklabels(categories, fontsize=10)
        ax2.legend(fontsize=10)
        ax2.grid(axis='y', alpha=0.3)
        ax2.set_ylim(0, 100)

        # Add value labels
        for bars in [bars1, bars2]:
            for bar in bars:
                height = bar.get_height()
                ax2.text(bar.get_x() + bar.get_width()/2., height,
                        f'{int(height)}%', ha='center', va='bottom', fontsize=9)

        # 3. Market Size Bubble Chart
        ax3 = fig.add_subplot(gs[0, 2])

        x = [1, 2, 3, 4]
        y = [3, 2.5, 2, 1.5]
        sizes = [600, 190, 200, 150]  # Millions of people
        colors_bubble = [self.colors['accent3'], self.colors['accent2'],
                        self.colors['primary'], self.colors['accent1']]
        labels_bubble = ['Healthcare\n600M', 'Finance\n190M', 'Education\n200M', 'Agriculture\n150M']

        for i in range(len(x)):
            circle = plt.Circle((x[i], y[i]), sizes[i]/1000,
                               color=colors_bubble[i], alpha=0.6)
            ax3.add_patch(circle)
            ax3.text(x[i], y[i], labels_bubble[i],
                    ha='center', va='center', fontsize=10, fontweight='bold')

        ax3.set_xlim(0, 5)
        ax3.set_ylim(0, 4)
        ax3.set_title('Underserved Population (Millions)', fontsize=14, fontweight='bold')
        ax3.axis('off')

        # 4. Infrastructure Gap Heatmap
        ax4 = fig.add_subplot(gs[1, :2])

        data = np.array([
            [85, 78, 82, 90],  # Urban
            [25, 35, 45, 40],  # Tier-2
            [15, 20, 30, 25],  # Tier-3
            [10, 15, 20, 15]   # Rural
        ])

        categories_heat = ['Doctors', 'Hospitals', 'Diagnostic', 'Pharmacy']
        regions = ['Urban', 'Tier-2', 'Tier-3', 'Rural']

        im = ax4.imshow(data, cmap='RdYlGn', aspect='auto', vmin=0, vmax=100)

        ax4.set_xticks(np.arange(len(categories_heat)))
        ax4.set_yticks(np.arange(len(regions)))
        ax4.set_xticklabels(categories_heat, fontsize=11)
        ax4.set_yticklabels(regions, fontsize=11)
        ax4.set_title('Healthcare Infrastructure Availability (%)',
                     fontsize=14, fontweight='bold', pad=10)

        # Add text annotations
        for i in range(len(regions)):
            for j in range(len(categories_heat)):
                text = ax4.text(j, i, f'{data[i, j]}%',
                              ha="center", va="center", color="white",
                              fontweight='bold', fontsize=10)

        # 5. Growth Trajectory
        ax5 = fig.add_subplot(gs[1, 2])

        years = ['2020', '2021', '2022', '2023', '2024', '2025']
        digital_health = [100, 150, 250, 400, 600, 900]
        telemedicine = [50, 120, 200, 350, 500, 750]

        ax5.plot(years, digital_health, marker='o', linewidth=3,
                color=self.colors['primary'], label='Digital Health', markersize=8)
        ax5.plot(years, telemedicine, marker='s', linewidth=3,
                color=self.colors['accent1'], label='Telemedicine', markersize=8)

        ax5.set_title('Market Growth (Index)', fontsize=14, fontweight='bold')
        ax5.set_ylabel('Growth Index', fontsize=11, fontweight='bold')
        ax5.legend(fontsize=10)
        ax5.grid(True, alpha=0.3)

        plt.suptitle('Tier-2/3 India: Digital Ready, Infrastructure Poor',
                    fontsize=16, fontweight='bold', y=1.02)

        # Save and add to slide
        img_stream = BytesIO()
        plt.tight_layout()
        plt.savefig(img_stream, format='PNG', dpi=150, bbox_inches='tight',
                   facecolor='white', edgecolor='none')
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, Inches(1), Inches(1.8), width=Inches(14))
        plt.close()

        # Add minimal footer
        footer = slide.shapes.add_textbox(Inches(0), Inches(8.5), Inches(16), Inches(0.4))
        footer.text_frame.text = "650M population • $0.17/GB data • 12B+ UPI transactions/month"
        footer.text_frame.paragraphs[0].font.size = Pt(12)
        footer.text_frame.paragraphs[0].font.color.rgb = self.rgb_colors['dark']
        footer.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def create_slide2_healthcare_infographic(self):
        """Slide 2: Healthcare Focus - Visual Heavy"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Minimal title
        title = slide.shapes.add_textbox(Inches(0), Inches(0.2), Inches(16), Inches(0.6))
        title.text_frame.text = "Healthcare: The Burning Platform"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.rgb_colors['accent3']
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Create healthcare infographic
        fig = plt.figure(figsize=(14, 7))
        gs = fig.add_gridspec(2, 3, hspace=0.3, wspace=0.3)

        # 1. Healthcare Spending Donut Chart
        ax1 = fig.add_subplot(gs[0, 0])
        sizes = [62, 30, 8]
        labels = ['Out-of-Pocket\n62%', 'Government\n30%', 'Insurance\n8%']
        colors = [self.colors['accent3'], self.colors['primary'], self.colors['accent1']]

        wedges, texts = ax1.pie(sizes, labels=labels, colors=colors,
                                startangle=90, counterclock=False,
                                wedgeprops=dict(width=0.5, edgecolor='white'))

        ax1.set_title('Healthcare Spending Pattern', fontsize=13, fontweight='bold')
        for text in texts:
            text.set_fontsize(10)
            text.set_fontweight('bold')

        # 2. Doctor Density Comparison
        ax2 = fig.add_subplot(gs[0, 1])
        regions = ['Urban', 'Tier-2', 'Tier-3', 'Rural']
        doctors = [1.2, 0.5, 0.3, 0.2]
        colors_bar = [self.colors['secondary'], self.colors['primary'],
                     self.colors['accent1'], self.colors['accent2']]

        bars = ax2.barh(regions, doctors, color=colors_bar, edgecolor='white', linewidth=2)
        ax2.set_xlabel('Doctors per 1000 people', fontsize=11, fontweight='bold')
        ax2.set_title('Doctor Distribution Crisis', fontsize=13, fontweight='bold')
        ax2.set_xlim(0, 1.5)

        for i, (bar, value) in enumerate(zip(bars, doctors)):
            ax2.text(value + 0.02, bar.get_y() + bar.get_height()/2,
                    f'{value}', va='center', fontweight='bold', fontsize=10)

        # 3. Distance to Healthcare
        ax3 = fig.add_subplot(gs[0, 2])

        distances = [5, 25, 50, 75]
        regions_dist = ['Urban', 'Tier-2', 'Tier-3', 'Rural']
        colors_dist = [self.colors['secondary'], self.colors['primary'],
                      self.colors['accent2'], self.colors['accent3']]

        for i, (dist, region, color) in enumerate(zip(distances, regions_dist, colors_dist)):
            circle = plt.Circle((2, 4-i), dist/100, color=color, alpha=0.7)
            ax3.add_patch(circle)
            ax3.text(3.2, 4-i, f'{region}: {dist}km',
                    va='center', fontsize=11, fontweight='bold')

        ax3.set_xlim(0, 4.5)
        ax3.set_ylim(0, 5)
        ax3.set_title('Average Distance to Hospital', fontsize=13, fontweight='bold')
        ax3.axis('off')

        # 4. Market Growth Funnel
        ax4 = fig.add_subplot(gs[1, 0])

        stages = ['Total Market\n$372B', 'Digital Health\n$45B', 'Telemedicine\n$5.4B']
        values = [372, 45, 5.4]
        colors_funnel = [self.colors['secondary'], self.colors['primary'], self.colors['accent1']]

        for i, (stage, value, color) in enumerate(zip(stages, values, colors_funnel)):
            width = value / values[0]
            left = (1 - width) / 2
            rect = patches.Rectangle((left, i), width, 0.8,
                                    facecolor=color, edgecolor='white', linewidth=2)
            ax4.add_patch(rect)
            ax4.text(0.5, i + 0.4, stage, ha='center', va='center',
                    fontsize=11, fontweight='bold', color='white')

        ax4.set_xlim(0, 1)
        ax4.set_ylim(-0.5, 3)
        ax4.set_title('Healthcare Market by 2025', fontsize=13, fontweight='bold')
        ax4.axis('off')

        # 5. Key Metrics Dashboard
        ax5 = fig.add_subplot(gs[1, 1:])
        ax5.axis('off')

        # Create metric cards
        metrics = [
            ('600M', 'Underserved Population', self.colors['accent3']),
            ('60M', 'Fall into Poverty Yearly', self.colors['accent2']),
            ('22%', 'Healthcare CAGR', self.colors['primary']),
            ('160M+', 'Teleconsultations Done', self.colors['accent1'])
        ]

        for i, (value, label, color) in enumerate(metrics):
            x = (i % 2) * 0.5
            y = 0.7 - (i // 2) * 0.4

            # Metric box
            rect = patches.FancyBboxPatch((x, y), 0.45, 0.3,
                                         boxstyle="round,pad=0.02",
                                         facecolor=color, alpha=0.2,
                                         edgecolor=color, linewidth=2)
            ax5.add_patch(rect)

            # Value
            ax5.text(x + 0.225, y + 0.2, value, ha='center', va='center',
                    fontsize=20, fontweight='bold', color=color)
            # Label
            ax5.text(x + 0.225, y + 0.08, label, ha='center', va='center',
                    fontsize=10, color='black')

        ax5.set_xlim(0, 1)
        ax5.set_ylim(0, 1)
        ax5.set_title('Critical Healthcare Metrics', fontsize=14, fontweight='bold', y=1.05)

        # Save and add to slide
        img_stream = BytesIO()
        plt.tight_layout()
        plt.savefig(img_stream, format='PNG', dpi=150, bbox_inches='tight',
                   facecolor='white', edgecolor='none')
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, Inches(1), Inches(1), width=Inches(14))
        plt.close()

    def create_slide3_competition_infographic(self):
        """Slide 3: Competition Matrix - Visual"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        title = slide.shapes.add_textbox(Inches(0), Inches(0.2), Inches(16), Inches(0.6))
        title.text_frame.text = "Competitive Landscape & White Space"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.rgb_colors['secondary']
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Create competition infographic
        fig = plt.figure(figsize=(14, 7))

        # Main competitive matrix
        ax = fig.add_subplot(111)

        # Draw quadrant lines
        ax.axhline(y=0, color='black', linewidth=2)
        ax.axvline(x=0, color='black', linewidth=2)

        # Add quadrant backgrounds
        ax.add_patch(patches.Rectangle((-100, -100), 100, 100,
                                       facecolor=self.colors['light'], alpha=0.3))
        ax.add_patch(patches.Rectangle((0, -100), 100, 100,
                                       facecolor=self.colors['light'], alpha=0.3))
        ax.add_patch(patches.Rectangle((-100, 0), 100, 100,
                                       facecolor=self.colors['light'], alpha=0.3))
        ax.add_patch(patches.Rectangle((0, 0), 100, 100,
                                       facecolor=self.colors['accent2'], alpha=0.2))

        # Plot competitors
        competitors = {
            'Practo': (-60, -50, self.colors['secondary'], 80),
            'Apollo 24/7': (-40, -40, self.colors['secondary'], 100),
            '1mg': (-50, 30, self.colors['primary'], 120),
            'PharmEasy': (-30, 40, self.colors['primary'], 110),
            'Tata Health': (-45, 50, self.colors['primary'], 90),
            'eSanjeevani': (30, 20, self.colors['accent1'], 150),
            'Local Clinics': (40, -30, self.colors['dark'], 60),
            'WHITE SPACE': (60, 60, self.colors['accent2'], 200)
        }

        for name, (x, y, color, size) in competitors.items():
            if name == 'WHITE SPACE':
                # Special styling for white space
                circle = plt.Circle((x, y), size/10, color=color, alpha=0.8,
                                   edgecolor='white', linewidth=3)
                ax.add_patch(circle)
                ax.text(x, y, name, ha='center', va='center',
                       fontsize=14, fontweight='bold', color='white')
            else:
                circle = plt.Circle((x, y), size/15, color=color, alpha=0.6,
                                   edgecolor='white', linewidth=2)
                ax.add_patch(circle)
                ax.text(x, y, name, ha='center', va='center',
                       fontsize=10, fontweight='bold')

        # Axis labels
        ax.set_xlim(-100, 100)
        ax.set_ylim(-100, 100)
        ax.set_xlabel('← Urban                    Geography                    Rural →',
                     fontsize=14, fontweight='bold')
        ax.set_ylabel('← Narrow                Service Breadth                Broad →',
                     fontsize=14, fontweight='bold')
        ax.set_title('Competitive Positioning Matrix', fontsize=16, fontweight='bold', pad=20)

        # Add annotations for quadrants
        ax.text(-50, 75, 'Urban + Broad\n(Established Players)',
               ha='center', fontsize=11, style='italic', alpha=0.7)
        ax.text(50, 75, 'Rural + Broad\n(Opportunity Zone)',
               ha='center', fontsize=11, style='italic', alpha=0.7)
        ax.text(-50, -75, 'Urban + Narrow\n(Specialists)',
               ha='center', fontsize=11, style='italic', alpha=0.7)
        ax.text(50, -75, 'Rural + Narrow\n(Traditional)',
               ha='center', fontsize=11, style='italic', alpha=0.7)

        # Add key insights as callout boxes
        insights = [
            'Vernacular-first approach needed',
            'Physical touchpoints critical',
            'Affordable pricing (<₹100)',
            'Trust through local partnerships'
        ]

        for i, insight in enumerate(insights):
            y_pos = 40 - i * 25
            ax.add_patch(patches.FancyBboxPatch((65, y_pos), 30, 12,
                                               boxstyle="round,pad=0.05",
                                               facecolor=self.colors['accent1'],
                                               alpha=0.2, edgecolor=self.colors['accent1'],
                                               linewidth=2))
            ax.text(80, y_pos + 6, f'✓ {insight}', ha='center', va='center',
                   fontsize=9, fontweight='bold')

        # Remove tick marks
        ax.set_xticks([])
        ax.set_yticks([])

        # Save and add to slide
        img_stream = BytesIO()
        plt.tight_layout()
        plt.savefig(img_stream, format='PNG', dpi=150, bbox_inches='tight',
                   facecolor='white', edgecolor='none')
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, Inches(1), Inches(1), width=Inches(14))
        plt.close()

    def create_slide4_solution_infographic(self):
        """Slide 4: MediChain Solution - Visual Architecture"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        title = slide.shapes.add_textbox(Inches(0), Inches(0.2), Inches(16), Inches(0.6))
        title.text_frame.text = "MediChain: Integrated Healthcare Solution"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.rgb_colors['accent1']
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Create solution architecture infographic
        fig = plt.figure(figsize=(14, 7))
        ax = fig.add_subplot(111)
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 10)
        ax.axis('off')

        # Central hub - MediChain
        center = plt.Circle((5, 5), 1.5, color=self.colors['accent1'],
                          alpha=0.9, edgecolor='white', linewidth=3)
        ax.add_patch(center)
        ax.text(5, 5, 'MediChain\nPlatform', ha='center', va='center',
               fontsize=16, fontweight='bold', color='white')

        # Four solution pillars around the center
        components = [
            (2, 8, 'AI Triage', '🤖\nVernacular\n<₹20', self.colors['secondary']),
            (8, 8, 'IoT Kiosks', '📊\nDiagnostics\n₹1L/unit', self.colors['primary']),
            (2, 2, 'Blockchain', '🔐\nHealth Records\nNDHM', self.colors['accent2']),
            (8, 2, 'Pharmacy', '🏪\nLast-mile\nDelivery', self.colors['accent3'])
        ]

        for x, y, title, desc, color in components:
            # Component circle
            comp_circle = plt.Circle((x, y), 1.2, color=color, alpha=0.7,
                                    edgecolor='white', linewidth=2)
            ax.add_patch(comp_circle)

            # Component text
            ax.text(x, y + 0.3, title, ha='center', va='center',
                   fontsize=13, fontweight='bold', color='white')
            ax.text(x, y - 0.3, desc, ha='center', va='center',
                   fontsize=9, color='white')

            # Connect to center
            ax.plot([x, 5], [y, 5], 'k--', alpha=0.3, linewidth=2)

        # User journey flow (bottom)
        journey_y = 0.5
        journey_stages = [
            (1, 'Symptom', self.colors['secondary']),
            (3, 'Triage', self.colors['primary']),
            (5, 'Consult', self.colors['accent1']),
            (7, 'Diagnose', self.colors['accent2']),
            (9, 'Deliver', self.colors['accent3'])
        ]

        for i, (x, stage, color) in enumerate(journey_stages):
            # Stage circle
            stage_circle = plt.Circle((x, journey_y), 0.4, color=color, alpha=0.8)
            ax.add_patch(stage_circle)
            ax.text(x, journey_y, str(i+1), ha='center', va='center',
                   fontsize=12, fontweight='bold', color='white')
            ax.text(x, journey_y - 0.7, stage, ha='center', va='center',
                   fontsize=10, fontweight='bold')

            # Connect stages
            if i < len(journey_stages) - 1:
                ax.arrow(x + 0.5, journey_y, 1, 0, head_width=0.1,
                        head_length=0.1, fc='gray', ec='gray', alpha=0.5)

        # Key metrics (top right)
        metrics_data = [
            ('₹499/year', 'Family Plan'),
            ('<₹100', 'Per Consult'),
            ('5 min', 'Response Time'),
            ('24x7', 'Availability')
        ]

        for i, (value, label) in enumerate(metrics_data):
            x = 9.5
            y = 9 - i * 0.8
            ax.text(x, y, f'{value}', ha='right', fontsize=11,
                   fontweight='bold', color=self.colors['accent1'])
            ax.text(x + 0.1, y, f' {label}', ha='left', fontsize=9,
                   color=self.colors['dark'])

        ax.set_title('Integrated Healthcare Ecosystem', fontsize=16,
                    fontweight='bold', y=1.05)

        # Save and add to slide
        img_stream = BytesIO()
        plt.tight_layout()
        plt.savefig(img_stream, format='PNG', dpi=150, bbox_inches='tight',
                   facecolor='white', edgecolor='none')
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, Inches(1), Inches(1), width=Inches(14))
        plt.close()

    def create_slide5_impact_infographic(self):
        """Slide 5: Impact & Scale - Visual Roadmap"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        title = slide.shapes.add_textbox(Inches(0), Inches(0.2), Inches(16), Inches(0.6))
        title.text_frame.text = "5-Year Impact Journey"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.rgb_colors['primary']
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Create impact visualization
        fig = plt.figure(figsize=(14, 7))
        gs = fig.add_gridspec(2, 2, height_ratios=[1, 1], width_ratios=[1, 1],
                             hspace=0.3, wspace=0.3)

        # 1. Growth Timeline
        ax1 = fig.add_subplot(gs[0, :])

        years = np.array([1, 2, 3, 4, 5])
        users = np.array([50, 300, 1000, 10000, 100000])  # in thousands
        kiosks = np.array([100, 500, 2500, 10000, 50000])

        # Plot with log scale
        ax1_2 = ax1.twinx()

        line1 = ax1.semilogy(years, users, marker='o', linewidth=3,
                            color=self.colors['primary'], markersize=10, label='Users (K)')
        line2 = ax1_2.semilogy(years, kiosks, marker='s', linewidth=3,
                              color=self.colors['accent1'], markersize=10, label='Kiosks')

        # Milestone annotations
        milestones = ['PILOT', 'EXPAND', 'SCALE', 'GROWTH', 'NATIONAL']
        for i, (yr, milestone) in enumerate(zip(years, milestones)):
            ax1.text(yr, users[i]*1.5, milestone, ha='center',
                    fontsize=10, fontweight='bold', color=self.colors['dark'])

        ax1.set_xlabel('Year', fontsize=12, fontweight='bold')
        ax1.set_ylabel('Users (Thousands)', fontsize=12, fontweight='bold',
                      color=self.colors['primary'])
        ax1_2.set_ylabel('Kiosks', fontsize=12, fontweight='bold',
                        color=self.colors['accent1'])
        ax1.set_title('Exponential Growth Trajectory', fontsize=14, fontweight='bold')
        ax1.grid(True, alpha=0.3)
        ax1.set_xticks(years)

        # Color y-axis labels
        ax1.tick_params(axis='y', labelcolor=self.colors['primary'])
        ax1_2.tick_params(axis='y', labelcolor=self.colors['accent1'])

        # 2. Economic Impact
        ax2 = fig.add_subplot(gs[1, 0])

        metrics = ['Revenue\n(₹ Cr)', 'EBITDA\nMargin', 'Breakeven\nTime']
        values = [500, 35, 18]  # Example values
        colors_metrics = [self.colors['secondary'], self.colors['accent1'], self.colors['accent2']]

        bars = ax2.bar(metrics, values, color=colors_metrics,
                      edgecolor='white', linewidth=2)

        # Add value labels
        for i, (bar, value) in enumerate(zip(bars, values)):
            height = bar.get_height()
            if i == 1:  # EBITDA Margin
                label = f'{value}%'
            elif i == 2:  # Breakeven Time
                label = f'{value}mo'
            else:  # Revenue
                label = f'₹{value}Cr'
            ax2.text(bar.get_x() + bar.get_width()/2., height,
                    label, ha='center', va='bottom', fontsize=11, fontweight='bold')

        ax2.set_title('Economic Sustainability', fontsize=13, fontweight='bold')
        ax2.set_ylim(0, max(values) * 1.2)
        ax2.set_ylabel('Value', fontsize=11)

        # 3. Social Impact Radar Chart
        ax3 = fig.add_subplot(gs[1, 1], projection='polar')

        categories = ['Lives\nImpacted', 'Women\nHealth', 'Preventive\nCare',
                     'Mental\nHealth', 'Rural\nReach']
        values_radar = [100, 50, 80, 60, 75]  # Percentage of target achieved

        angles = np.linspace(0, 2 * np.pi, len(categories), endpoint=False).tolist()
        values_radar += values_radar[:1]  # Complete the circle
        angles += angles[:1]

        ax3.plot(angles, values_radar, 'o-', linewidth=2,
                color=self.colors['accent1'], markersize=8)
        ax3.fill(angles, values_radar, alpha=0.25, color=self.colors['accent1'])

        ax3.set_xticks(angles[:-1])
        ax3.set_xticklabels(categories, fontsize=10)
        ax3.set_ylim(0, 100)
        ax3.set_title('Social Impact Metrics (%)', fontsize=13,
                     fontweight='bold', y=1.08)
        ax3.grid(True)

        # Overall title
        fig.suptitle('Creating Sustainable Healthcare for 100M+ Indians',
                    fontsize=16, fontweight='bold', y=1.02)

        # Save and add to slide
        img_stream = BytesIO()
        plt.tight_layout()
        plt.savefig(img_stream, format='PNG', dpi=150, bbox_inches='tight',
                   facecolor='white', edgecolor='none')
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, Inches(1), Inches(1), width=Inches(14))
        plt.close()

        # Add footer
        footer = slide.shapes.add_textbox(Inches(0), Inches(8.3), Inches(16), Inches(0.4))
        footer.text_frame.text = "Scalable • Sustainable • Socially Impactful"
        footer.text_frame.paragraphs[0].font.size = Pt(14)
        footer.text_frame.paragraphs[0].font.bold = True
        footer.text_frame.paragraphs[0].font.color.rgb = self.rgb_colors['primary']
        footer.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def generate_presentation(self):
        """Generate all slides"""
        print("\n📊 Creating MediChain Infographic Presentation...")
        print("=" * 60)
        print("Focus: Clean infographics with minimal text")
        print("=" * 60)

        self.create_slide1_infographic()
        print("✅ Slide 1: Opportunity - Multi-chart infographic")

        self.create_slide2_healthcare_infographic()
        print("✅ Slide 2: Healthcare - Visual metrics dashboard")

        self.create_slide3_competition_infographic()
        print("✅ Slide 3: Competition - Visual positioning matrix")

        self.create_slide4_solution_infographic()
        print("✅ Slide 4: Solution - Architecture diagram")

        self.create_slide5_impact_infographic()
        print("✅ Slide 5: Impact - Growth visualization")

        # Save
        filename = "PPT Generated/MediChain_Infographic_Presentation.pptx"
        self.prs.save(filename)

        print("=" * 60)
        print(f"✅ Saved: {filename}")
        print("\n🎨 Infographic Features:")
        print("  • Professional charts and graphs")
        print("  • Visual data representation")
        print("  • Minimal text, maximum visuals")
        print("  • Clean, modern design")
        print("  • Color-coded information")
        print("=" * 60)

        return filename

if __name__ == "__main__":
    creator = MediChainInfographicPresentation()
    creator.generate_presentation()