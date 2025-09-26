#!/usr/bin/env python3
"""
Ultra-Condensed PPT System for Case Competitions
Creates 3-slide or 5-slide presentations with maximum impact
Inspired by IIM Ranchi and XIMB SRC best practices
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd
import numpy as np
import io
from datetime import datetime

class UltraCondensedPPT:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 widescreen
        self.prs.slide_height = Inches(7.5)
        
        # Professional color scheme
        self.colors = {
            'primary': RGBColor(0, 102, 204),     # Professional Blue
            'secondary': RGBColor(255, 103, 31),  # Orange (XIMB style)
            'accent': RGBColor(0, 176, 240),      # Light Blue
            'success': RGBColor(0, 176, 80),      # Green
            'warning': RGBColor(255, 192, 0),     # Yellow
            'danger': RGBColor(220, 53, 69),      # Red
            'purple': RGBColor(102, 45, 145),     # Purple
            'dark': RGBColor(51, 51, 51),         # Dark Gray
            'gray': RGBColor(128, 128, 128),      # Medium Gray
            'light': RGBColor(241, 241, 241),     # Light Gray
            'white': RGBColor(255, 255, 255)      # White
        }
        
    def create_3_slide_presentation(self, case_data):
        """Create ultra-condensed 3-slide presentation"""
        
        # Slide 1: Problem & Opportunity
        self._create_problem_opportunity_slide(case_data)
        
        # Slide 2: Solution & Analysis
        self._create_solution_analysis_slide(case_data)
        
        # Slide 3: Implementation & Impact
        self._create_implementation_impact_slide(case_data)
        
        return self.prs
    
    def create_5_slide_presentation(self, case_data):
        """Create condensed 5-slide presentation"""
        
        # Slide 1: Title & Context
        self._create_title_context_slide(case_data)
        
        # Slide 2: Problem Analysis
        self._create_problem_analysis_slide(case_data)
        
        # Slide 3: Strategic Solution
        self._create_strategic_solution_slide(case_data)
        
        # Slide 4: Implementation Roadmap
        self._create_implementation_roadmap_slide(case_data)
        
        # Slide 5: Financial Impact & Recommendations
        self._create_impact_recommendations_slide(case_data)
        
        return self.prs
    
    def _create_title_context_slide(self, case_data):
        """5-slide format: Title with context setting"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title section
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(1))
        tf = title_box.text_frame
        tf.text = case_data['title'].upper()
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(0.6))
        tf = subtitle_box.text_frame
        tf.text = case_data['subtitle']
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = self.colors['secondary']
        p.alignment = PP_ALIGN.CENTER
        
        # Context stats boxes
        stat_width = 3.5
        stat_height = 1.2
        y_pos = 3.8
        
        for i, stat in enumerate(case_data['context_stats'][:3]):
            x_pos = 1.5 + (stat_width + 0.5) * i
            
            stat_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(y_pos),
                Inches(stat_width), Inches(stat_height)
            )
            stat_box.fill.solid()
            colors = [self.colors['primary'], self.colors['secondary'], self.colors['purple']]
            stat_box.fill.fore_color.rgb = colors[i]
            stat_box.line.fill.background()
            
            tf = stat_box.text_frame
            p = tf.paragraphs[0]
            p.text = stat
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            
        # Team info
        team_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.3), Inches(0.5))
        tf = team_box.text_frame
        tf.text = case_data['team']
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = self.colors['gray']
        p.alignment = PP_ALIGN.CENTER
        
    def _create_problem_analysis_slide(self, case_data):
        """5-slide format: Deep problem analysis"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6))
        tf = title_box.text_frame
        tf.text = "PROBLEM ANALYSIS"
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        # Problem statement banner
        statement_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1),
            Inches(12.3), Inches(0.8)
        )
        statement_box.fill.solid()
        statement_box.fill.fore_color.rgb = self.colors['danger']
        statement_box.line.fill.background()
        
        tf = statement_box.text_frame
        p = tf.paragraphs[0]
        p.text = case_data['problem_statement']
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # Problem factors (3 columns)
        factor_y = 2.2
        factor_width = 3.8
        
        colors = [self.colors['warning'], self.colors['danger'], self.colors['purple']]
        
        for i, (factor, items) in enumerate(case_data['problem_factors'].items()):
            x_pos = 0.5 + (factor_width + 0.3) * i
            
            # Factor box
            factor_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(factor_y),
                Inches(factor_width), Inches(2)
            )
            factor_box.fill.solid()
            factor_box.fill.fore_color.rgb = self.colors['light']
            factor_box.line.color.rgb = colors[i]
            factor_box.line.width = Pt(2)
            
            tf = factor_box.text_frame
            tf.margin_all = Inches(0.2)
            
            # Factor title
            p = tf.paragraphs[0]
            p.text = factor.upper()
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = colors[i]
            p.alignment = PP_ALIGN.CENTER
            
            # Factor items
            for item in items:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(10)
                p.font.color.rgb = self.colors['dark']
                p.line_spacing = 1.15
                
        # Impact section
        impact_y = 4.5
        impact_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(impact_y),
            Inches(12.3), Inches(1.8)
        )
        impact_box.fill.solid()
        impact_box.fill.fore_color.rgb = self.colors['white']
        impact_box.line.color.rgb = self.colors['dark']
        impact_box.line.width = Pt(1)
        
        tf = impact_box.text_frame
        tf.margin_all = Inches(0.3)
        
        # Impact header
        p = tf.paragraphs[0]
        p.text = "BUSINESS IMPACT"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors['dark']
        
        # Impact metrics in columns
        p = tf.add_paragraph()
        impact_text = "   |   ".join([f"{k}: {v}" for k, v in case_data['problem_impact'].items()])
        p.text = impact_text
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.colors['danger']
        p.alignment = PP_ALIGN.CENTER
        
    def _create_strategic_solution_slide(self, case_data):
        """5-slide format: Strategic solution details"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6))
        tf = title_box.text_frame
        tf.text = "STRATEGIC SOLUTION"
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        # Solution name banner
        solution_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1),
            Inches(12.3), Inches(0.8)
        )
        solution_box.fill.solid()
        solution_box.fill.fore_color.rgb = self.colors['success']
        solution_box.line.fill.background()
        
        tf = solution_box.text_frame
        p = tf.paragraphs[0]
        p.text = case_data['solution_name']
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # Solution pillars
        pillar_y = 2
        pillar_width = 3.8
        
        for i, pillar in enumerate(case_data['solution_pillars']):
            x_pos = 0.5 + (pillar_width + 0.3) * i
            
            # Pillar box
            pillar_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(pillar_y),
                Inches(pillar_width), Inches(2.3)
            )
            colors = [self.colors['primary'], self.colors['secondary'], self.colors['purple']]
            pillar_box.fill.solid()
            pillar_box.fill.fore_color.rgb = colors[i]
            pillar_box.line.fill.background()
            
            tf = pillar_box.text_frame
            tf.margin_all = Inches(0.2)
            
            # Pillar name
            p = tf.paragraphs[0]
            p.text = pillar['name'].upper()
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            
            # Components
            for component in pillar['components'][:4]:
                p = tf.add_paragraph()
                p.text = f"✓ {component}"
                p.font.size = Pt(10)
                p.font.color.rgb = self.colors['white']
                p.line_spacing = 1.15
                
        # Competitive advantages
        adv_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(4.5),
            Inches(12.3), Inches(1.5)
        )
        adv_box.fill.solid()
        adv_box.fill.fore_color.rgb = self.colors['light']
        adv_box.line.color.rgb = self.colors['success']
        adv_box.line.width = Pt(1)
        
        tf = adv_box.text_frame
        tf.margin_all = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = "COMPETITIVE ADVANTAGES"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors['dark']
        
        # Advantages in grid
        p = tf.add_paragraph()
        adv_text = "  •  ".join(case_data['competitive_advantages'])
        p.text = adv_text
        p.font.size = Pt(11)
        p.font.color.rgb = self.colors['success']
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
    def _create_implementation_roadmap_slide(self, case_data):
        """5-slide format: Detailed implementation roadmap"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6))
        tf = title_box.text_frame
        tf.text = "IMPLEMENTATION ROADMAP"
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        # Timeline phases
        phase_y = 1.2
        phase_height = 2.8
        phase_width = 3.8
        
        colors = [self.colors['warning'], self.colors['secondary'], self.colors['success']]
        
        for i, phase in enumerate(case_data['phases']):
            x_pos = 0.5 + (phase_width + 0.3) * i
            
            # Phase box
            phase_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(phase_y),
                Inches(phase_width), Inches(phase_height)
            )
            phase_box.fill.solid()
            phase_box.fill.fore_color.rgb = self.colors['white']
            phase_box.line.color.rgb = colors[i]
            phase_box.line.width = Pt(2)
            
            # Phase header
            header_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x_pos), Inches(phase_y),
                Inches(phase_width), Inches(0.6)
            )
            header_box.fill.solid()
            header_box.fill.fore_color.rgb = colors[i]
            header_box.line.fill.background()
            
            tf = header_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{phase['phase'].upper()} • {phase['timeline']}"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            
            # Phase content
            content_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.2), Inches(phase_y + 0.7),
                Inches(phase_width - 0.4), Inches(phase_height - 0.9)
            )
            tf = content_box.text_frame
            
            # Targets
            p = tf.paragraphs[0]
            p.text = "TARGET"
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = colors[i]
            
            p = tf.add_paragraph()
            p.text = phase['targets']
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.colors['dark']
            
            # Milestones
            p = tf.add_paragraph()
            p.text = "\nKEY MILESTONES"
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = colors[i]
            
            for milestone in phase['milestones']:
                p = tf.add_paragraph()
                p.text = f"• {milestone}"
                p.font.size = Pt(9)
                p.font.color.rgb = self.colors['dark']
                p.line_spacing = 1.1
                
        # Partnerships section
        partner_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(4.3),
            Inches(12.3), Inches(1.8)
        )
        partner_box.fill.solid()
        partner_box.fill.fore_color.rgb = self.colors['light']
        partner_box.line.color.rgb = self.colors['primary']
        partner_box.line.width = Pt(1)
        
        tf = partner_box.text_frame
        tf.margin_all = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = "STRATEGIC PARTNERSHIPS"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        p.alignment = PP_ALIGN.CENTER
        
        # Partnership grid
        for i, partner in enumerate(case_data['partnerships']):
            if i % 2 == 0:
                p = tf.add_paragraph()
                p.text = ""
            
            run = p.add_run()
            run.text = f"  •  {partner}"
            run.font.size = Pt(11)
            run.font.color.rgb = self.colors['dark']
            
    def _create_impact_recommendations_slide(self, case_data):
        """5-slide format: Financial impact and recommendations"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6))
        tf = title_box.text_frame
        tf.text = "FINANCIAL IMPACT & NEXT STEPS"
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        # Financial metrics (left side)
        fin_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.2),
            Inches(6), Inches(2.5)
        )
        fin_box.fill.solid()
        fin_box.fill.fore_color.rgb = self.colors['success']
        fin_box.line.fill.background()
        
        tf = fin_box.text_frame
        tf.margin_all = Inches(0.3)
        
        # Header
        p = tf.paragraphs[0]
        p.text = "FINANCIAL PROJECTIONS (YEAR 3)"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        # Key metrics
        metrics = case_data['financial_projections']['key_metrics']
        for key, value in metrics.items():
            p = tf.add_paragraph()
            p.text = f"{key}: {value}"
            p.font.size = Pt(11)
            p.font.color.rgb = self.colors['white']
            p.font.bold = True if 'Revenue' in key or 'EBITDA' in key else False
            
        # Investment & ROI (right side)
        roi_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.8), Inches(1.2),
            Inches(6), Inches(2.5)
        )
        roi_box.fill.solid()
        roi_box.fill.fore_color.rgb = self.colors['primary']
        roi_box.line.fill.background()
        
        tf = roi_box.text_frame
        tf.margin_all = Inches(0.3)
        
        # ROI header
        p = tf.paragraphs[0]
        p.text = "INVESTMENT RETURNS"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        # ROI details
        p = tf.add_paragraph()
        p.text = f"Investment Required: {case_data['investment_required']}"
        p.font.size = Pt(11)
        p.font.color.rgb = self.colors['white']
        
        p = tf.add_paragraph()
        p.text = f"ROI: {case_data['roi']}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = f"Payback Period: {case_data['payback']}"
        p.font.size = Pt(11)
        p.font.color.rgb = self.colors['white']
        
        # Social impact
        if 'social_impact' in case_data:
            p = tf.add_paragraph()
            p.text = "\nSOCIAL IMPACT"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            
            for metric, value in list(case_data['social_impact'].items())[:2]:
                p = tf.add_paragraph()
                p.text = f"• {metric}: {value}"
                p.font.size = Pt(10)
                p.font.color.rgb = self.colors['white']
                
        # Recommendations section
        rec_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(4),
            Inches(12.3), Inches(2.3)
        )
        rec_box.fill.solid()
        rec_box.fill.fore_color.rgb = self.colors['light']
        rec_box.line.color.rgb = self.colors['dark']
        rec_box.line.width = Pt(1)
        
        tf = rec_box.text_frame
        tf.margin_all = Inches(0.3)
        
        # Recommendations header
        p = tf.paragraphs[0]
        p.text = "IMMEDIATE ACTION ITEMS"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors['dark']
        
        # Recommendations in 2 columns
        for i, rec in enumerate(case_data['recommendations']):
            if i % 2 == 0:
                p = tf.add_paragraph()
            
            run = p.add_run()
            run.text = f"   {i+1}. {rec}   "
            run.font.size = Pt(11)
            run.font.color.rgb = self.colors['dark']
            run.font.bold = True
    
    def _create_problem_opportunity_slide(self, case_data):
        """Slide 1 for 3-slide format: Problem & Opportunity combined"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title with company name
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6))
        tf = title_box.text_frame
        tf.text = f"{case_data['company']} | {case_data['title']}"
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        # Add navigation bar (XIMB style)
        self._add_navigation_bar(slide, ['PROBLEM', 'SOLUTION', 'IMPACT'], 0)
        
        # Problem Box (Left side - 40%)
        problem_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.5),
            Inches(5), Inches(5.5)
        )
        problem_box.fill.solid()
        problem_box.fill.fore_color.rgb = self.colors['light']
        problem_box.line.color.rgb = self.colors['danger']
        problem_box.line.width = Pt(2)
        
        tf = problem_box.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.3)
        tf.margin_right = Inches(0.3)
        
        # Problem header
        p = tf.paragraphs[0]
        p.text = "CORE PROBLEM"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.colors['danger']
        
        # Problem statement
        p = tf.add_paragraph()
        p.text = case_data['problem_statement']
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors['dark']
        p.line_spacing = 1.2
        
        # Key challenges
        p = tf.add_paragraph()
        p.text = "Key Challenges:"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors['dark']
        
        for challenge in case_data['challenges'][:3]:
            p = tf.add_paragraph()
            p.text = f"• {challenge}"
            p.font.size = Pt(11)
            p.font.color.rgb = self.colors['dark']
            p.line_spacing = 1.15
        
        # Market data
        if 'market_stats' in case_data:
            p = tf.add_paragraph()
            p.text = "Market Reality:"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.colors['dark']
            
            for stat in case_data['market_stats'][:2]:
                p = tf.add_paragraph()
                p.text = f"• {stat}"
                p.font.size = Pt(11)
                p.font.color.rgb = self.colors['danger']
                p.font.bold = True
        
        # Opportunity Box (Right side - 55%)
        opp_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.8), Inches(1.5),
            Inches(7), Inches(5.5)
        )
        opp_box.fill.solid()
        opp_box.fill.fore_color.rgb = self.colors['white']
        opp_box.line.color.rgb = self.colors['success']
        opp_box.line.width = Pt(2)
        
        tf = opp_box.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.3)
        tf.margin_right = Inches(0.3)
        
        # Opportunity header
        p = tf.paragraphs[0]
        p.text = "MARKET OPPORTUNITY"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.colors['success']
        
        # Opportunity size
        p = tf.add_paragraph()
        p.text = case_data['opportunity_size']
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        p.alignment = PP_ALIGN.CENTER
        
        # Add visual chart if data provided
        if 'opportunity_chart_data' in case_data:
            self._add_mini_chart(slide, case_data['opportunity_chart_data'], 
                               Inches(6.3), Inches(3.5), Inches(6), Inches(3))
        
        # Growth drivers
        p = tf.add_paragraph()
        p.text = "Growth Drivers:"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors['dark']
        
        for driver in case_data['growth_drivers'][:3]:
            p = tf.add_paragraph()
            p.text = f"✓ {driver}"
            p.font.size = Pt(11)
            p.font.color.rgb = self.colors['success']
            p.line_spacing = 1.15
            
    def _create_solution_analysis_slide(self, case_data):
        """Slide 2 for 3-slide format: Solution & Analysis"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6))
        tf = title_box.text_frame
        tf.text = "STRATEGIC SOLUTION & FRAMEWORK"
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        # Navigation bar
        self._add_navigation_bar(slide, ['PROBLEM', 'SOLUTION', 'IMPACT'], 1)
        
        # Solution Overview (Top section)
        solution_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.5),
            Inches(12.3), Inches(1.2)
        )
        solution_box.fill.solid()
        solution_box.fill.fore_color.rgb = self.colors['primary']
        solution_box.line.fill.background()
        
        tf = solution_box.text_frame
        tf.margin_left = Inches(0.5)
        tf.margin_right = Inches(0.5)
        p = tf.paragraphs[0]
        p.text = case_data['solution_statement']
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # Framework or Analysis (Middle section)
        if case_data.get('framework_type') == '3_pillars':
            self._add_3_pillar_framework(slide, case_data['framework_data'], 
                                        Inches(0.5), Inches(3), Inches(12.3), Inches(2.5))
        elif case_data.get('framework_type') == 'matrix':
            self._add_2x2_matrix_condensed(slide, case_data['framework_data'],
                                          Inches(0.5), Inches(3), Inches(6), Inches(3.5))
            # Add key strategies on the right
            self._add_key_strategies(slide, case_data['strategies'],
                                   Inches(7), Inches(3), Inches(5.8), Inches(3.5))
        else:
            # Default: Key initiatives
            self._add_key_initiatives(slide, case_data['initiatives'],
                                    Inches(0.5), Inches(3), Inches(12.3), Inches(3.5))
    
    def _create_implementation_impact_slide(self, case_data):
        """Slide 3 for 3-slide format: Implementation & Impact"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6))
        tf = title_box.text_frame
        tf.text = "IMPLEMENTATION ROADMAP & EXPECTED IMPACT"
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        # Navigation bar
        self._add_navigation_bar(slide, ['PROBLEM', 'SOLUTION', 'IMPACT'], 2)
        
        # Implementation Timeline (Left 60%)
        timeline_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.5),
            Inches(7.5), Inches(2.2)
        )
        timeline_box.fill.solid()
        timeline_box.fill.fore_color.rgb = self.colors['white']
        timeline_box.line.color.rgb = self.colors['gray']
        timeline_box.line.width = Pt(1)
        
        # Add timeline visual
        self._add_condensed_timeline(slide, case_data['timeline_phases'],
                                   Inches(0.7), Inches(1.7), Inches(7.1), Inches(1.8))
        
        # Financial Impact (Right 40%)
        impact_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.3), Inches(1.5),
            Inches(4.5), Inches(2.2)
        )
        impact_box.fill.solid()
        impact_box.fill.fore_color.rgb = self.colors['success']
        impact_box.line.fill.background()
        
        tf = impact_box.text_frame
        tf.margin_all = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = "EXPECTED IMPACT"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        # ROI
        p = tf.add_paragraph()
        p.text = f"ROI: {case_data['roi']}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        # Payback
        p = tf.add_paragraph()
        p.text = f"Payback: {case_data['payback_period']}"
        p.font.size = Pt(12)
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        # Key Metrics Dashboard (Bottom section)
        metrics_y = 4
        
        # Create metric boxes
        metrics = case_data['key_metrics'][:4]  # Max 4 metrics
        metric_width = 2.8
        metric_spacing = 0.2
        
        for i, metric in enumerate(metrics):
            x_pos = 0.5 + (metric_width + metric_spacing) * i
            
            metric_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(metrics_y),
                Inches(metric_width), Inches(1.2)
            )
            metric_box.fill.solid()
            colors = [self.colors['primary'], self.colors['secondary'], 
                     self.colors['purple'], self.colors['warning']]
            metric_box.fill.fore_color.rgb = colors[i % len(colors)]
            metric_box.line.fill.background()
            
            tf = metric_box.text_frame
            tf.margin_all = Inches(0.15)
            
            # Metric name
            p = tf.paragraphs[0]
            p.text = metric['name']
            p.font.size = Pt(11)
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            
            # Metric value
            p = tf.add_paragraph()
            p.text = metric['value']
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            
        # Recommendations box (Bottom)
        rec_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(5.5),
            Inches(12.3), Inches(1.3)
        )
        rec_box.fill.solid()
        rec_box.fill.fore_color.rgb = self.colors['light']
        rec_box.line.color.rgb = self.colors['primary']
        rec_box.line.width = Pt(1)
        
        tf = rec_box.text_frame
        tf.margin_all = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = "KEY RECOMMENDATIONS: "
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        # Add recommendations as continuous text
        rec_text = " | ".join([f"→ {rec}" for rec in case_data['recommendations'][:3]])
        run = p.add_run()
        run.text = rec_text
        run.font.size = Pt(11)
        run.font.color.rgb = self.colors['dark']
        run.font.bold = False
        
    def _add_navigation_bar(self, slide, sections, current_index):
        """Add XIMB-style navigation bar"""
        bar_y = 1
        bar_height = 0.3
        total_width = 12.3
        section_width = total_width / len(sections)
        
        for i, section in enumerate(sections):
            x_pos = 0.5 + section_width * i
            
            # Create arrow shape
            if i == current_index:
                color = self.colors['secondary']  # Orange for current
            else:
                color = self.colors['gray']
                
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                Inches(x_pos), Inches(bar_y),
                Inches(section_width - 0.05), Inches(bar_height)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = color
            arrow.line.fill.background()
            
            # Add text
            tf = arrow.text_frame
            tf.margin_all = Inches(0)
            p = tf.paragraphs[0]
            p.text = section
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            
    def _add_3_pillar_framework(self, slide, pillars, x, y, width, height):
        """Add 3-pillar framework like IIM Ranchi"""
        pillar_width = width / 3 - Inches(0.1)
        
        colors = [self.colors['primary'], self.colors['success'], self.colors['secondary']]
        
        for i, pillar in enumerate(pillars):
            x_pos = x + (width / 3) * i
            
            # Pillar box
            pillar_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_pos, y,
                pillar_width, height
            )
            pillar_box.fill.solid()
            pillar_box.fill.fore_color.rgb = colors[i % len(colors)]
            pillar_box.line.fill.background()
            
            tf = pillar_box.text_frame
            tf.margin_all = Inches(0.2)
            
            # Pillar title
            p = tf.paragraphs[0]
            p.text = pillar['title'].upper()
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            
            # Pillar items
            for item in pillar['items'][:4]:  # Max 4 items
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(10)
                p.font.color.rgb = self.colors['white']
                p.line_spacing = 1.1
                
    def _add_2x2_matrix_condensed(self, slide, matrix_data, x, y, width, height):
        """Add condensed 2x2 matrix"""
        # Draw cross lines
        h_line = slide.shapes.add_connector(
            1, x + Inches(0.5), y + height/2, 
            x + width - Inches(0.5), y + height/2
        )
        h_line.line.color.rgb = self.colors['dark']
        h_line.line.width = Pt(1.5)
        
        v_line = slide.shapes.add_connector(
            1, x + width/2, y + Inches(0.5),
            x + width/2, y + height - Inches(0.5)
        )
        v_line.line.color.rgb = self.colors['dark']
        v_line.line.width = Pt(1.5)
        
        # Add quadrant boxes
        quadrants = [
            {'x': x, 'y': y, 'data': matrix_data[0]},  # Top Left
            {'x': x + width/2, 'y': y, 'data': matrix_data[1]},  # Top Right
            {'x': x, 'y': y + height/2, 'data': matrix_data[2]},  # Bottom Left
            {'x': x + width/2, 'y': y + height/2, 'data': matrix_data[3]}  # Bottom Right
        ]
        
        colors = [self.colors['warning'], self.colors['success'], 
                  self.colors['danger'], self.colors['primary']]
        
        for i, quad in enumerate(quadrants):
            # Create text box for each quadrant
            text_box = slide.shapes.add_textbox(
                quad['x'] + Inches(0.1), quad['y'] + Inches(0.1),
                width/2 - Inches(0.2), height/2 - Inches(0.2)
            )
            
            tf = text_box.text_frame
            tf.margin_all = Inches(0.1)
            
            # Quadrant title
            p = tf.paragraphs[0]
            p.text = quad['data']['title']
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = colors[i]
            p.alignment = PP_ALIGN.CENTER
            
            # Items (condensed)
            for item in quad['data']['items'][:2]:  # Max 2 items
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(9)
                p.font.color.rgb = self.colors['dark']
                
    def _add_mini_chart(self, slide, chart_data, x, y, width, height):
        """Add small chart visualization"""
        # Create matplotlib figure
        fig, ax = plt.subplots(figsize=(width.inches, height.inches))
        
        if chart_data['type'] == 'bar':
            bars = ax.bar(chart_data['labels'], chart_data['values'], 
                          color='#0066CC')
            
            # Add value labels on bars
            for bar in bars:
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2., height,
                       f'{height:.0f}', ha='center', va='bottom', fontsize=8)
                       
        elif chart_data['type'] == 'line':
            ax.plot(chart_data['x'], chart_data['y'], 
                   marker='o', linewidth=2, markersize=6,
                   color='#0066CC')
            ax.fill_between(chart_data['x'], chart_data['y'], alpha=0.3)
            
        # Remove spines and labels for clean look
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_visible(False)
        ax.spines['bottom'].set_visible(True)
        ax.yaxis.set_visible(False)
        ax.xaxis.set_ticks_position('none')
        plt.xticks(fontsize=8)
        
        # Grid
        ax.yaxis.grid(True, linestyle='-', alpha=0.2)
        ax.set_axisbelow(True)
        
        plt.tight_layout()
        
        # Save and add to slide
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=200, bbox_inches='tight', 
                   facecolor='white', edgecolor='none')
        plt.close()
        
        img_stream.seek(0)
        pic = slide.shapes.add_picture(img_stream, x, y, width=width)
        
    def _add_condensed_timeline(self, slide, phases, x, y, width, height):
        """Add condensed timeline visualization"""
        phase_width = width / len(phases)
        
        # Timeline line
        line_y = y + height/2
        timeline_line = slide.shapes.add_connector(
            1, x, line_y, x + width, line_y
        )
        timeline_line.line.color.rgb = self.colors['primary']
        timeline_line.line.width = Pt(3)
        
        for i, phase in enumerate(phases):
            x_pos = x + phase_width * i + phase_width/2 - Inches(0.3)
            
            # Phase circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x_pos, line_y - Inches(0.15),
                Inches(0.3), Inches(0.3)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.colors['secondary']
            circle.line.color.rgb = self.colors['white']
            circle.line.width = Pt(2)
            
            # Phase name (above)
            name_box = slide.shapes.add_textbox(
                x_pos - Inches(0.5), y,
                Inches(1.3), Inches(0.4)
            )
            tf = name_box.text_frame
            p = tf.paragraphs[0]
            p.text = phase['name']
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = self.colors['dark']
            p.alignment = PP_ALIGN.CENTER
            
            # Phase duration (below)
            duration_box = slide.shapes.add_textbox(
                x_pos - Inches(0.5), line_y + Inches(0.3),
                Inches(1.3), Inches(0.3)
            )
            tf = duration_box.text_frame
            p = tf.paragraphs[0]
            p.text = phase['duration']
            p.font.size = Pt(9)
            p.font.color.rgb = self.colors['gray']
            p.alignment = PP_ALIGN.CENTER
            
    def _add_key_initiatives(self, slide, initiatives, x, y, width, height):
        """Add key initiatives in a structured layout"""
        # Create container
        container = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x, y, width, height
        )
        container.fill.background()
        container.line.color.rgb = self.colors['gray']
        container.line.width = Pt(1)
        
        # Add initiatives in a grid
        cols = 3
        rows = (len(initiatives) + cols - 1) // cols
        
        init_width = width / cols - Inches(0.2)
        init_height = height / rows - Inches(0.2)
        
        for i, initiative in enumerate(initiatives[:6]):  # Max 6 initiatives
            row = i // cols
            col = i % cols
            
            x_pos = x + Inches(0.1) + (init_width + Inches(0.1)) * col
            y_pos = y + Inches(0.1) + (init_height + Inches(0.1)) * row
            
            # Initiative box
            init_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_pos, y_pos,
                init_width, init_height
            )
            init_box.fill.solid()
            colors = [self.colors['primary'], self.colors['secondary'], 
                     self.colors['purple'], self.colors['success'],
                     self.colors['warning'], self.colors['danger']]
            init_box.fill.fore_color.rgb = colors[i % len(colors)]
            init_box.line.fill.background()
            
            tf = init_box.text_frame
            tf.margin_all = Inches(0.15)
            
            # Initiative title
            p = tf.paragraphs[0]
            p.text = initiative['title']
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            
            # Initiative impact
            if 'impact' in initiative:
                p = tf.add_paragraph()
                p.text = initiative['impact']
                p.font.size = Pt(9)
                p.font.color.rgb = self.colors['white']
                p.alignment = PP_ALIGN.CENTER
                
    def _add_key_strategies(self, slide, strategies, x, y, width, height):
        """Add key strategies list"""
        strat_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x, y, width, height
        )
        strat_box.fill.solid()
        strat_box.fill.fore_color.rgb = self.colors['light']
        strat_box.line.color.rgb = self.colors['primary']
        strat_box.line.width = Pt(1)
        
        tf = strat_box.text_frame
        tf.margin_all = Inches(0.3)
        
        # Header
        p = tf.paragraphs[0]
        p.text = "KEY STRATEGIES"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        p.alignment = PP_ALIGN.CENTER
        
        # Strategies
        for i, strategy in enumerate(strategies[:5], 1):
            p = tf.add_paragraph()
            p.text = f"{i}. {strategy}"
            p.font.size = Pt(11)
            p.font.color.rgb = self.colors['dark']
            p.line_spacing = 1.2
            
    def save(self, filename):
        """Save the presentation"""
        self.prs.save(filename)
        return filename


# Example usage functions
def create_3_slide_example():
    """Create example 3-slide presentation"""
    ppt = UltraCondensedPPT()
    
    case_data = {
        'company': 'TechCorp India',
        'title': 'Digital Transformation in Traditional Retail',
        
        # Slide 1 data
        'problem_statement': 'Traditional retail losing 25% market share annually to e-commerce giants',
        'challenges': [
            'Legacy systems preventing omnichannel integration',
            '70% of customer touchpoints remain offline',
            'Inventory visibility limited to 30%'
        ],
        'market_stats': [
            '$150B e-commerce market by 2025',
            '45% CAGR in online retail'
        ],
        'opportunity_size': '$2.5 BILLION',
        'opportunity_chart_data': {
            'type': 'bar',
            'labels': ['2023', '2024', '2025', '2026'],
            'values': [500, 1000, 1800, 2500]
        },
        'growth_drivers': [
            'Smartphone penetration at 78%',
            'Digital payment adoption surge',
            'Post-COVID behavioral shift'
        ],
        
        # Slide 2 data
        'solution_statement': 'Unified Commerce Platform: Seamless integration of online & offline channels',
        'framework_type': '3_pillars',
        'framework_data': [
            {
                'title': 'Digital Core',
                'items': [
                    'Cloud-native architecture',
                    'Real-time inventory',
                    'AI-powered analytics',
                    'Mobile-first design'
                ]
            },
            {
                'title': 'Customer Experience',
                'items': [
                    'Unified loyalty program',
                    '360° customer view',
                    'Personalization engine',
                    'Omnichannel support'
                ]
            },
            {
                'title': 'Operations Excellence',
                'items': [
                    'Supply chain optimization',
                    'Smart warehousing',
                    'Last-mile delivery',
                    'Predictive maintenance'
                ]
            }
        ],
        
        # Slide 3 data
        'timeline_phases': [
            {'name': 'Foundation', 'duration': '0-3 months'},
            {'name': 'Pilot', 'duration': '3-6 months'},
            {'name': 'Scale', 'duration': '6-12 months'},
            {'name': 'Optimize', 'duration': '12-18 months'}
        ],
        'roi': '320%',
        'payback_period': '14 months',
        'key_metrics': [
            {'name': 'Revenue Uplift', 'value': '+45%'},
            {'name': 'Cost Reduction', 'value': '-30%'},
            {'name': 'Customer NPS', 'value': '+25pts'},
            {'name': 'Market Share', 'value': '+12%'}
        ],
        'recommendations': [
            'Establish Digital Transformation Office',
            'Partner with leading tech providers',
            'Launch pilot in top 5 metros'
        ]
    }
    
    ppt.create_3_slide_presentation(case_data)
    filename = "Ultra_Condensed_3_Slide_Presentation.pptx"
    ppt.save(filename)
    print(f"✓ Created 3-slide presentation: {filename}")
    return filename


def create_5_slide_example():
    """Create example 5-slide presentation"""
    ppt = UltraCondensedPPT()
    
    case_data = {
        'company': 'FinTech Solutions',
        'title': 'Democratizing Financial Services in Rural India',
        
        # Slide 1 - Title & Context
        'subtitle': 'Bridging the $380B financial inclusion gap',
        'team': 'SIMSREE Team Alpha',
        'context_stats': [
            '400M+ unbanked Indians',
            '65% rural population',
            'Only 18% credit penetration'
        ],
        
        # Slide 2 - Problem Analysis  
        'problem_statement': 'Traditional banking fails to serve 65% of India due to high costs and accessibility barriers',
        'problem_factors': {
            'Accessibility': ['2+ hour travel to nearest bank', 'Limited branch coverage', 'No 24/7 availability'],
            'Affordability': ['High minimum balance', 'Transaction fees', 'Documentation costs'],
            'Awareness': ['Low financial literacy', 'Language barriers', 'Complex products']
        },
        'problem_impact': {
            'Economic': '$380B credit gap',
            'Social': '250M excluded from formal economy',
            'Growth': '2.5% GDP loss annually'
        },
        
        # Slide 3 - Strategic Solution
        'solution_name': 'BharatFin: AI-Powered Mobile Banking Platform',
        'solution_pillars': [
            {
                'name': 'Technology Stack',
                'components': ['Voice-based AI interface', 'Offline transaction capability', 'Blockchain security', 'Biometric authentication']
            },
            {
                'name': 'Distribution Model', 
                'components': ['10,000 village entrepreneurs', 'Kirana store partnerships', 'Mobile van banking', 'WhatsApp banking']
            },
            {
                'name': 'Product Suite',
                'components': ['Micro-savings (₹10 start)', 'Instant micro-loans', 'Crop insurance', 'Digital payments']
            }
        ],
        'competitive_advantages': ['95% lower cost structure', 'Regional language support', 'No documentation needed', 'Instant account opening'],
        
        # Slide 4 - Implementation Roadmap
        'phases': [
            {
                'phase': 'Pilot',
                'timeline': 'Months 1-3',
                'targets': '5 districts, 50K users',
                'milestones': ['Technology deployment', 'Partner onboarding', 'Regulatory approval']
            },
            {
                'phase': 'Scale', 
                'timeline': 'Months 4-9',
                'targets': '50 districts, 1M users',
                'milestones': ['Feature expansion', 'Credit products launch', 'Break-even achieved']
            },
            {
                'phase': 'National',
                'timeline': 'Months 10-18',
                'targets': '500 districts, 10M users', 
                'milestones': ['Pan-India coverage', 'Profitability', 'IPO preparation']
            }
        ],
        'partnerships': ['NPCI for payments', 'Regional banks for capital', 'Telecom for connectivity', 'NGOs for literacy'],
        
        # Slide 5 - Financial Impact & Recommendations
        'financial_projections': {
            'revenue_streams': [
                {'source': 'Transaction fees', 'year3': '₹450Cr'},
                {'source': 'Lending income', 'year3': '₹850Cr'}, 
                {'source': 'Insurance commission', 'year3': '₹200Cr'}
            ],
            'key_metrics': {
                'Total Revenue Y3': '₹1,500 Cr',
                'EBITDA Margin': '35%',
                'CAC': '₹50',
                'LTV': '₹2,500'
            }
        },
        'investment_required': '₹250 Cr',
        'roi': '450%',
        'payback': '18 months',
        'social_impact': {
            'Financial Inclusion': '10M new accounts',
            'Credit Access': '₹5,000Cr disbursed',
            'Jobs Created': '50,000 agents'
        },
        'recommendations': [
            'Secure Series A funding of ₹250Cr by Q1',
            'Partner with government schemes (PM Jan Dhan)',
            'Launch pilot in Maharashtra & Karnataka',
            'Build tech team of 50+ engineers'
        ]
    }
    
    # Create the 5-slide presentation
    ppt._create_title_context_slide(case_data)
    ppt._create_problem_analysis_slide(case_data)
    ppt._create_strategic_solution_slide(case_data)
    ppt._create_implementation_roadmap_slide(case_data)
    ppt._create_impact_recommendations_slide(case_data)
    
    filename = "Ultra_Condensed_5_Slide_Presentation.pptx"
    ppt.save(filename)
    print(f"✓ Created 5-slide presentation: {filename}")
    return filename


if __name__ == "__main__":
    print("Creating ultra-condensed case competition presentations...")
    
    # Create 3-slide example
    create_3_slide_example()
    
    # Create 5-slide example
    create_5_slide_example()
    
    print("\n✓ Ultra-condensed PPT system ready!")
    print("✓ Created both 3-slide and 5-slide examples")
    print("✓ Maximizes information density while maintaining clarity")
    print("✓ Inspired by IIM Ranchi and XIMB SRC best practices")