#!/usr/bin/env python3
"""
Rich Visual PPT System - Creates visually dense, infographic-heavy presentations
No empty spaces - Every slide is packed with professional visuals
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, Circle, Rectangle, Arrow
import seaborn as sns
import pandas as pd
import numpy as np
from wordcloud import WordCloud
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots
import io
from PIL import Image, ImageDraw, ImageFont
import textwrap

class RichVisualPPT:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        # Rich color palette
        self.colors = {
            # Primary colors
            'navy': RGBColor(0, 32, 96),
            'blue': RGBColor(0, 102, 204),
            'light_blue': RGBColor(0, 176, 240),
            'sky_blue': RGBColor(135, 206, 235),
            
            # Secondary colors
            'orange': RGBColor(255, 103, 31),
            'yellow': RGBColor(255, 192, 0),
            'green': RGBColor(0, 176, 80),
            'teal': RGBColor(0, 128, 128),
            
            # Accent colors
            'purple': RGBColor(102, 45, 145),
            'pink': RGBColor(255, 105, 180),
            'red': RGBColor(220, 53, 69),
            'coral': RGBColor(255, 127, 80),
            
            # Neutral colors
            'dark': RGBColor(51, 51, 51),
            'gray': RGBColor(128, 128, 128),
            'light_gray': RGBColor(241, 241, 241),
            'white': RGBColor(255, 255, 255)
        }
        
        # Icon library (using shapes and symbols)
        self.icons = {
            'growth': '↗',
            'decline': '↘',
            'target': '◎',
            'star': '★',
            'check': '✓',
            'warning': '⚠',
            'info': 'ℹ',
            'dollar': '$',
            'percent': '%',
            'users': '👥',
            'globe': '🌍',
            'rocket': '🚀',
            'bulb': '💡',
            'chart': '📊',
            'time': '⏱',
            'lock': '🔒',
            'cloud': '☁',
            'mobile': '📱',
            'email': '✉'
        }
    
    def _inches_to_float(self, inches_obj):
        """Convert Inches object to float value"""
        return inches_obj / 914400  # EMU to inches conversion
    
    def create_rich_3_slide_presentation(self, case_data):
        """Create visually rich 3-slide presentation"""
        
        # Slide 1: Problem & Market Analysis (Dense Visual)
        self._create_rich_problem_slide(case_data)
        
        # Slide 2: Solution Architecture (Infographic Heavy)
        self._create_rich_solution_slide(case_data)
        
        # Slide 3: Impact Dashboard (Data Visualization)
        self._create_rich_impact_slide(case_data)
        
        return self.prs
    
    def _create_rich_problem_slide(self, case_data):
        """Create visually rich problem analysis slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background gradient effect
        self._add_gradient_background(slide, self.colors['navy'], self.colors['light_gray'])
        
        # Title with icon
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
        tf = title_box.text_frame
        tf.text = f"⚠ {case_data['title']} - CRITICAL MARKET CHALLENGE"
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        
        # Create visual grid layout
        # Left: Problem visualization (40%)
        self._create_problem_visualization(slide, case_data, Inches(0.3), Inches(1), Inches(5), Inches(6))
        
        # Center: Market data dashboard (35%)
        self._create_market_dashboard(slide, case_data, Inches(5.5), Inches(1), Inches(4.5), Inches(6))
        
        # Right: Impact metrics (25%)
        self._create_impact_sidebar(slide, case_data, Inches(10.2), Inches(1), Inches(3), Inches(6))
        
    def _create_problem_visualization(self, slide, case_data, x, y, width, height):
        """Create rich problem visualization"""
        # Main problem box with gradient
        problem_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height
        )
        problem_box.fill.solid()
        problem_box.fill.fore_color.rgb = self.colors['red']
        problem_box.fill.transparency = 0.1
        problem_box.line.color.rgb = self.colors['red']
        problem_box.line.width = Pt(2)
        
        # Problem statement header
        header_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, Inches(0.8)
        )
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = self.colors['red']
        header_box.line.fill.background()
        
        tf = header_box.text_frame
        tf.text = "CORE PROBLEM"
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # Problem factors as interconnected bubbles
        self._add_problem_bubble_chart(slide, case_data['problem_factors'], 
                                      x + Inches(0.2), y + Inches(1), 
                                      width - Inches(0.4), height - Inches(1.5))
        
    def _add_problem_bubble_chart(self, slide, factors, x, y, width, height):
        """Add interconnected bubble visualization"""
        # Create matplotlib figure
        fig, ax = plt.subplots(figsize=(self._inches_to_float(width), self._inches_to_float(height)))
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 10)
        ax.axis('off')
        
        # Define bubble positions
        positions = [(2, 7), (5, 8), (8, 7), (2, 3), (5, 2), (8, 3)]
        colors = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#FFA07A', '#98D8C8', '#6C5CE7']
        sizes = [2000, 1800, 2200, 1600, 1900, 2100]
        
        # Draw connections
        for i in range(len(positions)):
            for j in range(i+1, len(positions)):
                ax.plot([positions[i][0], positions[j][0]], 
                       [positions[i][1], positions[j][1]], 
                       'gray', alpha=0.3, linewidth=1)
        
        # Draw bubbles
        for i, (factor, desc) in enumerate(factors.items()):
            if i < len(positions):
                circle = plt.Circle(positions[i], radius=1.2, color=colors[i], alpha=0.8)
                ax.add_patch(circle)
                
                # Add text
                wrapped_text = textwrap.fill(factor, 10)
                ax.text(positions[i][0], positions[i][1], wrapped_text, 
                       ha='center', va='center', fontsize=10, fontweight='bold', color='white')
        
        plt.tight_layout()
        
        # Save and add to slide
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=300, transparent=True, bbox_inches='tight')
        plt.close()
        
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, x, y, width=width)
        
    def _create_market_dashboard(self, slide, case_data, x, y, width, height):
        """Create market data dashboard"""
        # Container
        dash_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height
        )
        dash_box.fill.solid()
        dash_box.fill.fore_color.rgb = self.colors['white']
        dash_box.line.color.rgb = self.colors['blue']
        dash_box.line.width = Pt(1)
        
        # Header
        header_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, Inches(0.6)
        )
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = self.colors['blue']
        header_box.line.fill.background()
        
        tf = header_box.text_frame
        tf.text = "📊 MARKET DYNAMICS"
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # Add multiple mini-charts
        # 1. Market size evolution
        self._add_mini_line_chart(slide, case_data['market_growth'], 
                                 x + Inches(0.2), y + Inches(0.8), 
                                 width - Inches(0.4), Inches(1.8))
        
        # 2. Competitor landscape
        self._add_competitor_matrix(slide, case_data['competitors'],
                                   x + Inches(0.2), y + Inches(2.8),
                                   width - Inches(0.4), Inches(1.5))
        
        # 3. Key metrics
        self._add_metric_cards(slide, case_data['market_metrics'],
                              x + Inches(0.2), y + Inches(4.5),
                              width - Inches(0.4), Inches(1.3))
        
    def _add_mini_line_chart(self, slide, data, x, y, width, height):
        """Add small line chart with area fill"""
        fig, ax = plt.subplots(figsize=(self._inches_to_float(width), self._inches_to_float(height)))
        
        years = data['years']
        values = data['values']
        
        # Plot with gradient fill
        ax.plot(years, values, color='#0066CC', linewidth=3, marker='o', markersize=8)
        ax.fill_between(years, values, alpha=0.3, color='#0066CC')
        
        # Styling
        ax.set_title(data['title'], fontsize=12, fontweight='bold', pad=10)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.grid(True, alpha=0.3, linestyle='--')
        
        # Add value labels
        for i, (year, value) in enumerate(zip(years, values)):
            ax.text(year, value + max(values)*0.05, f'${value}B', 
                   ha='center', fontsize=9, fontweight='bold')
        
        plt.tight_layout()
        
        # Save and add to slide
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=200, transparent=True, bbox_inches='tight')
        plt.close()
        
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, x, y, width=width)
        
    def _add_competitor_matrix(self, slide, competitors, x, y, width, height):
        """Add competitor comparison matrix"""
        # Create visual matrix
        fig, ax = plt.subplots(figsize=(self._inches_to_float(width), self._inches_to_float(height)))
        
        # Data
        companies = competitors['names']
        market_share = competitors['market_share']
        growth_rate = competitors['growth_rate']
        
        # Create scatter plot
        colors = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#FFA07A', '#98D8C8']
        
        scatter = ax.scatter(market_share, growth_rate, s=[1000]*len(companies), 
                           c=colors[:len(companies)], alpha=0.7, edgecolors='black', linewidth=2)
        
        # Add company names
        for i, company in enumerate(companies):
            ax.annotate(company, (market_share[i], growth_rate[i]), 
                       ha='center', va='center', fontweight='bold', fontsize=9)
        
        # Styling
        ax.set_xlabel('Market Share (%)', fontweight='bold')
        ax.set_ylabel('Growth Rate (%)', fontweight='bold')
        ax.set_title('Competitive Landscape', fontsize=11, fontweight='bold')
        ax.grid(True, alpha=0.3)
        
        plt.tight_layout()
        
        # Save and add to slide
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=200, transparent=True, bbox_inches='tight')
        plt.close()
        
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, x, y, width=width)
        
    def _add_metric_cards(self, slide, metrics, x, y, width, height):
        """Add metric cards in a row"""
        card_width = width / len(metrics) - Inches(0.1)
        
        colors = [self.colors['orange'], self.colors['green'], self.colors['purple'], self.colors['teal']]
        
        for i, metric in enumerate(metrics):
            x_pos = x + (card_width + Inches(0.1)) * i
            
            # Card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_pos, y, card_width, height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = colors[i % len(colors)]
            card.line.fill.background()
            
            tf = card.text_frame
            tf.margin_all = Inches(0.1)
            
            # Metric name
            p = tf.paragraphs[0]
            p.text = metric['name']
            p.font.size = Pt(9)
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            
            # Metric value
            p = tf.add_paragraph()
            p.text = metric['value']
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            
            # Change indicator
            if 'change' in metric:
                p = tf.add_paragraph()
                p.text = metric['change']
                p.font.size = Pt(8)
                p.font.color.rgb = self.colors['white']
                p.alignment = PP_ALIGN.CENTER
                
    def _create_impact_sidebar(self, slide, case_data, x, y, width, height):
        """Create impact metrics sidebar"""
        # Container
        sidebar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, height
        )
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = self.colors['dark']
        sidebar.line.fill.background()
        
        # Header
        tf = sidebar.text_frame
        tf.margin_all = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = "BUSINESS IMPACT"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.colors['yellow']
        p.alignment = PP_ALIGN.CENTER
        
        # Impact items with icons
        impacts = case_data['business_impacts']
        
        for impact in impacts:
            p = tf.add_paragraph()
            p.text = ""  # Empty line for spacing
            
            # Icon and title
            p = tf.add_paragraph()
            p.text = f"{impact['icon']} {impact['metric']}"
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            
            # Value
            p = tf.add_paragraph()
            p.text = impact['value']
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.colors['yellow']
            p.alignment = PP_ALIGN.CENTER
            
            # Description
            p = tf.add_paragraph()
            p.text = impact['desc']
            p.font.size = Pt(8)
            p.font.color.rgb = self.colors['light_gray']
            p.line_spacing = 1.1
            
    def _create_rich_solution_slide(self, case_data):
        """Create solution slide with rich infographics"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Gradient background
        self._add_gradient_background(slide, self.colors['white'], self.colors['light_gray'])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
        tf = title_box.text_frame
        tf.text = f"💡 {case_data['solution_name']} - COMPREHENSIVE SOLUTION ARCHITECTURE"
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = self.colors['navy']
        
        # Create three main sections
        # 1. Solution framework (center piece)
        self._create_solution_framework(slide, case_data, Inches(3), Inches(1), Inches(7.3), Inches(4))
        
        # 2. Technology stack (left)
        self._create_tech_stack_visual(slide, case_data['tech_stack'], Inches(0.3), Inches(1), Inches(2.5), Inches(6))
        
        # 3. Implementation phases (bottom)
        self._create_phase_timeline(slide, case_data['phases'], Inches(0.3), Inches(5.2), Inches(12.7), Inches(1.8))
        
        # 4. Key features (right)
        self._create_feature_cards(slide, case_data['key_features'], Inches(10.5), Inches(1), Inches(2.5), Inches(4))
        
    def _create_solution_framework(self, slide, case_data, x, y, width, height):
        """Create central solution framework visualization"""
        # Create interconnected system diagram
        fig, ax = plt.subplots(figsize=(self._inches_to_float(width), self._inches_to_float(height)))
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 10)
        ax.axis('off')
        
        # Central hub
        center = (5, 5)
        hub = Circle(center, 1.5, color='#0066CC', alpha=0.8)
        ax.add_patch(hub)
        ax.text(center[0], center[1], case_data['solution_core'], 
               ha='center', va='center', fontsize=12, fontweight='bold', color='white', wrap=True)
        
        # Surrounding components
        components = case_data['solution_components']
        angles = np.linspace(0, 2*np.pi, len(components), endpoint=False)
        radius = 3
        
        component_colors = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#FFA07A', '#98D8C8', '#6C5CE7']
        
        for i, (comp_name, comp_details) in enumerate(components.items()):
            # Calculate position
            x_pos = center[0] + radius * np.cos(angles[i])
            y_pos = center[1] + radius * np.sin(angles[i])
            
            # Draw component circle
            comp_circle = Circle((x_pos, y_pos), 1, color=component_colors[i % len(component_colors)], alpha=0.8)
            ax.add_patch(comp_circle)
            
            # Add text
            wrapped_text = textwrap.fill(comp_name, 12)
            ax.text(x_pos, y_pos, wrapped_text, ha='center', va='center', 
                   fontsize=9, fontweight='bold', color='white')
            
            # Draw connection to center
            ax.plot([center[0], x_pos], [center[1], y_pos], 
                   color='gray', linewidth=2, alpha=0.5, linestyle='--')
            
            # Add sub-components
            if 'sub_items' in comp_details:
                sub_radius = 1.5
                sub_angles = np.linspace(angles[i] - 0.3, angles[i] + 0.3, len(comp_details['sub_items']), endpoint=True)
                
                for j, sub_item in enumerate(comp_details['sub_items'][:3]):
                    sub_x = x_pos + sub_radius * np.cos(sub_angles[j])
                    sub_y = y_pos + sub_radius * np.sin(sub_angles[j])
                    
                    # Small circle for sub-item
                    sub_circle = Circle((sub_x, sub_y), 0.3, color='lightgray', alpha=0.8)
                    ax.add_patch(sub_circle)
                    
                    # Connection
                    ax.plot([x_pos, sub_x], [y_pos, sub_y], 
                           color='lightgray', linewidth=1, alpha=0.5)
        
        # Add title
        ax.text(5, 9, "INTEGRATED SOLUTION ECOSYSTEM", ha='center', fontsize=14, fontweight='bold')
        
        plt.tight_layout()
        
        # Save and add to slide
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=300, transparent=True, bbox_inches='tight')
        plt.close()
        
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, x, y, width=width)
        
    def _create_tech_stack_visual(self, slide, tech_stack, x, y, width, height):
        """Create technology stack visualization"""
        # Container
        stack_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height
        )
        stack_box.fill.solid()
        stack_box.fill.fore_color.rgb = self.colors['navy']
        stack_box.line.fill.background()
        
        # Header
        header_height = 0.6
        tf = stack_box.text_frame
        tf.margin_all = Inches(0.15)
        
        p = tf.paragraphs[0]
        p.text = "🛠 TECH STACK"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors['yellow']
        p.alignment = PP_ALIGN.CENTER
        
        # Stack layers
        layer_height = (height - Inches(header_height * 2)) / len(tech_stack)
        colors = [self.colors['light_blue'], self.colors['teal'], self.colors['green'], 
                  self.colors['orange'], self.colors['purple']]
        
        for i, (layer_name, technologies) in enumerate(tech_stack.items()):
            # Layer box
            layer_y = y + Inches(header_height * 1.5) + layer_height * i
            
            layer_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x + Inches(0.1), layer_y,
                width - Inches(0.2), layer_height - Inches(0.1)
            )
            layer_box.fill.solid()
            layer_box.fill.fore_color.rgb = colors[i % len(colors)]
            layer_box.line.color.rgb = self.colors['white']
            layer_box.line.width = Pt(1)
            
            tf_layer = layer_box.text_frame
            tf_layer.margin_all = Inches(0.1)
            
            # Layer name
            p = tf_layer.paragraphs[0]
            p.text = layer_name.upper()
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            
            # Technologies
            p = tf_layer.add_paragraph()
            p.text = " • ".join(technologies[:3])
            p.font.size = Pt(8)
            p.font.color.rgb = self.colors['white']
            p.line_spacing = 1
            
    def _create_phase_timeline(self, slide, phases, x, y, width, height):
        """Create implementation phase timeline"""
        # Create matplotlib figure for timeline
        fig, ax = plt.subplots(figsize=(self._inches_to_float(width), self._inches_to_float(height)))
        ax.set_xlim(0, 12)
        ax.set_ylim(0, 3)
        ax.axis('off')
        
        # Timeline base
        timeline_y = 1.5
        ax.plot([0.5, 11.5], [timeline_y, timeline_y], color='#333333', linewidth=4)
        
        # Phase segments
        segment_width = 11 / len(phases)
        colors = ['#FF6B6B', '#FFA07A', '#FFD700', '#98D8C8', '#4ECDC4']
        
        for i, phase in enumerate(phases):
            x_start = 0.5 + segment_width * i
            x_end = x_start + segment_width
            
            # Phase bar
            rect = Rectangle((x_start, timeline_y - 0.3), segment_width - 0.1, 0.6,
                           facecolor=colors[i % len(colors)], edgecolor='black', linewidth=2)
            ax.add_patch(rect)
            
            # Phase name
            ax.text(x_start + segment_width/2, timeline_y + 0.6, phase['name'],
                   ha='center', fontsize=10, fontweight='bold')
            
            # Duration
            ax.text(x_start + segment_width/2, timeline_y, phase['duration'],
                   ha='center', va='center', fontsize=8, color='white', fontweight='bold')
            
            # Key milestone
            ax.text(x_start + segment_width/2, timeline_y - 0.8, phase['milestone'],
                   ha='center', fontsize=8, style='italic', wrap=True)
            
            # Milestone marker
            ax.plot(x_start + segment_width/2, timeline_y - 0.3, 'wo', markersize=8, 
                   markeredgecolor='black', markeredgewidth=2)
        
        # Title
        ax.text(6, 2.5, "IMPLEMENTATION ROADMAP", ha='center', fontsize=12, fontweight='bold')
        
        plt.tight_layout()
        
        # Save and add to slide
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=300, transparent=True, bbox_inches='tight')
        plt.close()
        
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, x, y, width=width)
        
    def _create_feature_cards(self, slide, features, x, y, width, height):
        """Create feature cards on the right"""
        # Container
        feature_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height
        )
        feature_box.fill.solid()
        feature_box.fill.fore_color.rgb = self.colors['light_gray']
        feature_box.line.color.rgb = self.colors['gray']
        feature_box.line.width = Pt(1)
        
        # Header
        header_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, Inches(0.5)
        )
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = self.colors['green']
        header_box.line.fill.background()
        
        tf = header_box.text_frame
        tf.text = "✨ KEY FEATURES"
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # Feature items
        feature_height = (height - Inches(0.6)) / len(features)
        
        for i, feature in enumerate(features[:5]):  # Max 5 features
            feature_y = y + Inches(0.6) + feature_height * i
            
            # Feature item box
            item_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x + Inches(0.1), feature_y,
                width - Inches(0.2), feature_height - Inches(0.1)
            )
            item_box.fill.solid()
            item_box.fill.fore_color.rgb = self.colors['white']
            item_box.line.color.rgb = self.colors['green']
            item_box.line.width = Pt(1)
            
            tf = item_box.text_frame
            tf.margin_all = Inches(0.1)
            
            # Feature name with icon
            p = tf.paragraphs[0]
            p.text = f"{feature['icon']} {feature['name']}"
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = self.colors['dark']
            
            # Feature description
            if 'desc' in feature:
                p = tf.add_paragraph()
                p.text = feature['desc']
                p.font.size = Pt(7)
                p.font.color.rgb = self.colors['gray']
                p.line_spacing = 1
                
    def _create_rich_impact_slide(self, case_data):
        """Create impact dashboard slide with rich visualizations"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Dark background for dashboard effect
        self._add_gradient_background(slide, self.colors['dark'], self.colors['navy'])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
        tf = title_box.text_frame
        tf.text = f"📈 PROJECTED IMPACT & ROI DASHBOARD"
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        
        # Create dashboard grid
        # Top row: 3 KPI cards
        self._create_kpi_cards(slide, case_data['kpis'], Inches(0.3), Inches(1), Inches(12.7), Inches(1.5))
        
        # Middle left: Revenue projection chart
        self._create_revenue_chart(slide, case_data['revenue_projection'], 
                                  Inches(0.3), Inches(2.7), Inches(6), Inches(2.5))
        
        # Middle right: ROI gauge
        self._create_roi_gauge(slide, case_data['roi_metrics'],
                              Inches(6.5), Inches(2.7), Inches(3.5), Inches(2.5))
        
        # Right side: Success metrics
        self._create_success_metrics(slide, case_data['success_metrics'],
                                    Inches(10.2), Inches(2.7), Inches(2.8), Inches(2.5))
        
        # Bottom: Implementation timeline with milestones
        self._create_milestone_tracker(slide, case_data['milestones'],
                                      Inches(0.3), Inches(5.4), Inches(12.7), Inches(1.6))
        
    def _create_kpi_cards(self, slide, kpis, x, y, width, height):
        """Create KPI cards row"""
        card_width = width / len(kpis) - Inches(0.1)
        
        for i, kpi in enumerate(kpis[:4]):  # Max 4 KPIs
            x_pos = x + (card_width + Inches(0.1)) * i
            
            # Card with gradient effect
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_pos, y, card_width, height
            )
            
            # Different colors for each KPI
            colors = [self.colors['green'], self.colors['blue'], 
                     self.colors['orange'], self.colors['purple']]
            card.fill.solid()
            card.fill.fore_color.rgb = colors[i % len(colors)]
            card.line.fill.background()
            
            tf = card.text_frame
            tf.margin_all = Inches(0.15)
            
            # KPI Icon
            p = tf.paragraphs[0]
            p.text = kpi['icon']
            p.font.size = Pt(20)
            p.alignment = PP_ALIGN.CENTER
            
            # KPI Name
            p = tf.add_paragraph()
            p.text = kpi['name']
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            
            # KPI Value
            p = tf.add_paragraph()
            p.text = kpi['value']
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            
            # Change indicator
            if 'change' in kpi:
                p = tf.add_paragraph()
                p.text = kpi['change']
                p.font.size = Pt(9)
                p.font.color.rgb = self.colors['white']
                p.alignment = PP_ALIGN.CENTER
                
    def _create_revenue_chart(self, slide, revenue_data, x, y, width, height):
        """Create revenue projection chart"""
        fig, ax = plt.subplots(figsize=(self._inches_to_float(width), self._inches_to_float(height)), facecolor='none')
        
        years = revenue_data['years']
        revenue = revenue_data['revenue']
        costs = revenue_data['costs']
        profit = [r - c for r, c in zip(revenue, costs)]
        
        # Create bar chart with line overlay
        x_pos = np.arange(len(years))
        
        # Revenue bars
        bars1 = ax.bar(x_pos - 0.2, revenue, 0.4, label='Revenue', color='#4ECDC4', alpha=0.8)
        # Cost bars
        bars2 = ax.bar(x_pos + 0.2, costs, 0.4, label='Costs', color='#FF6B6B', alpha=0.8)
        
        # Profit line
        ax2 = ax.twinx()
        line = ax2.plot(x_pos, profit, color='#FFD700', marker='o', markersize=10, 
                       linewidth=3, label='Profit', markeredgecolor='black', markeredgewidth=2)
        
        # Add value labels
        for i, (r, c, p) in enumerate(zip(revenue, costs, profit)):
            ax.text(i - 0.2, r + max(revenue)*0.02, f'${r}M', ha='center', fontsize=8, fontweight='bold')
            ax.text(i + 0.2, c + max(revenue)*0.02, f'${c}M', ha='center', fontsize=8, fontweight='bold')
            ax2.text(i, p + max(profit)*0.05, f'${p}M', ha='center', fontsize=9, 
                    fontweight='bold', bbox=dict(boxstyle="round,pad=0.3", facecolor='yellow', alpha=0.7))
        
        # Styling
        ax.set_xlabel('Year', fontsize=10, fontweight='bold', color='white')
        ax.set_ylabel('Revenue & Costs ($M)', fontsize=10, fontweight='bold', color='white')
        ax2.set_ylabel('Profit ($M)', fontsize=10, fontweight='bold', color='#FFD700')
        ax.set_title('Financial Projections', fontsize=12, fontweight='bold', color='white', pad=10)
        
        ax.set_xticks(x_pos)
        ax.set_xticklabels(years)
        ax.tick_params(colors='white')
        ax2.tick_params(colors='#FFD700')
        
        # Grid
        ax.grid(True, alpha=0.3, color='white')
        ax.set_axisbelow(True)
        
        # Legend
        ax.legend(loc='upper left', framealpha=0.8)
        ax2.legend(loc='upper right', framealpha=0.8)
        
        # Remove spines
        for spine in ax.spines.values():
            spine.set_color('white')
            spine.set_alpha(0.3)
        for spine in ax2.spines.values():
            spine.set_visible(False)
            
        plt.tight_layout()
        
        # Save and add to slide
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=300, transparent=True, bbox_inches='tight')
        plt.close()
        
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, x, y, width=width)
        
    def _create_roi_gauge(self, slide, roi_data, x, y, width, height):
        """Create ROI gauge visualization using matplotlib"""
        # Container
        gauge_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height
        )
        gauge_box.fill.solid()
        gauge_box.fill.fore_color.rgb = self.colors['white']
        gauge_box.line.color.rgb = self.colors['gray']
        gauge_box.line.width = Pt(1)
        
        # Create gauge using matplotlib
        fig, ax = plt.subplots(figsize=(self._inches_to_float(width), self._inches_to_float(height)), 
                              subplot_kw={'projection': 'polar'})
        
        # Gauge parameters
        current_value = roi_data['current']
        max_value = 500
        
        # Create gauge segments
        theta = np.linspace(np.pi, 0, 1000)
        r = np.ones_like(theta)
        
        # Color segments
        colors_ranges = [(0, 100, '#FFE5E5'), (100, 200, '#FFFACD'), 
                        (200, 300, '#E5FFE5'), (300, 500, '#90EE90')]
        
        for start, end, color in colors_ranges:
            mask = (theta >= np.pi * (1 - end/max_value)) & (theta <= np.pi * (1 - start/max_value))
            ax.fill_between(theta[mask], 0.8, 1, color=color, alpha=0.8)
        
        # Add gauge needle
        angle = np.pi * (1 - current_value/max_value)
        ax.plot([angle, angle], [0, 0.85], color='darkblue', linewidth=4)
        ax.plot(angle, 0.85, 'o', color='darkblue', markersize=10)
        
        # Add center circle
        circle = plt.Circle((0, 0), 0.3, color='white', transform=ax.transProjectionAffine + ax.transAxes)
        ax.add_patch(circle)
        
        # Add value text
        ax.text(0, 0, f'{current_value}%', ha='center', va='center', 
               fontsize=20, fontweight='bold', color='darkblue',
               transform=ax.transAxes)
        
        # Add title
        ax.text(0.5, 1.1, 'ROI %', ha='center', va='center', 
               fontsize=16, fontweight='bold', transform=ax.transAxes)
        
        # Remove grid and labels
        ax.set_ylim(0, 1)
        ax.set_yticklabels([])
        ax.set_xticklabels([])
        ax.grid(False)
        ax.spines['polar'].set_visible(False)
        
        plt.tight_layout()
        
        # Save gauge
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=200, transparent=True, bbox_inches='tight')
        plt.close()
        
        img_stream.seek(0)
        
        # Add to slide
        gauge_pic = slide.shapes.add_picture(img_stream, x + Inches(0.1), y + Inches(0.1), 
                                           width=width - Inches(0.2))
        
        # Add text below gauge
        text_box = slide.shapes.add_textbox(x, y + height - Inches(0.5), width, Inches(0.4))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Payback: {roi_data['payback_period']}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = self.colors['dark']
        p.alignment = PP_ALIGN.CENTER
        
    def _create_success_metrics(self, slide, metrics, x, y, width, height):
        """Create success metrics panel"""
        # Container
        metrics_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height
        )
        metrics_box.fill.solid()
        metrics_box.fill.fore_color.rgb = self.colors['purple']
        metrics_box.line.fill.background()
        
        tf = metrics_box.text_frame
        tf.margin_all = Inches(0.15)
        
        # Header
        p = tf.paragraphs[0]
        p.text = "SUCCESS METRICS"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.colors['yellow']
        p.alignment = PP_ALIGN.CENTER
        
        # Metrics
        for metric in metrics:
            p = tf.add_paragraph()
            p.text = ""  # Spacing
            
            # Metric name
            p = tf.add_paragraph()
            p.text = metric['name']
            p.font.size = Pt(9)
            p.font.color.rgb = self.colors['white']
            
            # Progress bar visualization
            p = tf.add_paragraph()
            progress = int(metric['progress'] / 10)
            p.text = "█" * progress + "░" * (10 - progress)
            p.font.size = Pt(8)
            p.font.color.rgb = self.colors['yellow']
            
            # Value
            p = tf.add_paragraph()
            p.text = f"{metric['current']} / {metric['target']}"
            p.font.size = Pt(8)
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            
    def _create_milestone_tracker(self, slide, milestones, x, y, width, height):
        """Create milestone tracker at bottom"""
        # Container with gradient
        tracker_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, height
        )
        tracker_box.fill.solid()
        tracker_box.fill.fore_color.rgb = self.colors['light_gray']
        tracker_box.fill.transparency = 0.5
        tracker_box.line.color.rgb = self.colors['gray']
        tracker_box.line.width = Pt(1)
        
        # Create milestone visualization
        milestone_width = width / len(milestones) - Inches(0.1)
        
        for i, milestone in enumerate(milestones):
            x_pos = x + (milestone_width + Inches(0.1)) * i + Inches(0.05)
            
            # Milestone card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_pos, y + Inches(0.1),
                milestone_width, height - Inches(0.2)
            )
            
            # Status-based coloring
            status_colors = {
                'completed': self.colors['green'],
                'in_progress': self.colors['yellow'],
                'upcoming': self.colors['gray']
            }
            card.fill.solid()
            card.fill.fore_color.rgb = status_colors.get(milestone['status'], self.colors['gray'])
            card.line.fill.background()
            
            tf = card.text_frame
            tf.margin_all = Inches(0.1)
            
            # Milestone name
            p = tf.paragraphs[0]
            p.text = milestone['name']
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            
            # Date
            p = tf.add_paragraph()
            p.text = milestone['date']
            p.font.size = Pt(8)
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            
            # Status icon
            status_icons = {
                'completed': '✓',
                'in_progress': '◈',
                'upcoming': '○'
            }
            p = tf.add_paragraph()
            p.text = status_icons.get(milestone['status'], '○')
            p.font.size = Pt(12)
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            
    def _add_gradient_background(self, slide, color1, color2):
        """Add gradient background to slide"""
        # Create gradient effect using shapes
        gradient_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0), 
            self.prs.slide_width, self.prs.slide_height
        )
        gradient_box.fill.solid()
        gradient_box.fill.fore_color.rgb = color1
        gradient_box.fill.transparency = 0.3
        gradient_box.line.fill.background()
        
        # Send to back
        slide.shapes._spTree.remove(gradient_box._element)
        slide.shapes._spTree.insert(2, gradient_box._element)
        
    def save(self, filename):
        """Save the presentation"""
        self.prs.save(filename)
        return filename


# Example usage with rich data
def create_rich_visual_example():
    """Create example rich visual presentation"""
    ppt = RichVisualPPT()
    
    case_data = {
        # Slide 1 - Problem Analysis
        'title': 'Healthcare Digital Divide Crisis',
        'problem_factors': {
            'Access Barriers': 'Rural areas lack infrastructure',
            'Technology Gap': '60% hospitals use legacy systems',
            'Cost Constraints': '$500B annual inefficiency',
            'Patient Experience': '4-hour average wait times',
            'Data Silos': 'No interoperability between systems',
            'Skill Shortage': '200K healthcare IT professionals needed'
        },
        'market_growth': {
            'title': 'Digital Health Market Growth',
            'years': [2020, 2021, 2022, 2023, 2024, 2025],
            'values': [150, 180, 250, 380, 520, 750]
        },
        'competitors': {
            'names': ['TeleHealth+', 'MedConnect', 'HealthHub', 'CareCloud', 'Our Solution'],
            'market_share': [25, 20, 15, 12, 5],
            'growth_rate': [15, 25, 30, 20, 85]
        },
        'market_metrics': [
            {'name': 'TAM', 'value': '$750B', 'change': '↗ +35%'},
            {'name': 'CAGR', 'value': '28.5%', 'change': '↗ +5pp'},
            {'name': 'Penetration', 'value': '18%', 'change': '↗ +7%'},
            {'name': 'Adoption', 'value': '2.5M', 'change': '↗ +125%'}
        ],
        'business_impacts': [
            {'icon': '💰', 'metric': 'Revenue Loss', 'value': '$500B', 'desc': 'Annual healthcare inefficiency'},
            {'icon': '⏱', 'metric': 'Time Waste', 'value': '4.2hrs', 'desc': 'Average patient wait time'},
            {'icon': '📉', 'metric': 'Satisfaction', 'value': '42%', 'desc': 'Patient satisfaction score'},
            {'icon': '⚠️', 'metric': 'Error Rate', 'value': '23%', 'desc': 'Medical record errors'}
        ],
        
        # Slide 2 - Solution Architecture
        'solution_name': 'HealthConnect AI Platform',
        'solution_core': 'AI-Powered\nUnified Health\nEcosystem',
        'solution_components': {
            'Patient Portal': {
                'sub_items': ['Mobile App', 'Web Access', 'Voice Assistant']
            },
            'Provider Dashboard': {
                'sub_items': ['Real-time Analytics', 'Patient History', 'Prescription AI']
            },
            'AI Engine': {
                'sub_items': ['Diagnosis Support', 'Treatment Recommendations', 'Risk Prediction']
            },
            'Data Integration': {
                'sub_items': ['EHR Connect', 'Lab Systems', 'Insurance API']
            },
            'Telehealth Suite': {
                'sub_items': ['Video Consults', 'Remote Monitoring', 'Digital Triage']
            },
            'Admin Console': {
                'sub_items': ['Resource Planning', 'Cost Analytics', 'Compliance Tools']
            }
        },
        'tech_stack': {
            'Frontend': ['React Native', 'Flutter', 'Progressive Web'],
            'Backend': ['Node.js', 'Python', 'GraphQL'],
            'AI/ML': ['TensorFlow', 'PyTorch', 'AutoML'],
            'Infrastructure': ['AWS', 'Kubernetes', 'Terraform'],
            'Security': ['OAuth2', 'Blockchain', 'Zero Trust']
        },
        'phases': [
            {'name': 'Foundation', 'duration': 'Q1 2024', 'milestone': 'Core platform launch'},
            {'name': 'Integration', 'duration': 'Q2 2024', 'milestone': 'Connect 100 hospitals'},
            {'name': 'AI Rollout', 'duration': 'Q3 2024', 'milestone': 'Deploy ML models'},
            {'name': 'Scale', 'duration': 'Q4 2024', 'milestone': '1M active users'},
            {'name': 'Optimize', 'duration': 'Q1 2025', 'milestone': 'Full automation'}
        ],
        'key_features': [
            {'icon': '🏥', 'name': 'Instant Access', 'desc': '24/7 healthcare services'},
            {'icon': '🤖', 'name': 'AI Diagnosis', 'desc': '95% accuracy rate'},
            {'icon': '💊', 'name': 'Smart Prescriptions', 'desc': 'Automated refills'},
            {'icon': '📱', 'name': 'Mobile First', 'desc': 'Complete care on-the-go'},
            {'icon': '🔐', 'name': 'Secure & Private', 'desc': 'HIPAA compliant'}
        ],
        
        # Slide 3 - Impact Dashboard
        'kpis': [
            {'icon': '📈', 'name': 'Revenue Growth', 'value': '+285%', 'change': 'vs. baseline'},
            {'icon': '⏰', 'name': 'Wait Time', 'value': '15 min', 'change': '↘ -73%'},
            {'icon': '😊', 'name': 'Satisfaction', 'value': '94%', 'change': '↗ +52pp'},
            {'icon': '💰', 'name': 'Cost Savings', 'value': '$125M', 'change': 'per hospital'}
        ],
        'revenue_projection': {
            'years': ['Y1', 'Y2', 'Y3', 'Y4', 'Y5'],
            'revenue': [50, 150, 380, 750, 1200],
            'costs': [80, 120, 180, 250, 320]
        },
        'roi_metrics': {
            'current': 380,
            'target': 250,
            'payback_period': '14 months'
        },
        'success_metrics': [
            {'name': 'User Adoption', 'current': 850, 'target': 1000, 'progress': 85},
            {'name': 'System Uptime', 'current': 99.8, 'target': 99.9, 'progress': 95},
            {'name': 'AI Accuracy', 'current': 94, 'target': 97, 'progress': 90},
            {'name': 'Cost Reduction', 'current': 28, 'target': 35, 'progress': 80}
        ],
        'milestones': [
            {'name': 'MVP Launch', 'date': 'Jan 2024', 'status': 'completed'},
            {'name': 'First Hospital', 'date': 'Mar 2024', 'status': 'completed'},
            {'name': '100K Users', 'date': 'Jun 2024', 'status': 'in_progress'},
            {'name': 'Series A', 'date': 'Sep 2024', 'status': 'upcoming'},
            {'name': 'National Scale', 'date': 'Dec 2024', 'status': 'upcoming'}
        ]
    }
    
    ppt.create_rich_3_slide_presentation(case_data)
    filename = "Rich_Visual_3_Slide_Presentation.pptx"
    ppt.save(filename)
    print(f"✓ Created rich visual presentation: {filename}")
    return filename


if __name__ == "__main__":
    print("Creating rich visual presentation with maximum infographics...")
    create_rich_visual_example()
    print("\n✓ Rich visual PPT system ready!")
    print("✓ Every slide packed with infographics and visuals")
    print("✓ No empty spaces - maximum visual impact")