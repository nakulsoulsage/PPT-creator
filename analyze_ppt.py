#!/usr/bin/env python3
"""
Analyze existing PowerPoint presentations to identify issues
"""

from pptx import Presentation
import os

def analyze_presentation(file_path):
    """Analyze a PowerPoint presentation for issues"""
    print(f"\n{'='*60}")
    print(f"Analyzing: {os.path.basename(file_path)}")
    print(f"{'='*60}")

    try:
        prs = Presentation(file_path)
        print(f"Total slides: {len(prs.slides)}")
        print(f"Slide dimensions: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")

        for i, slide in enumerate(prs.slides, 1):
            print(f"\nSlide {i}:")

            # Count shapes
            shapes_count = len(slide.shapes)
            print(f"  Total shapes: {shapes_count}")

            # Analyze text boxes
            text_boxes = []
            charts = []
            pictures = []
            tables = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_boxes.append(shape)
                    # Check for overlapping text
                    text = shape.text_frame.text
                    if text:
                        print(f"  - Text box: {text[:50]}..." if len(text) > 50 else f"  - Text box: {text}")
                        print(f"    Position: ({shape.left.inches:.2f}\", {shape.top.inches:.2f}\")")
                        print(f"    Size: ({shape.width.inches:.2f}\" x {shape.height.inches:.2f}\")")

                if shape.has_chart:
                    charts.append(shape)
                    print(f"  - Chart found at position ({shape.left.inches:.2f}\", {shape.top.inches:.2f}\")")

                if shape.shape_type == 13:  # Picture
                    pictures.append(shape)
                    print(f"  - Picture/Image at position ({shape.left.inches:.2f}\", {shape.top.inches:.2f}\")")
                    print(f"    Size: ({shape.width.inches:.2f}\" x {shape.height.inches:.2f}\")")

                if shape.has_table:
                    tables.append(shape)
                    print(f"  - Table with {shape.table.rows} rows and {shape.table.columns} columns")

            print(f"  Summary: {len(text_boxes)} text boxes, {len(charts)} charts, {len(pictures)} images, {len(tables)} tables")

            # Check for potential overlaps
            shapes_positions = []
            for shape in slide.shapes:
                if hasattr(shape, 'left') and hasattr(shape, 'top'):
                    shapes_positions.append({
                        'left': shape.left.inches,
                        'top': shape.top.inches,
                        'right': shape.left.inches + shape.width.inches,
                        'bottom': shape.top.inches + shape.height.inches
                    })

            # Simple overlap detection
            overlaps = []
            for i in range(len(shapes_positions)):
                for j in range(i+1, len(shapes_positions)):
                    s1 = shapes_positions[i]
                    s2 = shapes_positions[j]

                    # Check if rectangles overlap
                    if not (s1['right'] < s2['left'] or s1['left'] > s2['right'] or
                           s1['bottom'] < s2['top'] or s1['top'] > s2['bottom']):
                        overlaps.append((i, j))

            if overlaps:
                print(f"  ⚠️  Warning: Detected {len(overlaps)} potential overlapping elements")

    except Exception as e:
        print(f"Error analyzing {file_path}: {e}")

# Analyze all presentations in PPT Generated folder
ppt_folder = "/mnt/e/AI and Projects/Case Comp PPT/PPT Generated"
for file in os.listdir(ppt_folder):
    if file.endswith('.pptx') and not file.startswith('~$'):
        analyze_presentation(os.path.join(ppt_folder, file))

# Also analyze presentations in main directory
main_folder = "/mnt/e/AI and Projects/Case Comp PPT"
main_ppts = [
    "Case_Competition_Presentation.pptx",
    "DisruptX_Round1_MediChain.pptx",
    "Professional_Case_Competition_Presentation.pptx",
    "Rich_Visual_3_Slide_Presentation.pptx",
    "Ultra_Condensed_3_Slide_Presentation.pptx",
    "Ultra_Condensed_5_Slide_Presentation.pptx"
]

for file in main_ppts:
    file_path = os.path.join(main_folder, file)
    if os.path.exists(file_path):
        analyze_presentation(file_path)